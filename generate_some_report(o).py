"""
SoMe Report Generator
=====================
Usage:
    python generate_some_report.py
    python generate_some_report.py "SoMe_Data(2).xlsx"

Make sure your Excel file is in the same folder as this script.
Output: SoMe_Report.pptx  (saved in same folder)

Excel structure expected:
  Sheet1 — 2025 Totals
    A1: Social Media, B1: Followers, C1: Impressions, D1: Engagement
    A2: Facebook, A3: Instagram, A4: LinkedIn (data in rows 2-4)

  Sheet2 — Monthly 2026 Data (add one row per month at the bottom)
    Row 1: Headers
    A: Month | B: FB Followers | C: FB Impressions | D: FB Engagement
             | E: IG Followers | F: IG Impressions | G: IG Engagement
             | H: LI Followers | I: LI Impressions | J: LI Engagement
    Row 2 onwards: one row per month (January, February, ...)
    → To add a new month: just add a new row at the bottom and run the script!
"""

import sys, os, io, re
from copy import deepcopy
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import zipfile, shutil, tempfile

# ── Try to locate the Excel file ─────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Accept filename as command-line argument, otherwise use default
if len(sys.argv) > 1:
    EXCEL_PATH = os.path.join(SCRIPT_DIR, sys.argv[1])
else:
    EXCEL_PATH = os.path.join(SCRIPT_DIR, "SoMe_Data(2).xlsx")

if not os.path.exists(EXCEL_PATH):
    print(f"❌  Could not find {os.path.basename(EXCEL_PATH)} in the same folder as this script.")
    print(f"    Please make sure the file is in: {SCRIPT_DIR}")
    sys.exit(1)

OUTPUT_PATH = os.path.join(SCRIPT_DIR, "SoMe_Report.pptx")

print(f"Reading: {EXCEL_PATH}")

# ── Read Excel ────────────────────────────────────────────────────────────────
wb = openpyxl.load_workbook(EXCEL_PATH, data_only=True, read_only=False)

print(f"Sheets found: {[ws.title for ws in wb.worksheets]}")

ws1 = wb.worksheets[0]   # 2025 totals
ws2 = wb.worksheets[1]   # all 2026 monthly data

PLATFORMS = ["Facebook", "Instagram", "LinkedIn"]

# ── Read Sheet1: 2025 totals ──────────────────────────────────────────────────
# A2:D2=Facebook, A3:D3=Instagram, A4:D4=LinkedIn
# Cols: B=Followers, C=Impressions, D=Engagement
data_2025 = {}
for i, platform in enumerate(PLATFORMS):
    row = i + 2
    data_2025[platform] = {
        "followers":   float(ws1.cell(row, 2).value or 0),
        "impressions": float(ws1.cell(row, 3).value or 0),
        "engagement":  float(ws1.cell(row, 4).value or 0),
    }

# ── Read Sheet2: all monthly 2026 data ───────────────────────────────────────
# Row 1 = headers, Row 2 onwards = one row per month
# A=Month | B=FB Fol | C=FB Imp | D=FB Eng
#          | E=IG Fol | F=IG Imp | G=IG Eng
#          | H=LI Fol | I=LI Imp | J=LI Eng
monthly = {"Facebook": {"months": [], "followers": [], "impressions": [], "engagement": []},
           "Instagram": {"months": [], "followers": [], "impressions": [], "engagement": []},
           "LinkedIn":  {"months": [], "followers": [], "impressions": [], "engagement": []}}

# Platform column offsets in Sheet2 (B=2, E=5, H=8)
platform_cols = {"Facebook": 2, "Instagram": 5, "LinkedIn": 8}

row = 2
while True:
    month_name = ws2.cell(row, 1).value
    if not month_name:
        break
    month_name = str(month_name).strip()
    for p, col_start in platform_cols.items():
        monthly[p]["months"].append(month_name)
        monthly[p]["followers"].append(  float(ws2.cell(row, col_start).value     or 0))
        monthly[p]["impressions"].append(float(ws2.cell(row, col_start + 1).value or 0))
        monthly[p]["engagement"].append( float(ws2.cell(row, col_start + 2).value or 0))
    row += 1

# ── Calculate 2026 totals for Slide 2 ────────────────────────────────────────
# Followers  = latest row only (point-in-time snapshot per platform)
# Impressions = sum of all months
# Engagement  = sum of all months
data_2026 = {}
for p in PLATFORMS:
    data_2026[p] = {
        "followers":   monthly[p]["followers"][-1]          if monthly[p]["followers"]   else 0,
        "impressions": sum(monthly[p]["impressions"])        if monthly[p]["impressions"] else 0,
        "engagement":  sum(monthly[p]["engagement"])         if monthly[p]["engagement"]  else 0,
    }

# ── Debug output ──────────────────────────────────────────────────────────────
print(f"\n📋 Months read from Sheet2: {monthly['Facebook']['months']}")
print(f"\n📋 2026 totals (Slide 2):")
for p in PLATFORMS:
    print(f"   {p}: Followers={int(data_2026[p]['followers']):,}  "
          f"Impressions={int(data_2026[p]['impressions']):,}  "
          f"Engagement={int(data_2026[p]['engagement']):,}")


# ── Colours ───────────────────────────────────────────────────────────────────
TEAL        = RGBColor(0x00, 0x85, 0x7C)   # #00857C
DARK_NAVY   = RGBColor(0x0E, 0x28, 0x41)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
GREEN_UP    = RGBColor(0x2E, 0x7D, 0x32)
RED_DOWN    = RGBColor(0xC6, 0x28, 0x28)
COL_FOL     = RGBColor(0x0E, 0x28, 0x41)   # followers bar
COL_IMP     = RGBColor(0x6F, 0xBE, 0x4A)   # impressions line
COL_ENG     = RGBColor(0x15, 0x85, 0x7C)   # engagement line

def rgb_hex(rgb: RGBColor):
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

# ── python-pptx helpers ───────────────────────────────────────────────────────
def inches(n): return Inches(n)
def pt(n):     return Pt(n)

def add_textbox(slide, text, x, y, w, h, size=12, bold=False, italic=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, font_face="Calibri"):
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_face
    return txBox

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE=1
        inches(x), inches(y), inches(w), inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def fmt_num(n):
    return f"{int(n):,}"

def pct_change(v2026, v2025):
    if v2025 == 0: return "+0.00%"
    c = (v2026 - v2025) / v2025 * 100
    return f"+{c:.2f}%" if c >= 0 else f"{c:.2f}%"

def pct_color(v2026, v2025):
    return GREEN_UP if v2026 >= v2025 else RED_DOWN

# ── SVG icons (inline base64) ─────────────────────────────────────────────────
import base64

FB_SVG = b"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 50 50" width="128" height="128">
  <circle cx="25" cy="25" r="25" fill="#1877F2"/>
  <path d="M33 25h-6v-4c0-1.7 1.1-2 1.8-2H33v-6h-5c-5.5 0-7 4.1-7 6.7V25h-4v6h4v16h7V31h4.5L33 25z" fill="white"/>
</svg>"""

IG_SVG = b"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 50 50" width="128" height="128">
  <defs>
    <linearGradient id="g" x1="0%" y1="100%" x2="100%" y2="0%">
      <stop offset="0%" style="stop-color:#F58529"/>
      <stop offset="50%" style="stop-color:#DD2A7B"/>
      <stop offset="100%" style="stop-color:#8134AF"/>
    </linearGradient>
  </defs>
  <rect width="50" height="50" rx="12" fill="url(#g)"/>
  <rect x="13" y="13" width="24" height="24" rx="6" fill="none" stroke="white" stroke-width="2.5"/>
  <circle cx="25" cy="25" r="7" fill="none" stroke="white" stroke-width="2.5"/>
  <circle cx="34" cy="16" r="1.8" fill="white"/>
</svg>"""

LI_SVG = b"""<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 50 50" width="128" height="128">
  <rect width="50" height="50" rx="8" fill="#0A66C2"/>
  <rect x="10" y="20" width="8" height="21" fill="white"/>
  <circle cx="14" cy="14" r="4.5" fill="white"/>
  <path d="M22 20h7v3c1-2 3.5-4 7-4 7 0 8 4.6 8 10.5V41h-8V31c0-2.3-.5-5.5-3.5-5.5-3.5 0-4 3-4 5.5v10H22V20z" fill="white"/>
</svg>"""

ICONS = {"Facebook": FB_SVG, "Instagram": IG_SVG, "LinkedIn": LI_SVG}

def add_svg_icon(slide, svg_bytes, x, y, size=0.35):
    """Add SVG icon to slide via PNG conversion using Pillow if available, else embed SVG."""
    try:
        from PIL import Image as PILImage
        import cairosvg
        png_data = cairosvg.svg2png(bytestring=svg_bytes, output_width=128, output_height=128)
        img_stream = io.BytesIO(png_data)
        slide.shapes.add_picture(img_stream, inches(x), inches(y), inches(size), inches(size))
    except Exception:
        # Fallback: embed as SVG (works in modern PowerPoint)
        img_stream = io.BytesIO(svg_bytes)
        try:
            slide.shapes.add_picture(img_stream, inches(x), inches(y), inches(size), inches(size))
        except Exception:
            pass  # silently skip if SVG not supported

# ── Build native chart XML helper ─────────────────────────────────────────────
def emu(inches_val):
    return int(inches_val * 914400)

def build_combo_chart_xml(months, followers, engagement, impressions,
                           col_fol, col_imp, col_eng):
    """
    Combo chart (CORRECTED axis logic):
      - BAR  (Followers)   → PRIMARY Y-axis (left, visible) — small numbers, bars always show
      - LINE (Engagement)  → SECONDARY Y-axis (right, hidden) — large numbers, own scale
      - LINE (Impressions) → SECONDARY Y-axis (right, hidden) — large numbers, own scale
    Axes:
      catAx  = 10  (bottom, used by bar)
      valAx  = 20  (left,  PRIMARY   – bar uses this, visible)
      catAx  = 40  (bottom, hidden,  used by lines)
      valAx  = 30  (right, SECONDARY – lines use this, hidden)
    """
    n = len(months)
    wb_xml = build_chart_workbook(months, followers, engagement, impressions)

    def hex3(rgb: RGBColor):
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

    month_refs = " ".join(f'<c:pt idx="{i}"><c:v>{m}</c:v></c:pt>' for i, m in enumerate(months))
    fol_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(followers))
    eng_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(engagement))
    imp_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(impressions))

    # Tight dynamic scale for lines (SFMC-style) so lines always zig-zag
    all_line_vals = list(impressions) + list(engagement)
    line_max = max(all_line_vals) if all_line_vals else 1
    line_min = min(all_line_vals) if all_line_vals else 0
    line_pad = (line_max - line_min) * 0.15 if (line_max - line_min) > 0 else line_max * 0.05
    sec_max  = line_max + line_pad
    sec_min  = max(0, line_min - line_pad)

    chart_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:date1904 val="0"/>
  <c:lang val="en-US"/>
  <c:chart>
    <c:autoTitleDeleted val="1"/>
    <c:plotArea>
      <c:layout/>

      <!-- ── BAR chart: Followers on PRIMARY axis (left, visible) ── -->
      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:varyColors val="0"/>
        <c:ser>
          <c:idx val="0"/>
          <c:order val="0"/>
          <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Followers</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{hex3(col_fol)}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_fol)}"/></a:solidFill>
              <a:ln><a:noFill/></a:ln>
            </c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="inEnd"/>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$B$2:$B${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{fol_refs}</c:numCache>
          </c:numRef></c:val>
        </c:ser>
        <!-- bar uses primary axes: catAx=10, valAx=20 -->
        <c:axId val="10"/>
        <c:axId val="20"/>
      </c:barChart>

      <!-- ── LINE chart: Engagement + Impressions on SECONDARY axis (right, hidden) ── -->
      <c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>

        <!-- Engagement -->
        <c:ser>
          <c:idx val="1"/>
          <c:order val="1"/>
          <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Engagement</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr><a:ln w="19050"><a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill></a:ln></c:spPr>
          <c:marker>
            <c:symbol val="triangle"/><c:size val="5"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill>
              <a:ln><a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill></a:ln>
            </c:spPr>
          </c:marker>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill>
              <a:ln><a:noFill/></a:ln>
            </c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="b"/>
            <c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$C$2:$C${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{eng_refs}</c:numCache>
          </c:numRef></c:val>
          <c:smooth val="0"/>
        </c:ser>

        <!-- Impressions -->
        <c:ser>
          <c:idx val="2"/>
          <c:order val="2"/>
          <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Impressions</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr><a:ln w="19050"><a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill></a:ln></c:spPr>
          <c:marker>
            <c:symbol val="circle"/><c:size val="5"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill>
              <a:ln><a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill></a:ln>
            </c:spPr>
          </c:marker>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill>
              <a:ln><a:noFill/></a:ln>
            </c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="t"/>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$D$2:$D${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{imp_refs}</c:numCache>
          </c:numRef></c:val>
          <c:smooth val="0"/>
        </c:ser>

        <!-- lines use secondary axes: catAx=40, valAx=30 -->
        <c:axId val="40"/>
        <c:axId val="30"/>
      </c:lineChart>

      <!-- ── PRIMARY category axis (bottom) – used by bar ── -->
      <c:catAx>
        <c:axId val="10"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:solidFill><a:srgbClr val="DDDDDD"/></a:solidFill></a:ln></c:spPr>
        <c:txPr><a:bodyPr rot="-2700000" vert="horz"/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="600" b="0">
          <a:solidFill><a:srgbClr val="666666"/></a:solidFill>
        </a:defRPr></a:pPr></a:p></c:txPr>
        <c:crossAx val="20"/>
        <c:auto val="1"/>
        <c:lblAlgn val="ctr"/>
        <c:noMultiLvlLbl val="1"/>
      </c:catAx>

      <!-- ── PRIMARY value axis (left) – hide numbers but keep axis alive ── -->
      <c:valAx>
        <c:axId val="20"/>
        <c:scaling>
          <c:orientation val="minMax"/>
          <c:min val="0"/>
          <c:max val="{max(followers) * 1.30 if max(followers) > 0 else 1}"/>
        </c:scaling>
        <c:delete val="0"/>
        <c:axPos val="l"/>
        <c:majorGridlines>
          <c:spPr><a:ln w="9525"><a:solidFill><a:srgbClr val="E8E8E8"/></a:solidFill></a:ln></c:spPr>
        </c:majorGridlines>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="700" b="0">
          <a:solidFill><a:srgbClr val="666666"/></a:solidFill>
        </a:defRPr></a:pPr></a:p></c:txPr>
        <c:crossAx val="10"/>
      </c:valAx>

      <!-- ── SECONDARY category axis (bottom, hidden) – used by lines ── -->
      <c:catAx>
        <c:axId val="40"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="1"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t/></a:r></a:p></c:txPr>
        <c:crossAx val="30"/>
      </c:catAx>

      <!-- ── SECONDARY value axis (right, VISIBLE) — tight dynamic range so lines zig-zag ── -->
      <c:valAx>
        <c:axId val="30"/>
        <c:scaling>
          <c:orientation val="minMax"/>
          <c:min val="{sec_min}"/>
          <c:max val="{sec_max}"/>
        </c:scaling>
        <c:delete val="0"/>
        <c:axPos val="r"/>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="700" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
        <c:crosses val="autoZero"/>
        <c:crossAx val="40"/>
      </c:valAx>

    </c:plotArea>
    <c:legend>
      <c:legendPos val="b"/>
      <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
      <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="800" b="0">
        <a:solidFill><a:srgbClr val="444444"/></a:solidFill>
      </a:defRPr></a:pPr></a:p></c:txPr>
      <c:overlay val="0"/>
    </c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </c:spPr>
  <c:externalData r:id="rId1"><c:autoUpdate val="0"/></c:externalData>
</c:chartSpace>"""
    return chart_xml.encode("utf-8"), wb_xml



def build_chart_workbook(months, followers, engagement, impressions):
    """Build a minimal xlsx workbook as bytes for the embedded chart data."""
    wb2 = openpyxl.Workbook()
    ws = wb2.active
    ws.title = "Sheet1"
    ws.append(["Month", "Followers", "Engagement", "Impressions"])
    for i, m in enumerate(months):
        ws.append([m, followers[i], engagement[i], impressions[i]])
    buf = io.BytesIO()
    wb2.save(buf)
    return buf.getvalue()


def build_fb_ig_chart_xml(months, followers, engagement, impressions,
                           col_fol, col_imp, col_eng):
    """
    FB/IG chart — SFMC pattern:
      - BAR  (Followers)              -> primary axes  10/20
      - LINE (Impressions+Engagement) -> secondary axes 11/30 (shared, dynamic tight range)
    Single barChart + single lineChart, exactly like the SFMC overview chart
    (schema-safe, no repair). Labels centered/sized like SFMC.
    """
    n = len(months)
    wb_xml = build_chart_workbook(months, followers, engagement, impressions)

    def hex3(rgb: RGBColor):
        return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

    month_refs = " ".join(f'<c:pt idx="{i}"><c:v>{m}</c:v></c:pt>' for i, m in enumerate(months))
    fol_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(followers))
    eng_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(engagement))
    imp_refs   = " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(impressions))

    bar_max = max(followers) * 1.18 if max(followers) > 0 else 1

    all_line_vals = list(impressions) + list(engagement)
    line_max = max(all_line_vals) if all_line_vals else 1
    line_min = min(all_line_vals) if all_line_vals else 0
    line_pad = (line_max - line_min) * 0.15 if (line_max - line_min) > 0 else line_max * 0.05
    sec_max  = line_max + line_pad
    sec_min  = max(0, line_min - line_pad)

    chart_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:date1904 val="0"/>
  <c:lang val="en-US"/>
  <c:chart>
    <c:autoTitleDeleted val="1"/>
    <c:plotArea>
      <c:layout/>

      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:varyColors val="0"/>
        <c:ser>
          <c:idx val="0"/><c:order val="0"/>
          <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Followers</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{hex3(col_fol)}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr><a:solidFill><a:srgbClr val="{hex3(col_fol)}"/></a:solidFill><a:ln><a:noFill/></a:ln></c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="inEnd"/>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$B$2:$B${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{fol_refs}</c:numCache>
          </c:numRef></c:val>
        </c:ser>
        <c:gapWidth val="150"/>
        <c:overlap val="0"/>
        <c:axId val="10"/>
        <c:axId val="20"/>
      </c:barChart>

      <c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>

        <c:ser>
          <c:idx val="1"/><c:order val="1"/>
          <c:tx><c:strRef><c:f>Sheet1!$D$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Impressions</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr><a:ln w="19050"><a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill></a:ln></c:spPr>
          <c:marker>
            <c:symbol val="circle"/><c:size val="5"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill>
              <a:ln><a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill></a:ln>
            </c:spPr>
          </c:marker>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr><a:solidFill><a:srgbClr val="{hex3(col_imp)}"/></a:solidFill><a:ln><a:noFill/></a:ln></c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="t"/>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$D$2:$D${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{imp_refs}</c:numCache>
          </c:numRef></c:val>
          <c:smooth val="0"/>
        </c:ser>

        <c:ser>
          <c:idx val="2"/><c:order val="2"/>
          <c:tx><c:strRef><c:f>Sheet1!$C$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Engagement</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr><a:ln w="19050"><a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill></a:ln></c:spPr>
          <c:marker>
            <c:symbol val="triangle"/><c:size val="5"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill>
              <a:ln><a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill></a:ln>
            </c:spPr>
          </c:marker>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr><a:solidFill><a:srgbClr val="{hex3(col_eng)}"/></a:solidFill><a:ln><a:noFill/></a:ln></c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr b="1" sz="460"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
            <c:dLblPos val="b"/>
            <c:showLegendKey val="0"/><c:showVal val="0"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{month_refs}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$C$2:$C${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{eng_refs}</c:numCache>
          </c:numRef></c:val>
          <c:smooth val="0"/>
        </c:ser>

        <c:marker val="1"/>
        <c:axId val="11"/>
        <c:axId val="30"/>
      </c:lineChart>

      <c:catAx>
        <c:axId val="10"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:solidFill><a:srgbClr val="DDDDDD"/></a:solidFill></a:ln></c:spPr>
        <c:txPr><a:bodyPr rot="-2700000" vert="horz"/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="600" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
        <c:crossAx val="20"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

      <c:valAx>
        <c:axId val="20"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="0"/><c:max val="{bar_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="l"/>
        <c:majorGridlines><c:spPr><a:ln w="9525"><a:solidFill><a:srgbClr val="E8E8E8"/></a:solidFill></a:ln></c:spPr></c:majorGridlines>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
        <c:crossAx val="10"/>
      </c:valAx>

      <c:valAx>
        <c:axId val="30"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="{sec_min}"/><c:max val="{sec_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="r"/>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="700" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
        <c:crossAx val="11"/>
      </c:valAx>

      <c:catAx>
        <c:axId val="11"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="1"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:r><a:t/></a:r></a:p></c:txPr>
        <c:crossAx val="30"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

    </c:plotArea>
    <c:legend>
      <c:legendPos val="b"/>
      <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
      <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr><a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="444444"/></a:solidFill></a:defRPr></a:pPr></a:p></c:txPr>
      <c:overlay val="0"/>
    </c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </c:spPr>
  <c:externalData r:id="rId1"><c:autoUpdate val="0"/></c:externalData>
</c:chartSpace>"""
    return chart_xml.encode("utf-8"), wb_xml



def insert_chart_into_slide(prs, slide, x, y, w, h, months, followers, engagement, impressions,
                             chart_type="linkedin"):
    """
    Inject a native combo chart into the slide at position (x,y) inches, size (w,h) inches.
    The chart is editable in Excel when opened in PowerPoint.
    """
    if chart_type in ("facebook", "instagram"):
        chart_xml_bytes, wb_bytes = build_fb_ig_chart_xml(
            months, followers, engagement, impressions,
            COL_FOL, COL_IMP, COL_ENG
        )
    else:
        chart_xml_bytes, wb_bytes = build_combo_chart_xml(
            months, followers, engagement, impressions,
            COL_FOL, COL_IMP, COL_ENG
        )

    # Add chart via python-pptx's GraphicFrame mechanism
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    from lxml import etree
    import uuid

    slide_part = slide.part
    chart_part = slide_part.part_related_by  # not available directly; use workaround

    # Use python-pptx's add_chart helper if available
    try:
        from pptx.util import Inches
        from pptx.chart.data import ChartData
        from pptx import chart as pptx_chart
        from pptx.enum.chart import XL_CHART_TYPE

        # We'll use a workaround: add a placeholder chart then replace its XML
        chart_data = ChartData()
        chart_data.categories = months
        chart_data.add_series("Followers",   followers)
        chart_data.add_series("Engagement",  engagement)
        chart_data.add_series("Impressions", impressions)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            inches(x), inches(y), inches(w), inches(h),
            chart_data
        )
        chart = graphic_frame.chart
        chart_part_obj = chart.part

        # Now replace the chart XML with our custom combo XML
        chart_part_obj._element = etree.fromstring(chart_xml_bytes)

        # Replace embedded workbook
        xlsx_part = chart_part_obj.part_related_by("http://schemas.openxmlformats.org/officeDocument/2006/relationships/package")
        xlsx_part._blob = wb_bytes

    except Exception as e:
        print(f"  ⚠️  Chart injection warning: {e}")
        # Fallback: simple clustered bar chart
        from pptx.chart.data import ChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_data = ChartData()
        chart_data.categories = months
        chart_data.add_series("Followers",   followers)
        chart_data.add_series("Engagement",  engagement)
        chart_data.add_series("Impressions", impressions)

        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            inches(x), inches(y), inches(w), inches(h),
            chart_data
        )
        chart = graphic_frame.chart

        # Style the series colours
        from pptx.dml.color import RGBColor as RGB
        colors = [COL_FOL, COL_ENG, COL_IMP]
        for i, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[i]

    return graphic_frame


# ═════════════════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ═════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)

blank_layout = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)

# Background
bg = add_rect(s1, 0, 0, 10, 5.625, fill_color=TEAL)

# White circle (MSO_SHAPE_TYPE oval = 9)
circle = s1.shapes.add_shape(9, inches(0.42), inches(1.82), inches(0.72), inches(0.72))
circle.fill.solid()
circle.fill.fore_color.rgb = WHITE
circle.line.fill.background()

# Number in circle — Roboto Condensed bold
add_textbox(s1, "1", 0.42, 1.82, 0.72, 0.72, size=22, bold=True, color=TEAL,
            align=PP_ALIGN.CENTER, font_face="Roboto Condensed")

# Title — Invention font for "SoMe", Invention for subtitle
add_textbox(s1, "SoMe", 1.28, 1.80, 4, 0.50, size=36, bold=True, color=WHITE,
            font_face="Invention")
add_textbox(s1, "Performance & KPIs", 1.28, 2.30, 4, 0.30, size=14, color=WHITE,
            font_face="Invention")

print("Slide 1 done")

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2 — SoMe Overview
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)

# Teal accent bar next to title
add_rect(s2, 0.28, 0.20, 0.07, 0.88, fill_color=TEAL)

add_textbox(s2, "SoMe Overview", 0.45, 0.17, 6, 0.55, size=30, bold=True, color=DARK_TEXT,
            font_face="Invention")
add_textbox(s2, "Social Media Platform Metrics", 0.45, 0.72, 6, 0.26, size=12, color=MID_GREY,
            font_face="Invention")

# Column headers
hdr_y = 1.12
for label, cx in [("Followers", 3.1), ("Impressions", 5.7), ("Engagement", 8.3)]:
    add_textbox(s2, label, cx - 0.6, hdr_y, 1.5, 0.26, size=11, bold=True,
                color=DARK_TEXT, align=PP_ALIGN.CENTER, font_face="Invention")

# Separator line
line = s2.shapes.add_shape(1, inches(0.28), inches(1.42), inches(9.4), inches(0.01))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
line.line.fill.background()

platform_rows = [
    ("Facebook",  1.52, FB_SVG),
    ("Instagram", 2.87, IG_SVG),
    ("LinkedIn",  4.22, LI_SVG),
]

for plat, icon_y, svg in platform_rows:
    d25 = data_2025[plat]
    d26 = data_2026[plat]

    add_svg_icon(s2, svg, 0.28, icon_y, size=0.36)
    add_textbox(s2, plat, 0.72, icon_y + 0.02, 2, 0.30, size=12, bold=True, color=DARK_TEXT,
                font_face="Invention")

    ry25 = icon_y + 0.44
    ry26 = ry25 + 0.28
    rypc = ry26 + 0.28

    add_textbox(s2, "2025",     0.72, ry25, 1.2, 0.25, size=11, color=MID_GREY, font_face="Invention")
    add_textbox(s2, "2026",     0.72, ry26, 1.2, 0.25, size=11, color=MID_GREY, font_face="Invention")
    add_textbox(s2, "% change", 0.72, rypc, 1.2, 0.25, size=10, italic=True, color=MID_GREY,
                font_face="Invention")

    metrics = [
        (3.1, d25["followers"],   d26["followers"]),
        (5.7, d25["impressions"], d26["impressions"]),
        (8.3, d25["engagement"],  d26["engagement"]),
    ]
    for cx, v25, v26 in metrics:
        add_textbox(s2, fmt_num(v25), cx-0.7, ry25, 1.6, 0.25, size=9,
                    color=DARK_TEXT, align=PP_ALIGN.CENTER, font_face="Invention")
        add_textbox(s2, fmt_num(v26), cx-0.7, ry26, 1.6, 0.25, size=9,
                    color=DARK_TEXT, align=PP_ALIGN.CENTER, font_face="Invention")
        pv = pct_change(v26, v25)
        pc = pct_color(v26, v25)
        add_textbox(s2, pv, cx-0.7, rypc, 1.6, 0.25, size=8, italic=True,
                    color=pc, align=PP_ALIGN.CENTER, font_face="Invention")

    # Divider
    if icon_y < 4.22:
        div = s2.shapes.add_shape(1, inches(0.28), inches(icon_y + 1.22), inches(9.4), inches(0.01))
        div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        div.line.fill.background()

print("Slide 2 done")

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3 — SoMe 2026 Breakdown (native editable charts)
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)

add_rect(s3, 0.20, 0.08, 0.06, 0.46, fill_color=TEAL)
add_textbox(s3, "SoMe 2026 Breakdown", 0.34, 0.06, 7, 0.38, size=20, bold=True, color=DARK_TEXT)
add_textbox(s3, "Social Media Platform Metrics - 2026", 0.34, 0.44, 7, 0.20, size=9, color=MID_GREY)

# Slide height = 5.625"
# Header row:  0.00 – 0.68"
# FB/IG label: 0.68"  chart: 0.90 – 3.00" (h=2.10)
# Legend sits ~0.25" below chart bottom = 3.25"
# LinkedIn label: 3.32"  chart: 3.52 – 5.52" (h=2.00) — fits within 5.625
chart_configs = [
    ("Facebook",  FB_SVG, 0.18, 1.05, 4.60, 2.00),
    ("Instagram", IG_SVG, 5.10, 1.05, 4.60, 2.00),
    ("LinkedIn",  LI_SVG, 2.55, 3.60, 4.75, 1.85),
]

for plat, svg, cx, cy, cw, ch in chart_configs:
    md = monthly[plat]
    add_svg_icon(s3, svg, cx, cy - 0.34, size=0.26)
    add_textbox(s3, plat, cx + 0.32, cy - 0.36, 2.5, 0.30, size=11, bold=True, color=DARK_TEXT)

    insert_chart_into_slide(
        prs, s3,
        cx, cy, cw, ch,
        md["months"], md["followers"], md["engagement"], md["impressions"],
        chart_type=plat.lower()
    )

print("Slide 3 done")

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4 — Top Performing Posts
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)

add_textbox(s4, "Top Performing Posts By Engagement",
            0.3, 0.18, 9.4, 0.60,
            size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font_face="Invention")

col_configs = [
    ("Instagram", IG_SVG, 0.30),
    ("Facebook",  FB_SVG, 3.55),
    ("LinkedIn",  LI_SVG, 6.80),
]

for plat, svg, px in col_configs:
    add_svg_icon(s4, svg, px, 0.95, size=0.30)
    add_textbox(s4, plat, px + 0.36, 0.95, 2.5, 0.30,
                size=18, bold=False, color=TEAL, align=PP_ALIGN.LEFT, font_face="Invention")

    # Placeholder box
    box = add_rect(s4, px, 1.38, 2.90, 3.90,
                   fill_color=RGBColor(0xF5, 0xF5, 0xF5),
                   line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=0.75)

    add_textbox(s4, "[ Paste post image here ]",
                px, 1.38, 2.90, 3.90,
                size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

print("Slide 4 done")

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT_PATH)
print(f"\nSaved: {OUTPUT_PATH}")
