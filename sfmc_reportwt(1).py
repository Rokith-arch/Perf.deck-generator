import sys
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ─────────────────────────────────────────────────────────────────
DARK_TEAL  = RGBColor(0x0D, 0x6E, 0x6E)
LIGHT_TEAL = RGBColor(0xA7, 0xD7, 0xD2)
NAVY       = RGBColor(0x1A, 0x2E, 0x44)
GREY       = RGBColor(0x6B, 0x72, 0x80)
CARD_BG    = RGBColor(0xF0, 0xF4, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INSIGHT_BG = RGBColor(0xEA, 0xF4, 0xF4)
LINE1      = RGBColor(0x1A, 0x2E, 0x44)
LINE2      = RGBColor(0x5E, 0xAD, 0x9E)
BAR_DELIVERED = RGBColor(0x15, 0x25, 0x48)
BAR_OPENS     = RGBColor(0x15, 0x85, 0x7C)
BAR_CLICKS    = RGBColor(0x92, 0xD0, 0x50)

MONTH_ORDER = ["January","February","March","April","May","June",
               "July","August","September","October","November","December"]

TA_MAP = {"Vaccines":"VACC","Oncology":"ONCO","Diabetes":"DIA","HAC":"HAC"}
TA_ORDER = ["DIA","HAC","ONCO","VACC"]

FONT = "Trebuchet MS"
AXIS_FONT = "Aptos"

# ── Embedded icon PNGs (base64) — same icons as REE report, no external assets needed ──
import os
import base64 as _base64
import tempfile as _tempfile

_ICON_B64 = {
    "delivered": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAC8UlEQVR4nO2cwY6bQBQE21E+Ofco9/wzOUQjsaywAWOYeVV1WWkPrNdd048B48c0TREuP+5+AXIvCgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8",
    "opens": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAADhklEQVR4nO2dzZKiMBRGYWoe2b3l3ndmFlZXOTaRBJLcn++cpdIY7ne8kGDrum3bArr8sR4A2IIA4iCAOAggDgKI89d6ALWsz3u46cp2e6zWYzgiRAeIGP6yxBi3ewEiFPEb3sfvWgDvxavF83G4FcBz0c7g9XjcCgBzQABxEEAcBBAnzELQHp4WWrxe5B1BBxAHAcRBAHEQQBwEECf0LOAdi6twT7OQs9ABxEEAcRBAHAQQBwHEQQBxEECcNOsAGebkFtABxJkiQNR75dbMqNtwAX4OAgnamFW3oQJ8Dh4J6phZt2EClAaNBN+ZXbchs4Cjwa7P+9b7qj3D3UCLunXvALVB0An+x6puXQVoHRwSvLCsWzcBzgyKxZsXZ+rQS4IuAhD+dawkuCwA4ffDQoJLAhB+f2ZLcFoAwh/HTAlOrQN4DD+bXNvtsZ6ZHbTWobkDeAw/KzM6QZMAhD8fsy3Ya7lvuQ8u3YPOqFLYqfgVoAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAMOvm1MurNd8sLt9++nG+VdoP0NxGAAAABJRU5ErkJggg==",
    "pct": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAD/UlEQVR4nO2dy3GbQBQFj1L+5Oxd2fuflUXKjrCFmMd939N7w0A3I4wQ3O73O0hffnkPgPjCAJrDAJrDAJrDAJrDAJrDAJrzprnw28f75UWG++8/N80xkNfcpC8EjUg/gzHYIxbAjvjvMAQ7RM4BJOVrLI+cszUDWIjibKDL8gxgdZRyNtBlKQBrKYxAj+kAvGQwAh14ISg42uFPnQSuDObVSZz08qrxuH+0tns4gFlZMwPWXHZWnu0Tje1W+QiYHWgHoTOcHRAaHwdDAcyseFWm5oyRiattk9520Rlg90juPhOMypWMQCwAKXldI5iVKhXBZQBRp9uo41phdVsk9oHIDCB91HaaBXYl7v49LwQ5InEE7x4sDMCJCPIBBuBCFPkAAzAnknyAAZgSTT7AAMyIKB/gPYEmRJUPDAQQ9X/yqOP6TmT5gOBHgNRRW+nojy4fED4H8L6qFYkM8oHBACy+qrX4ytmKLPIBpf8CvL7ZikAm+QDvCRQlm3xg4ZdBnkcr5cvDC0ECZJUPLATgNdCoR39m+cDiDGA9YMrXY/kjwGrg3jvojAryAaEHRGicGEbYOWdUkQ/wnsBpKskH+IygKarJBxQCOCy80FPCKsoHlAOoQlX5AC8EXVJZPsAAXlJdPsAATukgH2AAT+kiH2AAP+gkH2AAB7rJBxjAFx3lAwwAQF/5AANoLR9oHkB3+UDjACj/Hy0DoPz/tAuA8o+0CoDyf9ImAMp/TosAKP+c8gFQ/mtKB0D515QNgPLHKBkA5Y9TLgDKn6NUAJQ/T5kAKH+NEgFQ/jrpA6D8PVIHQPn7pA2A8mVIGQDly5EuAMqXJVUAlC9PmgAoX4cUAVC+HuEDoHxdQgdA+fq8eQ/gjE7yPZ+lFPIZQR3kR3maWrgAqsuXfKhmuTeHUr798sLMAJXlW7xjYXXbQ8wAlO+3HvcAKN93fa4BUL7/et0CqCw/Ey4ngdXlZ3q7mnkAlH9E86WcI8s2/QioLn+W2W0J+/r4ETrIt3j9rfSMYRJAB/kz7G6L5L5QD4Dyj0hti9RyVAPoJD/qC7CvxqUWQCf5o0R8u1qI18c/o5r8qIgHQPm5EA2A8vMhFgDl50QkAMrPy3YAlJ+b7QAiXdWKTsR7Al3fHl5JftRtuRqX2ElghG+2MiA1C0gtR/TfwFGpXeV/sitP8qNE/ELQldzK8jVv7lj5O7cbQs5WXFn+CrMRaHzhpHpL2OOAO8nnPYGPK/h4v3eS/4nn18Mz+1v9hpCO8jPh/sugqniFP7teBqCIdQQr62MAylhFsLqeMD8P74DGieFuYJwBDIl4TyBnACf4jCDyBZ8SRtzgOUBzGEBzGEBzGEBzGEBzGEBz/gKFAOnHjQLKDAAAAABJRU5ErkJggg==",
    "clicks": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAADlklEQVR4nO2dzVakMBBGC4+P7Ho8rsd3ZhYelGlpSCAkX/Ldu1OhW6puKgk/YZrnObaYPt83/zC/fUybO+zsk7IvtOFl65d7idz7Gwnuj18CHLXiZZuU7c58NtTllwDz28eU2pJJaP9Mz8YAEXkJXkvzuB9dgy67AnxvlCjCkujp830m6X2QJEDE+WoA2iQL8L0DU72h2JwG7rGXYJLfH9kV4L+dV9WA5PdJdgVYsySd5PfLtwBn5/Qkv2+meZ4l5+2ccq7DpgARbYJ8pgohwzWm+PtHYlp35bQyEpznUICFO4Nc4poCEpwjeRZw14WfUp/LhalzJFeANaVa2x1JoxLkUfRMYA7qFcWFl15G+0qfPxIvEek3gVBex+N1/cP6ev7jhi2Sv/WdtO6ybI4BFKZ8z/4Hblcry9NB4LpbqN36j76Prqgch7MAgj02ly4HQ/8ggDkIYA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOYggDkIYA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOa8Hm8yHqxA9oONAKkPi7qtfmrRBZx9UtjhCePhBbiaxNElGFoAViA7ZmgBSjKqBDaDwBI8vgpnhNkEApwg9dV6Efoi0AUksiyZk9sVqHcdQ1eA3FXPjra7Mp1UrQSXXhlz6gsvrhJ25TNSyvKzZI26oHV1ASKuSVBCoFxKlnE1CaS7APX+cwSkBVBhbyyRs78iTbqAiHtb990Bb9EN3UUzASLGeF/A1smhHhK/0FSACL9+Xk2O5gJE+EmwoCBDcwFck7+mpQhNTwWT/C9axkH6RJAjtatB9QpA8vepHR+uBppTVQBafxo141RNAJKfR614yV4LUJgjl0axEVSZBeQc+IiJf0QpHlKDQIfkR2gdp4wASkGpgcrx3i6AYr/XE3fHT6ICqLSG2igct4QA0A4EMAcBzEEAcxDAHAQwBwHMQQBzEMAcBDAHAcxBAHMQwBwEMAcBzEEAcxDAHAQwBwHMQQBzJARwvXNY4bhvF0DhzteesXkySKE11ETleGUEiNAJyt0oHWe1JWJyD3rErkMxBrKPhyu1kpGp1gWM2KLvpFa8qo4BkCCNmnGSGgRCfaoLQBXYp/pi16OtFt4rrRpG0y6AavBFyzg0Xyx6wbEaKDQAGQHWjCyDQtLXSAoA9WAaaA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOYggDn/AJsh0uzbsQ3aAAAAAElFTkSuQmCC",
}

def _icon_path(key):
    """Decode a base64 icon to a temp PNG file and return its path."""
    data = _base64.b64decode(_ICON_B64[key])
    tmp = _tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.write(data)
    tmp.close()
    return tmp.name

# Map KPI label → icon key (same as REE)
KPI_ICONS = {
    "Total Delivered": "delivered",
    "Total Opens":     "opens",
    "OR":              "pct",
    "Total Clicks":    "clicks",
    "CTR":             "pct",
}

def add_kpi_icon(slide, label, x, y, size):
    """Add an embedded icon to the left of a KPI block. Always available (no external file)."""
    icon_key = KPI_ICONS.get(label)
    if icon_key:
        slide.shapes.add_picture(_icon_path(icon_key), Inches(x), Inches(y), Inches(size), Inches(size))
        return True
    return False

# ── Data ───────────────────────────────────────────────────────────────────
def load_data(path):
    df = pd.read_excel(path, sheet_name=0)
    df.columns = df.columns.str.strip()
    for c in ["Month","Campaign","TA","Local/CMM"]:
        if c in df.columns:
            df[c] = df[c].astype(str).str.strip()
    for c in ["Total Sent","Total Delivered","Total Opens","Total Clicks"]:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    df = df[df["Total Delivered"] > 0].copy()
    df["or_rate"] = df["Total Opens"]  / df["Total Delivered"] * 100
    df["ctr"]     = df["Total Clicks"] / df["Total Delivered"] * 100
    df["ta_code"] = df["TA"].map(TA_MAP).fillna(df["TA"])
    months_present = [m for m in MONTH_ORDER if m in df["Month"].values]
    latest = months_present[-1] if months_present else df["Month"].iloc[0]
    return df, latest

def agg_group(df, key):
    g = df.groupby(key, as_index=False).agg(
        delivered=("Total Delivered","sum"),
        opens=("Total Opens","sum"),
        clicks=("Total Clicks","sum"),
    )
    g.rename(columns={key:"name"}, inplace=True)
    g["or_rate"] = g.apply(lambda r: r.opens/r.delivered*100 if r.delivered>0 else 0, axis=1)
    g["ctr"]     = g.apply(lambda r: r.clicks/r.delivered*100 if r.delivered>0 else 0, axis=1)
    return g

# ── Insight generator ──────────────────────────────────────────────────────
def generate_insights(camp_name, delivered, opens, clicks, or_rate, ctr, all_or, all_ctr):
    insights = []
    avg_or  = sum(all_or)  / len(all_or)  if all_or  else 0
    avg_ctr = sum(all_ctr) / len(all_ctr) if all_ctr else 0

    if or_rate > avg_or * 1.1:
        insights.append(f"Strong open rate of {or_rate:.2f}% — above campaign average of {avg_or:.2f}%.")
    elif or_rate < avg_or * 0.9:
        insights.append(f"Open rate of {or_rate:.2f}% is below average ({avg_or:.2f}%) — subject line may need review.")
    else:
        insights.append(f"Open rate of {or_rate:.2f}% is in line with the campaign average.")

    if ctr > avg_ctr * 1.1:
        insights.append(f"CTR of {ctr:.2f}% shows strong click engagement — content resonating well.")
    elif ctr < avg_ctr * 0.9:
        insights.append(f"CTR of {ctr:.2f}% is below average — CTA clarity could be improved.")
    else:
        insights.append(f"CTR of {ctr:.2f}% is consistent with average performance.")

    if delivered > 0:
        insights.append(f"{delivered:,} emails delivered with {opens:,} opens and {clicks:,} clicks.")

    return insights[:3]

# ── PPT helpers ────────────────────────────────────────────────────────────
def txt(slide, text, x, y, w, h, size=11, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = FONT
    return tb

def rect(slide, x, y, w, h, fill=CARD_BG, line_color=None, lw=Pt(0.75)):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def title_bar(slide, title, subtitle=""):
    rect(slide, 0.1, 0.13, 0.07, 0.6, fill=DARK_TEAL)
    txt(slide, title, 0.28, 0.13, 9.5, 0.5, size=22, bold=True, color=NAVY)
    if subtitle:
        txt(slide, subtitle, 0.28, 0.6, 9.5, 0.28, size=10, italic=True, color=GREY)

def kpi_block(slide, label, value, x, y, w):
    txt(slide, label, x, y,      w, 0.2,  size=8,  color=GREY,  align=PP_ALIGN.CENTER)
    txt(slide, value, x, y+0.19, w, 0.32, size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def _set_datalabels(ser_elem, font_size=900, font_name="Aptos Body", bar_color_hex="0D6E6E"):
    """Add data labels to a chart series via XML.
    Solid fill background = bar colour, white font, label inside top of bar.
    Child order MUST follow OOXML schema: numFmt, spPr, txPr, showLegendKey,
    showVal, showCatName, showSerName, showPercent, showBubbleSize, dLblPos.
    """
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    dLbls = etree.SubElement(ser_elem, qn("c:dLbls"))

    # 1. numFmt
    etree.SubElement(dLbls, qn("c:numFmt"),
                     attrib={"formatCode": "General", "sourceLinked": "0"})

    # 2. spPr — solid fill matching bar colour, no border line
    spPr = etree.SubElement(dLbls, qn("c:spPr"))
    sf = etree.SubElement(spPr, f"{{{A}}}solidFill")
    etree.SubElement(sf, f"{{{A}}}srgbClr", attrib={"val": bar_color_hex})
    ln = etree.SubElement(spPr, f"{{{A}}}ln")
    etree.SubElement(ln, f"{{{A}}}noFill")

    # 3. txPr — white font
    txPr = etree.SubElement(dLbls, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                               attrib={"sz": str(font_size), "b": "0"})
    solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
    etree.SubElement(solidFill, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
    etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": font_name})

    # 4-9. show* flags — must come AFTER txPr
    etree.SubElement(dLbls, qn("c:showLegendKey"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showVal"),         attrib={"val": "1"})
    etree.SubElement(dLbls, qn("c:showCatName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showSerName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showPercent"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showBubbleSize"),  attrib={"val": "0"})

    # 10. position — inEnd = inside top of bar
    etree.SubElement(dLbls, qn("c:dLblPos"),         attrib={"val": "inEnd"})

def _style_axis_font(chart, font_name="Aptos", font_size=8):
    """Set axis tick label fonts"""
    try:
        for axis in [chart.category_axis, chart.value_axis]:
            tf = axis.tick_labels.font
            tf.name = font_name
            tf.size = Pt(font_size)
    except Exception:
        pass



def _style_legend_font(chart, font_name="Aptos", font_size=7):
    """Set legend font size via XML."""
    try:
        legend_elem = chart._element.find(".//" + qn("c:legend"))
        if legend_elem is None:
            return
        txPr = legend_elem.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.SubElement(legend_elem, qn("c:txPr"))
        # clear existing
        for ch in list(txPr):
            txPr.remove(ch)
        etree.SubElement(txPr, qn("a:bodyPr"))
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                                   attrib={"sz": str(int(font_size * 100)), "b": "0"})
        etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": font_name})
    except Exception:
        pass


def add_bar_chart(slide, labels, series_data, x, y, w, h,
                  colors=None, show_legend=True, data_labels=True):
    if not labels or not series_data:
        return None
    cd = ChartData()
    cd.categories = labels
    for sname, svals in series_data:
        cd.add_series(sname, [float(v) for v in svals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    default_colors = [DARK_TEAL, LIGHT_TEAL, NAVY,
                      RGBColor(0x5E,0xAD,0x9E), RGBColor(0x2E,0x8B,0x8B)]
    if colors is None: colors = default_colors
    for i, ser in enumerate(chart.series):
        c = colors[i % len(colors)]
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = c
        if data_labels:
            hex_col = f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"
            _set_datalabels(ser._element, bar_color_hex=hex_col)
    chart.has_title = False
    chart.plots[0].gap_width = 80
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        _style_legend_font(chart, font_name=AXIS_FONT, font_size=9)
    else:
        chart.has_legend = False
    try:
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass
    # Light gridlines via XML — very pale grey, hairline weight
    try:
        A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        plot_area = chart._element.find(".//" + qn("c:plotArea"))
        val_ax = plot_area.find(qn("c:valAx"))
        if val_ax is not None:
            mj = val_ax.find(qn("c:majorGridlines"))
            if mj is None:
                mj = etree.SubElement(val_ax, qn("c:majorGridlines"))
            sp = mj.find(qn("c:spPr"))
            if sp is None:
                sp = etree.SubElement(mj, qn("c:spPr"))
            for old_ln in sp.findall(f"{{{A}}}ln"):
                sp.remove(old_ln)
            ln_gl = etree.SubElement(sp, f"{{{A}}}ln", attrib={"w": "6350"})
            sf_gl = etree.SubElement(ln_gl, f"{{{A}}}solidFill")
            etree.SubElement(sf_gl, f"{{{A}}}srgbClr", attrib={"val": "E8E8E8"})
    except Exception:
        pass
    _style_axis_font(chart, AXIS_FONT)
    return chart

# ── SLIDE 0: Title slide ───────────────────────────────────────────────────
def slide0(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = RGBColor(0x00, 0x85, 0x7C)
    # Circle with number
    circle = sl.shapes.add_shape(9, Inches(1.2), Inches(2.8), Inches(0.7), Inches(0.7))
    circle.fill.solid(); circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "5"
    r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = DARK_TEAL; r.font.name = FONT
    # Title
    txt(sl, "Email Performance", 2.1, 2.7, 9.0, 0.7,
        size=32, bold=True, color=WHITE)
    txt(sl, "Campaigns & KPI's", 2.1, 3.45, 9.0, 0.5,
        size=14, color=RGBColor(0xCC,0xEE,0xEE))


def _add_overview_native_chart(sl, labels, d_vals, or_vals, ct_vals):
    """
    Combo chart built using the full XML-replacement approach (same as SoMe report).
    - Delivered: clustered column bars, PRIMARY axis (left), hex #15857C, labels inside
    - OR %:      line, SECONDARY axis (right, hidden), hex #152548, solid fill labels
    - CTR %:     line, SECONDARY axis (right, hidden), hex #92D050, solid fill labels
    OR and CTR values are already in percent (e.g. 25.6), formatted as 25.6% not 25600%.
    Delivered is on the secondary (right) axis so lines sit proportionally in the chart.
    No repair prompt — chart XML and embedded workbook are replaced atomically.
    """
    import io
    import openpyxl as _openpyxl
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from lxml import etree

    n = len(labels)

    COL_DEL = "15857C"
    COL_OR  = "152548"
    COL_CTR = "92D050"

    def _pts(vals):
        return " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(vals))

    def _str_pts(vals):
        return " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(vals))

    lbl_pts   = _str_pts(labels)
    del_pts   = _pts(d_vals)
    or_pts    = _pts(or_vals)
    ctr_pts   = _pts(ct_vals)

    max_del   = max(d_vals)   if d_vals   else 1
    min_del   = min(d_vals)   if d_vals   else 0
    all_pct_vals = (or_vals or [0]) + (ct_vals or [0])
    max_pct   = max(all_pct_vals)
    min_pct   = min(all_pct_vals)
    del_max   = max_del * 1.18
    pct_pad   = (max_pct - min_pct) * 0.35 if (max_pct - min_pct) > 0 else max_pct * 0.1
    pct_max   = max_pct + pct_pad
    pct_min   = max(0, min_pct - pct_pad)

    def _line_series(idx, order, name, col, pts, ref_col_letter):
        return f"""
      <c:ser>
        <c:idx val="{idx}"/><c:order val="{order}"/>
        <c:tx><c:strRef><c:f>Sheet1!${ref_col_letter}$1</c:f>
          <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>{name}</c:v></c:pt></c:strCache>
        </c:strRef></c:tx>
        <c:spPr>
          <a:ln w="22860"><a:solidFill><a:srgbClr val="{col}"/></a:solidFill></a:ln>
        </c:spPr>
        <c:marker><c:symbol val="circle"/><c:size val="5"/>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{col}"/></a:solidFill>
            <a:ln><a:solidFill><a:srgbClr val="{col}"/></a:solidFill></a:ln>
          </c:spPr>
        </c:marker>
        <c:dLbls>
          <c:numFmt formatCode='0.00"%"' sourceLinked="0"/>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{col}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
            <a:defRPr b="1" sz="900"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr>
          </a:pPr></a:p></c:txPr>
          <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
          <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          <c:dLblPos val="t"/>
        </c:dLbls>
        <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
          <c:strCache><c:ptCount val="{n}"/>{lbl_pts}</c:strCache>
        </c:strRef></c:cat>
        <c:val><c:numRef><c:f>Sheet1!${ref_col_letter}$2:${ref_col_letter}${n+1}</c:f>
          <c:numCache><c:formatCode>0.00</c:formatCode><c:ptCount val="{n}"/>{pts}</c:numCache>
        </c:numRef></c:val>
        <c:smooth val="0"/>
      </c:ser>"""

    chart_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:date1904 val="0"/>
  <c:lang val="en-US"/>
  <c:chart>
    <c:autoTitleDeleted val="1"/>
    <c:plotArea>
      <c:layout/>

      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:varyColors val="0"/>
        <c:ser>
          <c:idx val="0"/><c:order val="0"/>
          <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Delivered</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{COL_DEL}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{COL_DEL}"/></a:solidFill>
              <a:ln><a:noFill/></a:ln>
            </c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
              <a:defRPr b="1" sz="900"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr>
            </a:pPr></a:p></c:txPr>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
            <c:dLblPos val="ctr"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{lbl_pts}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$B$2:$B${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{del_pts}</c:numCache>
          </c:numRef></c:val>
        </c:ser>
        <c:gapWidth val="60"/>
        <c:overlap val="0"/>
        <c:axId val="10"/>
        <c:axId val="20"/>
      </c:barChart>

      <c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>
        {_line_series(1, 1, "OR %",  COL_OR,  or_pts,  "C")}
        {_line_series(2, 2, "CTR %", COL_CTR, ctr_pts, "D")}
        <c:marker val="1"/>
        <c:axId val="11"/>
        <c:axId val="30"/>
      </c:lineChart>

      <c:catAx>
        <c:axId val="10"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:solidFill><a:srgbClr val="DDDDDD"/></a:solidFill></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="900" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="20"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

      <c:valAx>
        <c:axId val="20"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="0"/><c:max val="{del_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="l"/>
        <c:majorGridlines>
          <c:spPr><a:ln w="6350"><a:solidFill><a:srgbClr val="E8E8E8"/></a:solidFill></a:ln></c:spPr>
        </c:majorGridlines>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="10"/>
      </c:valAx>

      <c:valAx>
        <c:axId val="30"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="{pct_min}"/><c:max val="{pct_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="r"/>
        <c:numFmt formatCode='0.00"%"' sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="11"/>
      </c:valAx>

      <c:catAx>
        <c:axId val="11"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="1"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:crossAx val="30"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

    </c:plotArea>
    <c:legend>
      <c:legendPos val="b"/>
      <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
      <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
        <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="444444"/></a:solidFill></a:defRPr>
      </a:pPr></a:p></c:txPr>
      <c:overlay val="0"/>
    </c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </c:spPr>
  <c:externalData r:id="rId1"><c:autoUpdate val="0"/></c:externalData>
</c:chartSpace>""".encode("utf-8")

    wb2 = _openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "Sheet1"
    ws2.append(["Month", "Delivered", "OR %", "CTR %"])
    for i in range(n):
        ws2.append([labels[i], d_vals[i], or_vals[i], ct_vals[i]])
    buf = io.BytesIO()
    wb2.save(buf)
    wb_bytes = buf.getvalue()

    cd = ChartData()
    cd.categories = labels
    cd.add_series("OR %",      [float(v) for v in or_vals])
    cd.add_series("CTR %",     [float(v) for v in ct_vals])
    cd.add_series("Delivered", [float(v) for v in d_vals])

    chart_frame = sl.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.2), Inches(1.75), Inches(12.9), Inches(5.45), cd)

    chart_part = chart_frame.chart.part
    chart_part._element = etree.fromstring(chart_xml)

    xlsx_part = chart_part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package")
    xlsx_part._blob = wb_bytes

    return chart_frame


def slide1(prs, df):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE

    # Title bar — Invention font
    rect(sl, 0.1, 0.13, 0.07, 0.6, fill=DARK_TEAL)
    tb = sl.shapes.add_textbox(Inches(0.28), Inches(0.13), Inches(9.5), Inches(0.5))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "SFMC Overview"
    r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = NAVY; r.font.name = "Invention"

    tot_del   = int(df["Total Delivered"].sum())
    tot_opens = int(df["Total Opens"].sum())
    tot_clks  = int(df["Total Clicks"].sum())
    overall_or  = tot_opens/tot_del*100 if tot_del else 0
    overall_ctr = tot_clks/tot_del*100  if tot_del else 0

    kpis = [
        ("Total Delivered", f"{tot_del:,}"),
        ("Total Opens",     f"{tot_opens:,}"),
        ("OR",              f"{overall_or:.2f}%"),
        ("Total Clicks",    f"{tot_clks:,}"),
        ("CTR",             f"{overall_ctr:.2f}%"),
    ]
    cw = 13.0/len(kpis)
    icon_size = 0.32
    for i,(lbl,val) in enumerate(kpis):
        x = 0.15+i*cw
        y = 0.98
        has_icon = add_kpi_icon(sl, lbl, x+0.05, y-0.02, icon_size)
        text_x  = x + (icon_size+0.12 if has_icon else 0)
        text_w  = cw - (icon_size+0.12 if has_icon else 0)
        text_align = PP_ALIGN.LEFT if has_icon else PP_ALIGN.CENTER
        # Label — Roboto bold
        tb_lbl = sl.shapes.add_textbox(Inches(text_x), Inches(y), Inches(text_w), Inches(0.2))
        tf_lbl = tb_lbl.text_frame
        p_lbl = tf_lbl.paragraphs[0]; p_lbl.alignment = text_align
        r_lbl = p_lbl.add_run(); r_lbl.text = lbl
        r_lbl.font.size = Pt(10); r_lbl.font.bold = True
        r_lbl.font.color.rgb = GREY; r_lbl.font.name = "Roboto"
        # Value — Roboto Condensed
        tb_val = sl.shapes.add_textbox(Inches(text_x), Inches(y+0.19), Inches(text_w), Inches(0.32))
        tf_val = tb_val.text_frame
        p_val = tf_val.paragraphs[0]; p_val.alignment = text_align
        r_val = p_val.add_run(); r_val.text = val
        r_val.font.size = Pt(20); r_val.font.bold = True
        r_val.font.color.rgb = NAVY; r_val.font.name = "Roboto Condensed"

    # Build per-month data for the native chart
    by_month = df.groupby("Month", as_index=False).agg(
        delivered=("Total Delivered","sum"),
        opens=("Total Opens","sum"),
        clicks=("Total Clicks","sum"),
    )
    by_month["or_rate"] = by_month.apply(
        lambda r: r.opens/r.delivered*100 if r.delivered>0 else 0, axis=1)
    by_month["ctr"]     = by_month.apply(
        lambda r: r.clicks/r.delivered*100 if r.delivered>0 else 0, axis=1)
    by_month["Month"]   = pd.Categorical(by_month["Month"],
                                          categories=MONTH_ORDER, ordered=True)
    by_month = by_month.sort_values("Month")

    labels  = [m[:3] for m in by_month["Month"]]
    d_vals  = [int(v)         for v in by_month["delivered"]]
    or_vals = [round(float(v),2) for v in by_month["or_rate"]]
    ct_vals = [round(float(v),2) for v in by_month["ctr"]]

    _add_overview_native_chart(sl, labels, d_vals, or_vals, ct_vals)

# ── SLIDE 2: SFMC Overview TAs ─────────────────────────────────────────────
def slide2(prs, df):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, "SFMC Overview – TAs")

    positions = [(0.15,0.85),(6.75,0.85),(0.15,4.05),(6.75,4.05)]
    cw, ch = 6.3, 3.0

    for i, ta_code in enumerate(TA_ORDER):
        x, y = positions[i]
        # Find actual TA name in data
        ta_name = next((k for k,v in TA_MAP.items() if v==ta_code), ta_code)
        ta_df = df[df["TA"]==ta_name] if ta_name in df["TA"].values else pd.DataFrame()

        td  = int(ta_df["Total Delivered"].sum()) if len(ta_df)>0 else 0
        to  = int(ta_df["Total Opens"].sum())     if len(ta_df)>0 else 0
        tc  = int(ta_df["Total Clicks"].sum())    if len(ta_df)>0 else 0
        tor = to/td*100 if td>0 else 0
        tctr= tc/td*100 if td>0 else 0
        camps = ta_df["Campaign"].nunique() if len(ta_df)>0 else 0

        rect(sl, x, y, cw, ch, fill=CARD_BG, line_color=RGBColor(0xD0,0xE8,0xE6))
        txt(sl, ta_code, x+0.1, y+0.06, 1.5, 0.28, size=13, bold=True, color=DARK_TEAL)
        txt(sl, f"{camps} Campaign{'s' if camps!=1 else ''}",
            x+1.65, y+0.07, 3.5, 0.22, size=10, color=GREY)

        kpis = [
            ("Total Delivered", f"{td:,}"),
            ("Total Opens",     f"{to:,}"),
            ("OR",              f"{tor:.2f}%"),
            ("Total Clicks",    f"{tc:,}"),
            ("CTR",             f"{tctr:.2f}%"),
        ]
        kw = cw/len(kpis)
        icon_sz  = 0.20   # same size as REE mini KPI icons
        icon_gap = 0.04   # gap between icon and label text
        for ki,(lbl,val) in enumerate(kpis):
            icon_key = KPI_ICONS.get(lbl)
            kx = x + ki * kw
            row_t = y + 0.32
            if icon_key:
                # Icon — same placement pattern as REE slide2 mini KPI row
                sl.shapes.add_picture(
                    _icon_path(icon_key),
                    Inches(kx),
                    Inches(row_t + (0.22 - icon_sz) / 2),
                    Inches(icon_sz), Inches(icon_sz))
                lbl_x   = kx + icon_sz + icon_gap
                lbl_w   = kw - icon_sz - icon_gap
                lbl_align = PP_ALIGN.LEFT
            else:
                lbl_x   = kx
                lbl_w   = kw
                lbl_align = PP_ALIGN.CENTER
            # Label — Roboto
            tb_lbl = sl.shapes.add_textbox(Inches(lbl_x), Inches(row_t), Inches(lbl_w), Inches(0.18))
            tf_lbl = tb_lbl.text_frame
            p_lbl = tf_lbl.paragraphs[0]; p_lbl.alignment = lbl_align
            r_lbl = p_lbl.add_run(); r_lbl.text = lbl
            r_lbl.font.size = Pt(6); r_lbl.font.bold = False
            r_lbl.font.color.rgb = GREY; r_lbl.font.name = "Roboto"
            # Value — Roboto Condensed (aligned under label, past the icon)
            tb_val = sl.shapes.add_textbox(Inches(lbl_x), Inches(y+0.48), Inches(lbl_w), Inches(0.24))
            tf_val = tb_val.text_frame
            p_val = tf_val.paragraphs[0]; p_val.alignment = lbl_align
            r_val = p_val.add_run(); r_val.text = val
            r_val.font.size = Pt(10); r_val.font.bold = True
            r_val.font.color.rgb = NAVY; r_val.font.name = "Roboto Condensed"

        if len(ta_df)>0 and td>0:
            camp_agg = agg_group(ta_df, "Campaign").sort_values("or_rate", ascending=False).head(5)
            labels = [n[:14]+"…" if len(n)>14 else n for n in camp_agg["name"]]
            add_bar_chart(sl, labels,
                [("Delivered",[int(v) for v in camp_agg["delivered"]]),
                 ("Opens",    [int(v) for v in camp_agg["opens"]]),
                 ("Clicks",   [int(v) for v in camp_agg["clicks"]])],
                x+0.05, y+0.76, cw-0.1, ch-0.82,
                colors=[BAR_DELIVERED, BAR_OPENS, BAR_CLICKS], data_labels=True)
        else:
            rect(sl, x+0.05, y+0.76, cw-0.1, ch-0.82,
                 fill=RGBColor(0xF0,0xF4,0xF4), line_color=RGBColor(0xD0,0xE8,0xE6))
            txt(sl, "No data available", x+0.05, y+0.76+(ch-0.82)/2-0.12,
                cw-0.1, 0.25, size=9, color=GREY, align=PP_ALIGN.CENTER, italic=True)

# ── SLIDE 3 & 6: Asset cards ───────────────────────────────────────────────
def slide_assets(prs, df, latest_month, sort_by="or_rate", subtitle="Open Rate"):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, "Best Performing Campaigns – SFMC – Assets", f"Best by {subtitle} ({latest_month})")

    latest    = df[df["Month"]==latest_month]
    camp_agg  = agg_group(latest, "Campaign").sort_values(sort_by, ascending=False).head(3)
    cw, ch, sy = 3.85, 4.1, 1.0
    xs = [0.15, 4.57, 9.0] if len(camp_agg)==3 else [0.15+i*4.5 for i in range(len(camp_agg))]

    for i, row in enumerate(camp_agg.itertuples()):
        x = xs[i]
        ta = df[df["Campaign"]==row.name]["TA"].iloc[0] if len(df[df["Campaign"]==row.name])>0 else ""
        ta_code = TA_MAP.get(ta, ta)
        rect(sl, x, sy, cw, ch, fill=CARD_BG, line_color=RGBColor(0xD0,0xE8,0xE6))
        rect(sl, x+0.1, sy+0.1, cw-0.2, 2.4,
             fill=RGBColor(0xE8,0xF4,0xF3), line_color=RGBColor(0xD0,0xE8,0xE6))
        txt(sl, "[ Email Asset ]", x+0.1, sy+1.05, cw-0.2, 0.3,
            size=9, color=GREY, align=PP_ALIGN.CENTER, italic=True)
        val = f"OR – {row.or_rate:.2f}%" if sort_by=="or_rate" else f"CTR – {row.ctr:.2f}%"
        txt(sl, val, x+0.1, sy+2.62, cw-0.2, 0.38, size=13, bold=True, color=DARK_TEAL, align=PP_ALIGN.CENTER)
        txt(sl, ta_code.upper(), x+0.1, sy+3.02, cw-0.2, 0.22, size=9, color=GREY, align=PP_ALIGN.CENTER)
        txt(sl, row.name.upper(), x+0.1, sy+3.24, cw-0.2, 0.38, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDE 10 & 11: Best Performing Campaigns – Overall (all months) ────────
def slide_overall_assets(prs, df, sort_by="or_rate", subtitle="Open Rate"):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, "Best Performing Campaigns – SFMC – Assets", f"Best by {subtitle} (Overall)")

    camp_agg = agg_group(df, "Campaign").sort_values(sort_by, ascending=False).head(3)

    if len(camp_agg)==0:
        txt(sl, "No campaign data available.", 0.5, 2.0, 9, 0.5, size=12, color=GREY)
        return

    cw, ch, sy = 3.85, 4.3, 0.95
    xs = [0.15, 4.57, 9.0] if len(camp_agg)==3 else [0.15+i*4.5 for i in range(len(camp_agg))]

    for i, row in enumerate(camp_agg.itertuples()):
        x = xs[i]
        ta = df[df["Campaign"]==row.name]["TA"].iloc[0] if len(df[df["Campaign"]==row.name])>0 else ""
        ta_code = TA_MAP.get(ta, ta)
        rect(sl, x, sy, cw, ch, fill=CARD_BG, line_color=RGBColor(0xD0,0xE8,0xE6))
        # Bigger email-asset placeholder for poster image
        rect(sl, x+0.1, sy+0.1, cw-0.2, 3.2,
             fill=RGBColor(0xE8,0xF4,0xF3), line_color=RGBColor(0xD0,0xE8,0xE6))
        txt(sl, "[ Email Asset ]", x+0.1, sy+1.5, cw-0.2, 0.3,
            size=9, color=GREY, align=PP_ALIGN.CENTER, italic=True)
        val = f"OR – {row.or_rate:.2f}%" if sort_by=="or_rate" else f"CTR – {row.ctr:.2f}%"
        txt(sl, val, x+0.1, sy+3.35, cw-0.2, 0.32, size=13, bold=True, color=DARK_TEAL, align=PP_ALIGN.CENTER)
        txt(sl, ta_code.upper(), x+0.1, sy+3.70, cw-0.2, 0.20, size=9, color=GREY, align=PP_ALIGN.CENTER)
        txt(sl, row.name.upper(), x+0.1, sy+3.92, cw-0.2, 0.38, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── SLIDES 4,5,7,8: Campaign charts ───────────────────────────────────────
def slide_camp_charts(prs, df, latest_month, sort_by="or_rate",
                      subtitle="Open Rate", local_only=False):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title = ("Best Performing Local Campaigns – SFMC" if local_only
             else "Best Performing Campaigns – SFMC")
    title_bar(sl, title, f"Best by {subtitle} ({latest_month})")

    latest = df[df["Month"]==latest_month].copy()
    if local_only:
        latest = latest[latest["Local/CMM"]=="Local"]
    if len(latest)==0:
        txt(sl, "No data for this filter.", 0.5, 2.0, 9, 0.5, size=12, color=GREY)
        return

    camp_agg = agg_group(latest, "Campaign").sort_values(sort_by, ascending=False).head(3)
    all_or  = camp_agg["or_rate"].tolist()
    all_ctr = camp_agg["ctr"].tolist()
    count   = len(camp_agg)
    cw = (13.1/count) - 0.2
    xs = [0.15 + i*(cw+0.2) for i in range(count)]
    sy, ch = 0.9, 6.3

    for i, row in enumerate(camp_agg.itertuples()):
        x = xs[i]
        ta = df[df["Campaign"]==row.name]["TA"].iloc[0] if len(df[df["Campaign"]==row.name])>0 else ""
        ta_code = TA_MAP.get(ta, ta)

        rect(sl, x, sy, cw, ch, fill=CARD_BG, line_color=RGBColor(0xD0,0xE8,0xE6))
        txt(sl, ta_code.upper(), x+0.1, sy+0.12, cw-0.2, 0.22,
            size=8, color=GREY, align=PP_ALIGN.CENTER)
        txt(sl, row.name, x+0.1, sy+0.35, cw-0.2, 0.38,
            size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # Chart: latest month only — Delivered, Opens, Clicks with new colors
        camp_latest = df[(df["Campaign"]==row.name) & (df["Month"]==latest_month)]
        delivered = int(camp_latest["Total Delivered"].sum())
        opens     = int(camp_latest["Total Opens"].sum())
        clicks    = int(camp_latest["Total Clicks"].sum())
        add_bar_chart(sl,
            [latest_month[:3]],
            [("Delivered",[delivered]),
             ("Opens",    [opens]),
             ("Clicks",   [clicks])],
            x+0.07, sy+0.82, cw-0.14, 2.6,
            colors=[BAR_DELIVERED, BAR_OPENS, BAR_CLICKS], data_labels=True)

        # KPIs — pushed down to match chart bottom
        ky = sy + 3.55
        kw = (cw-0.2)/3
        kpis = [("Total Delivered", f"{int(row.delivered):,}"),
                ("OR",              f"{row.or_rate:.2f}%"),
                ("CTR",             f"{row.ctr:.2f}%")]
        for ki,(lbl,val) in enumerate(kpis):
            kx = x+0.1+ki*kw
            # Label — Roboto
            tb_lbl = sl.shapes.add_textbox(Inches(kx), Inches(ky), Inches(kw), Inches(0.2))
            tf_lbl = tb_lbl.text_frame
            p_lbl = tf_lbl.paragraphs[0]; p_lbl.alignment = PP_ALIGN.CENTER
            r_lbl = p_lbl.add_run(); r_lbl.text = lbl
            r_lbl.font.size = Pt(7); r_lbl.font.bold = False
            r_lbl.font.color.rgb = GREY; r_lbl.font.name = "Roboto"
            # Value — Roboto Condensed
            tb_val = sl.shapes.add_textbox(Inches(kx), Inches(ky+0.2), Inches(kw), Inches(0.32))
            tf_val = tb_val.text_frame
            p_val = tf_val.paragraphs[0]; p_val.alignment = PP_ALIGN.CENTER
            r_val = p_val.add_run(); r_val.text = val
            r_val.font.size = Pt(14); r_val.font.bold = True
            r_val.font.color.rgb = NAVY; r_val.font.name = "Roboto Condensed"

        # Insights — tightened up
        iy = sy + 4.22
        insights = generate_insights(
            row.name, int(row.delivered), int(row.opens), int(row.clicks),
            row.or_rate, row.ctr, all_or, all_ctr)
        rect(sl, x+0.07, iy, cw-0.14, 1.98,
             fill=INSIGHT_BG, line_color=RGBColor(0xD0,0xE8,0xE6))
        txt(sl, "INSIGHTS", x+0.07, iy+0.06, cw-0.14, 0.2,
            size=7, bold=True, color=GREY, align=PP_ALIGN.CENTER)
        for ii, ins in enumerate(insights):
            txt(sl, f"• {ins}", x+0.12, iy+0.28+ii*0.55, cw-0.24, 0.5,
                size=7, color=NAVY)

# ── EXCEL ──────────────────────────────────────────────────────────────────
def build_excel(df, latest_month, output_path):
    wb = openpyxl.Workbook()
    H_FILL = PatternFill("solid", fgColor="0D6E6E")
    S_FILL = PatternFill("solid", fgColor="1A4A4A")
    ALT    = PatternFill("solid", fgColor="EAF4F4")
    WH     = PatternFill("solid", fgColor="FFFFFF")
    H_FONT = Font(name=FONT, bold=True, color="FFFFFF", size=11)
    N_FONT = Font(name=FONT, size=10)
    CTR_A  = Alignment(horizontal="center", vertical="center")
    LFT_A  = Alignment(horizontal="left",   vertical="center")
    thin   = Side(style="thin", color="C0D8D8")
    BDR    = Border(left=thin, right=thin, top=thin, bottom=thin)

    def write_table(ws, start_row, section_title, headers, rows, col_widths=None):
        ws.merge_cells(start_row=start_row, start_column=1,
                       end_row=start_row, end_column=len(headers))
        c = ws.cell(start_row, 1, section_title)
        c.font = Font(name=FONT, bold=True, color="FFFFFF", size=10)
        c.fill = S_FILL; c.alignment = LFT_A
        for ci in range(2, len(headers)+1):
            ws.cell(start_row, ci).fill = S_FILL
        hr = start_row+1
        for ci, h in enumerate(headers, 1):
            cell = ws.cell(hr, ci, h)
            cell.font = H_FONT; cell.fill = H_FILL
            cell.alignment = CTR_A; cell.border = BDR
        for ri, row in enumerate(rows):
            dr = hr+1+ri
            for ci, val in enumerate(row, 1):
                cell = ws.cell(dr, ci, val)
                cell.font = N_FONT
                cell.fill = ALT if ri%2==1 else WH
                cell.alignment = CTR_A; cell.border = BDR
        if col_widths:
            for ci, cw in enumerate(col_widths, 1):
                ws.column_dimensions[get_column_letter(ci)].width = cw
        return hr+1+len(rows)

    ws = wb.active; ws.title = "Dashboard"
    ws.sheet_view.showGridLines = False
    ws.merge_cells("A1:J2")
    c = ws["A1"]
    c.value = "MSD SFMC Email Campaign Performance Report"
    c.font = Font(name=FONT, bold=True, color="FFFFFF", size=15)
    c.fill = PatternFill("solid", fgColor="0D6E6E")
    c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 28; ws.row_dimensions[2].height = 8

    tot_del = int(df["Total Delivered"].sum())
    tot_op  = int(df["Total Opens"].sum())
    tot_clk = int(df["Total Clicks"].sum())
    gor  = tot_op/tot_del*100  if tot_del else 0
    gctr = tot_clk/tot_del*100 if tot_del else 0
    kpis_xl = [
        ("Total Delivered",f"{tot_del:,}","0D6E6E"),
        ("Total Opens",    f"{tot_op:,}", "1A7B6B"),
        ("OR",             f"{gor:.2f}%", "2E8B8B"),
        ("Total Clicks",   f"{tot_clk:,}","5EAD9E"),
        ("CTR",            f"{gctr:.2f}%","A7D7D2"),
    ]
    kpi_cols = [1,3,5,7,9]
    for idx,(lbl,val,col) in enumerate(kpis_xl):
        c1 = ws.cell(4, kpi_cols[idx], lbl)
        c1.font = Font(name=FONT, bold=True, color="FFFFFF", size=9)
        c1.fill = PatternFill("solid", fgColor=col); c1.alignment = CTR_A
        ws.merge_cells(start_row=4, start_column=kpi_cols[idx], end_row=4, end_column=kpi_cols[idx]+1)
        ws.cell(4, kpi_cols[idx]+1).fill = PatternFill("solid", fgColor=col)
        c2 = ws.cell(5, kpi_cols[idx], val)
        c2.font = Font(name=FONT, bold=True, color=col, size=14)
        c2.fill = WH; c2.alignment = CTR_A
        ws.merge_cells(start_row=5, start_column=kpi_cols[idx], end_row=5, end_column=kpi_cols[idx]+1)
        ws.cell(5, kpi_cols[idx]+1).fill = WH
    ws.row_dimensions[4].height = 20; ws.row_dimensions[5].height = 26

    by_month = df.groupby("Month", as_index=False).agg(
        delivered=("Total Delivered","sum"),
        opens=("Total Opens","sum"),
        clicks=("Total Clicks","sum"))
    by_month["or_rate"] = by_month.apply(lambda r: r.opens/r.delivered*100 if r.delivered>0 else 0, axis=1)
    by_month["ctr"]     = by_month.apply(lambda r: r.clicks/r.delivered*100 if r.delivered>0 else 0, axis=1)
    by_month["Month"]   = pd.Categorical(by_month["Month"], categories=MONTH_ORDER, ordered=True)
    by_month = by_month.sort_values("Month")
    m_rows = [[r.Month, f"{int(r.delivered):,}", f"{int(r.opens):,}", f"{int(r.clicks):,}",
               f"{r.or_rate:.2f}%", f"{r.ctr:.2f}%"] for r in by_month.itertuples()]
    next_r = write_table(ws, 7, "Monthly Performance Summary",
        ["Month","Total Delivered","Total Opens","Total Clicks","OR","CTR"],
        m_rows, [14,17,14,14,10,10])

    camp_agg = agg_group(df, "Campaign").sort_values("or_rate", ascending=False)
    c_rows = [[r.name, f"{int(r.delivered):,}", f"{int(r.opens):,}", f"{int(r.clicks):,}",
               f"{r.or_rate:.2f}%", f"{r.ctr:.2f}%"] for r in camp_agg.itertuples()]
    next_r = write_table(ws, next_r+1, "Campaign Performance Summary",
        ["Campaign","Total Delivered","Total Opens","Total Clicks","OR","CTR"],
        c_rows, [35,17,14,14,10,10])

    ta_agg = agg_group(df, "TA").sort_values("or_rate", ascending=False)
    ta_rows = [[r.name, f"{int(r.delivered):,}", f"{int(r.opens):,}", f"{int(r.clicks):,}",
                f"{r.or_rate:.2f}%", f"{r.ctr:.2f}%"] for r in ta_agg.itertuples()]
    write_table(ws, next_r+1, "TA Performance Summary",
        ["TA","Total Delivered","Total Opens","Total Clicks","OR","CTR"],
        ta_rows, [20,17,14,14,10,10])

    ws2 = wb.create_sheet(f"Latest – {latest_month}")
    ws2.sheet_view.showGridLines = False
    latest = df[df["Month"]==latest_month]
    lat_agg = agg_group(latest, "Campaign").sort_values("or_rate", ascending=False)
    lr = [[r.name,
           df[df["Campaign"]==r.name]["TA"].iloc[0] if len(df[df["Campaign"]==r.name])>0 else "",
           df[df["Campaign"]==r.name]["Local/CMM"].iloc[0] if len(df[df["Campaign"]==r.name])>0 else "",
           f"{int(r.delivered):,}", f"{int(r.opens):,}", f"{int(r.clicks):,}",
           f"{r.or_rate:.2f}%", f"{r.ctr:.2f}%"] for r in lat_agg.itertuples()]
    write_table(ws2, 1, f"{latest_month} – All Campaigns",
        ["Campaign","TA","Local/CMM","Total Delivered","Total Opens","Total Clicks","OR","CTR"],
        lr, [35,15,12,17,14,14,10,10])

    ws3 = wb.create_sheet("Raw Data")
    cols = [c for c in ["Month","Campaign","TA","Local/CMM",
            "Total Sent","Total Delivered","Total Opens","Total Clicks"] if c in df.columns]
    for ci, col in enumerate(cols, 1):
        cell = ws3.cell(1, ci, col)
        cell.font = H_FONT; cell.fill = H_FILL
        cell.alignment = CTR_A; cell.border = BDR
        ws3.column_dimensions[get_column_letter(ci)].width = max(14, len(col)+2)
    for ri, row in enumerate(df[cols].itertuples(index=False), 2):
        for ci, val in enumerate(row, 1):
            cell = ws3.cell(ri, ci, val)
            cell.font = N_FONT
            cell.fill = ALT if ri%2==0 else WH
            cell.alignment = CTR_A; cell.border = BDR

    # ── Native Excel Charts (live — update when you change numbers) ────────
    # These charts are linked to the data tables above, so editing a cell
    # in the Monthly/Campaign/TA tables will instantly refresh the chart.
    from openpyxl.chart import BarChart, LineChart, Reference, Series
    from openpyxl.chart.series import SeriesLabel

    ws_chart = wb.create_sheet("Charts")
    ws_chart.sheet_view.showGridLines = False

    # ── 1. Monthly Delivered + OR trend (linked to Dashboard table) ─────
    # Find where the monthly table starts in ws (row 8 = header, data from row 9)
    # by_month was sorted and written to rows 9 onwards in write_table call
    # Header row = 8, data starts row 9
    m_count = len(m_rows)
    m_hdr   = 8   # header row of monthly table
    m_start = 9   # first data row
    m_end   = m_start + m_count - 1

    # Bar chart — Total Delivered (col B = column 2)
    bar_m = BarChart()
    bar_m.type    = "col"
    bar_m.grouping = "clustered"
    bar_m.title   = "Monthly Email Delivered"
    bar_m.y_axis.title = "Emails Delivered"
    bar_m.x_axis.title = "Month"
    bar_m.shape   = 4
    bar_m.width   = 20
    bar_m.height  = 12

    data_del = Reference(ws, min_col=2, max_col=2,
                         min_row=m_hdr, max_row=m_end)
    cats_m   = Reference(ws, min_col=1, min_row=m_start, max_row=m_end)
    bar_m.add_data(data_del, titles_from_data=True)
    bar_m.set_categories(cats_m)
    bar_m.series[0].graphicalProperties.solidFill = "0D6E6E"
    bar_m.series[0].graphicalProperties.line.solidFill = "0D6E6E"

    # Line chart — OR % (col E = column 5) — secondary axis
    line_or = LineChart()
    line_or.grouping = "standard"
    data_or = Reference(ws, min_col=5, max_col=5,
                        min_row=m_hdr, max_row=m_end)
    line_or.add_data(data_or, titles_from_data=True)
    line_or.series[0].graphicalProperties.line.solidFill = "1A2E44"
    line_or.series[0].graphicalProperties.line.width = 20000
    line_or.y_axis.axId = 200
    line_or.y_axis.crosses = "max"

    bar_m += line_or
    ws_chart.add_chart(bar_m, "A1")

    # ── 2. Monthly CTR trend bar chart ──────────────────────────────────
    bar_ctr = BarChart()
    bar_ctr.type      = "col"
    bar_ctr.grouping  = "clustered"
    bar_ctr.title     = "Monthly Click-Through Rate (%)"
    bar_ctr.y_axis.title = "CTR %"
    bar_ctr.x_axis.title = "Month"
    bar_ctr.width     = 20
    bar_ctr.height    = 12

    data_ctr = Reference(ws, min_col=6, max_col=6,
                         min_row=m_hdr, max_row=m_end)
    bar_ctr.add_data(data_ctr, titles_from_data=True)
    bar_ctr.set_categories(cats_m)
    bar_ctr.series[0].graphicalProperties.solidFill = "5EAD9E"
    bar_ctr.series[0].graphicalProperties.line.solidFill = "5EAD9E"
    ws_chart.add_chart(bar_ctr, "A25")

    # ── 3. Campaign OR bar chart (horizontal) ───────────────────────────
    # Campaign table starts at next_r_camp_start (recalculate from known offset)
    # monthly table: header=8, data=9..9+m_count-1, section title row before = 7
    # next_r after monthly = 9+m_count  → campaign section title at 9+m_count,
    # header at 9+m_count+1, data from 9+m_count+2
    c_hdr   = 9 + m_count + 1   # campaign header row
    c_start = c_hdr + 1
    c_end   = c_start + len(c_rows) - 1

    bar_camp = BarChart()
    bar_camp.type      = "bar"   # horizontal
    bar_camp.grouping  = "clustered"
    bar_camp.title     = "Campaign Open Rate (%)"
    bar_camp.x_axis.title = "OR %"
    bar_camp.width     = 20
    bar_camp.height    = max(10, len(c_rows) * 0.6 + 4)

    data_camp_or = Reference(ws, min_col=5, max_col=5,
                             min_row=c_hdr, max_row=c_end)
    cats_camp    = Reference(ws, min_col=1, min_row=c_start, max_row=c_end)
    bar_camp.add_data(data_camp_or, titles_from_data=True)
    bar_camp.set_categories(cats_camp)
    bar_camp.series[0].graphicalProperties.solidFill = "0D6E6E"
    bar_camp.series[0].graphicalProperties.line.solidFill = "0D6E6E"
    ws_chart.add_chart(bar_camp, "L1")

    # ── 4. TA delivered + OR grouped bar chart ──────────────────────────
    ta_hdr   = c_end + 2 + 1   # TA section title + header
    ta_start = ta_hdr + 1
    ta_end   = ta_start + len(ta_rows) - 1

    bar_ta = BarChart()
    bar_ta.type      = "col"
    bar_ta.grouping  = "clustered"
    bar_ta.title     = "TA Performance – Delivered, Opens & Clicks"
    bar_ta.y_axis.title = "Volume"
    bar_ta.x_axis.title = "Therapeutic Area"
    bar_ta.width     = 20
    bar_ta.height    = 12

    data_ta = Reference(ws, min_col=2, max_col=4,
                        min_row=ta_hdr, max_row=ta_end)
    cats_ta  = Reference(ws, min_col=1, min_row=ta_start, max_row=ta_end)
    bar_ta.add_data(data_ta, titles_from_data=True)
    bar_ta.set_categories(cats_ta)
    ta_fill_colors = ["A7D7D2", "0D6E6E", "1A2E44"]
    for si, col in enumerate(ta_fill_colors):
        if si < len(bar_ta.series):
            bar_ta.series[si].graphicalProperties.solidFill = col
            bar_ta.series[si].graphicalProperties.line.solidFill = col
    ws_chart.add_chart(bar_ta, "L25")

    wb.save(output_path)

# ── MAIN ──────────────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("Usage: py -3.10 sfmc_report.py <input.xlsx>")
        print("Outputs: SFMC_Report.pptx and SFMC_Report.xlsx")
        sys.exit(1)
    inp     = sys.argv[1]
    ppt_out = sys.argv[2] if len(sys.argv)>2 else "SFMC_Report.pptx"
    xls_out = ppt_out.replace(".pptx", ".xlsx")

    print("📂 Reading data...")
    df, latest = load_data(inp)
    print(f"✅ {len(df)} rows | Latest month: {latest}")

    print("🔨 Building PPT...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide0(prs);                                                      print("  ✓ Slide 1: Title")
    slide1(prs, df);                                                  print("  ✓ Slide 2: SFMC Overview")
    slide2(prs, df);                                                  print("  ✓ Slide 3: SFMC TAs")
    slide_assets(prs, df, latest, "or_rate", "Open Rate");            print("  ✓ Slide 3: Best OR Assets")
    slide_camp_charts(prs, df, latest, "or_rate", "Open Rate");       print("  ✓ Slide 4: Best OR Charts")
    slide_camp_charts(prs, df, latest, "or_rate", "Open Rate", True); print("  ✓ Slide 5: Best Local OR")
    slide_assets(prs, df, latest, "ctr", "Click Rate");               print("  ✓ Slide 6: Best CTR Assets")
    slide_camp_charts(prs, df, latest, "ctr", "Click Rate");          print("  ✓ Slide 7: Best CTR Charts")
    slide_camp_charts(prs, df, latest, "ctr", "Click Rate", True);    print("  ✓ Slide 8: Best Local CTR")
    slide_overall_assets(prs, df, "or_rate", "Open Rate");            print("  ✓ Slide 10: Best OR Assets (Overall)")
    slide_overall_assets(prs, df, "ctr", "Click Rate");               print("  ✓ Slide 11: Best CTR Assets (Overall)")

    prs.save(ppt_out)
    print(f"  ✅ Saved: {ppt_out}")

    print("📊 Building Excel...")
    build_excel(df, latest, xls_out)
    print(f"  ✅ Saved: {xls_out}")
    print(f"\n🎉 Done!")

if __name__ == "__main__":
    main()
