#!/usr/bin/env python3
"""
GCC Pulse Automation Script
Usage: py -3.10 gcc_pulse_automation.py <input.xlsx>
Outputs: <input>_GCC_Pulse_Report.pptx and <input>_GCC_Pulse_Report.xlsx
"""
import sys, os
import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree

# ── Colors ──────────────────────────────────────────────────────────────────
TEAL   = RGBColor(0x00, 0x85, 0x7C)
TEAL_D = RGBColor(0x00, 0x66, 0x66)
TEAL_L = RGBColor(0xE0, 0xF5, 0xF5)
NAVY   = RGBColor(0x1B, 0x2A, 0x4A)
GREY   = RGBColor(0x64, 0x74, 0x8B)
CARD   = RGBColor(0xF4, 0xF6, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Calibri"

# ── Native chart helpers ─────────────────────────────────────────────────────
def _set_datalabels(ser_elem, font_size=700, font_name="Calibri",
                    fill_rgb=None, font_color_rgb=None, label_position=None):
    """Add data labels. Child order follows OOXML schema exactly to avoid repair prompt.
    label_position: None=default, 'inEnd'=inside end, 'ctr'=center, 'outEnd'=outside end, 'inBase'=inside base.
    """
    dLbls = etree.SubElement(ser_elem, qn("c:dLbls"))
    etree.SubElement(dLbls, qn("c:numFmt"),
                     attrib={"formatCode": "General", "sourceLinked": "0"})
    # Solid fill for individual data label background (applied via extLst per-label or spPr on dLbls)
    if fill_rgb is not None:
        spPr = etree.SubElement(dLbls, qn("c:spPr"))
        solidFill = etree.SubElement(etree.SubElement(spPr, qn("a:solidFill")), qn("a:srgbClr"))
        solidFill.attrib["val"] = fill_rgb
    txPr = etree.SubElement(dLbls, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"), attrib={"algn": "ctr"})
    defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                               attrib={"sz": str(font_size), "b": "0"})
    # solidFill MUST come before a:latin in OOXML schema or PowerPoint ignores the color
    if font_color_rgb is not None:
        solidFillFont = etree.SubElement(defRPr, qn("a:solidFill"))
        etree.SubElement(solidFillFont, qn("a:srgbClr")).attrib["val"] = font_color_rgb
    etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": font_name})
    etree.SubElement(dLbls, qn("c:showLegendKey"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showVal"),         attrib={"val": "1"})
    etree.SubElement(dLbls, qn("c:showCatName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showSerName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showPercent"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showBubbleSize"),  attrib={"val": "0"})
    if label_position is not None:
        etree.SubElement(dLbls, qn("c:dLblPos"), attrib={"val": label_position})
    return dLbls


def _patch_low_bar_labels(dLbls, values, bar_fill_hex, font_size=700,
                           font_name="Calibri", threshold_pct=0.08):
    """For bars below threshold_pct of max value, add individual c:dLbl overrides
    with a solid fill pill (bar color) and white font so the label is always readable.
    Must be called AFTER _set_datalabels so dLbls already exists in the XML.
    """
    if not values:
        return
    max_val = max(v for v in values if v is not None and v > 0) if any(v and v > 0 for v in values) else 1
    threshold = max_val * threshold_pct

    for idx, val in enumerate(values):
        if val is None:
            continue
        if val < threshold:
            dLbl = etree.SubElement(dLbls, qn("c:dLbl"))
            etree.SubElement(dLbl, qn("c:idx"), attrib={"val": str(idx)})
            # Override position to outside end so the pill sits above tiny bar
            etree.SubElement(dLbl, qn("c:dLblPos"), attrib={"val": "outEnd"})
            # Solid fill background on label
            spPr = etree.SubElement(dLbl, qn("c:spPr"))
            sf = etree.SubElement(spPr, qn("a:solidFill"))
            etree.SubElement(sf, qn("a:srgbClr")).attrib["val"] = bar_fill_hex
            ln = etree.SubElement(spPr, qn("a:ln"))
            etree.SubElement(ln, qn("a:noFill"))
            # White font
            txPr = etree.SubElement(dLbl, qn("c:txPr"))
            etree.SubElement(txPr, qn("a:bodyPr"))
            etree.SubElement(txPr, qn("a:lstStyle"))
            lp = etree.SubElement(txPr, qn("a:p"))
            lpPr = etree.SubElement(lp, qn("a:pPr"), attrib={"algn": "ctr"})
            rpr = etree.SubElement(lpPr, qn("a:defRPr"),
                                   attrib={"sz": str(font_size), "b": "0"})
            wf = etree.SubElement(rpr, qn("a:solidFill"))
            etree.SubElement(wf, qn("a:srgbClr")).attrib["val"] = "FFFFFF"
            etree.SubElement(rpr, qn("a:latin"), attrib={"typeface": font_name})
            etree.SubElement(dLbl, qn("c:showLegendKey"), attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showVal"),        attrib={"val": "1"})
            etree.SubElement(dLbl, qn("c:showCatName"),    attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showSerName"),    attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showPercent"),    attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showBubbleSize"), attrib={"val": "0"})

def _always_fill_labels(dLbls, values, fill_hex, font_size=700, font_name="Calibri"):
    """Unconditionally add c:dLbl overrides for EVERY data point with a solid fill pill
    (fill_hex color) and white font, positioned above the bar (outEnd).
    Use when the series has low/variable bars and you always want a visible label.
    """
    for idx, val in enumerate(values):
        dLbl = etree.SubElement(dLbls, qn("c:dLbl"))
        etree.SubElement(dLbl, qn("c:idx"), attrib={"val": str(idx)})
        etree.SubElement(dLbl, qn("c:dLblPos"), attrib={"val": "outEnd"})
        # Solid fill background
        spPr = etree.SubElement(dLbl, qn("c:spPr"))
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr")).attrib["val"] = fill_hex
        ln = etree.SubElement(spPr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
        # White font
        txPr = etree.SubElement(dLbl, qn("c:txPr"))
        etree.SubElement(txPr, qn("a:bodyPr"))
        etree.SubElement(txPr, qn("a:lstStyle"))
        lp = etree.SubElement(txPr, qn("a:p"))
        lpPr = etree.SubElement(lp, qn("a:pPr"), attrib={"algn": "ctr"})
        rpr = etree.SubElement(lpPr, qn("a:defRPr"),
                               attrib={"sz": str(font_size), "b": "0"})
        wf = etree.SubElement(rpr, qn("a:solidFill"))
        etree.SubElement(wf, qn("a:srgbClr")).attrib["val"] = "FFFFFF"
        etree.SubElement(rpr, qn("a:latin"), attrib={"typeface": font_name})
        etree.SubElement(dLbl, qn("c:showLegendKey"), attrib={"val": "0"})
        etree.SubElement(dLbl, qn("c:showVal"),        attrib={"val": "1"})
        etree.SubElement(dLbl, qn("c:showCatName"),    attrib={"val": "0"})
        etree.SubElement(dLbl, qn("c:showSerName"),    attrib={"val": "0"})
        etree.SubElement(dLbl, qn("c:showPercent"),    attrib={"val": "0"})
        etree.SubElement(dLbl, qn("c:showBubbleSize"), attrib={"val": "0"})


def _style_axis_font(chart, font_name="Calibri", font_size=8,
                     cat_font_name=None, cat_font_size=None, cat_bold=False, cat_color=None):
    try:
        # Value axis
        tf = chart.value_axis.tick_labels.font
        tf.name = font_name
        tf.size = Pt(font_size)
    except Exception:
        pass
    try:
        # Category axis — use overrides if provided
        cf = chart.category_axis.tick_labels.font
        cf.name = cat_font_name if cat_font_name else font_name
        cf.size = Pt(cat_font_size if cat_font_size else font_size)
        cf.bold = cat_bold
        if cat_color:
            cf.color.rgb = cat_color
    except Exception:
        pass

def _add_bar_chart(slide, labels, series_data, x, y, w, h,
                   colors=None, show_legend=True, data_labels=True,
                   label_fill_rgb=None, label_font_color=None, label_font_name=None,
                   cat_font_name=None, cat_font_size=None, cat_bold=False, cat_color=None,
                   gridline_color=None, show_gridlines=True,
                   label_positions=None, label_font_sizes=None,
                   legend_font_name=None, legend_font_size=None,
                   low_bar_patch_series=None,
                   always_fill_series=None):
    """Add a native clustered column chart. No matplotlib, fully editable in PPT.
    label_positions: list of position strings per series (e.g. ['inEnd','outEnd'])
    label_font_sizes: list of font sizes (in 1/100 pt) per series
    gridline_color: hex string e.g. 'D0E8E6' for faded gridlines
    show_gridlines: set False to remove gridlines entirely
    legend_font_name/legend_font_size: font for legend text
    low_bar_patch_series: list of series indices to apply low-bar label patching on
    always_fill_series: dict of {series_index: hex_color} to unconditionally fill all labels
    """
    if not labels or not series_data:
        return None
    cd = ChartData()
    cd.categories = labels
    for sname, svals in series_data:
        cd.add_series(sname, [float(v) for v in svals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    default_colors = [RGBColor(0x00,0x96,0x88), RGBColor(0x1B,0x2A,0x4A),
                      RGBColor(0x00,0x80,0x80), RGBColor(0xE0,0xF5,0xF5)]
    if colors is None:
        colors = default_colors
    for i, ser in enumerate(chart.series):
        c = colors[i % len(colors)]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = c
        if data_labels:
            pos = label_positions[i] if label_positions and i < len(label_positions) else None
            fsz = label_font_sizes[i] if label_font_sizes and i < len(label_font_sizes) else 700
            dLbls = _set_datalabels(ser._element,
                            font_size=fsz,
                            font_name=label_font_name if label_font_name else FONT,
                            fill_rgb=label_fill_rgb,
                            font_color_rgb=label_font_color,
                            label_position=pos)
            # Patch low bars: add per-label solid fill pill so label stays visible
            if low_bar_patch_series and i in low_bar_patch_series:
                bar_hex = str(c).upper()  # RGBColor.__str__ returns 6-char hex e.g. '009688'
                vals = [v for v in series_data[i][1]]
                _patch_low_bar_labels(dLbls, vals, bar_hex, font_size=fsz,
                                      font_name=label_font_name if label_font_name else FONT)
            # Always-fill: unconditionally put every label in solid fill pill above bar
            if always_fill_series and i in always_fill_series:
                fill_hex = always_fill_series[i]
                vals = [v for v in series_data[i][1]]
                _always_fill_labels(dLbls, vals, fill_hex, font_size=fsz,
                                    font_name=label_font_name if label_font_name else FONT)
    chart.has_title = False
    chart.plots[0].gap_width = 70
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        if legend_font_name or legend_font_size:
            try:
                leg_elem = chart.legend._element
                txPr = etree.SubElement(leg_elem, qn("c:txPr"))
                etree.SubElement(txPr, qn("a:bodyPr"))
                etree.SubElement(txPr, qn("a:lstStyle"))
                lp = etree.SubElement(txPr, qn("a:p"))
                lpPr = etree.SubElement(lp, qn("a:pPr"))
                rpr_attribs = {}
                if legend_font_size:
                    rpr_attribs["sz"] = str(int(legend_font_size * 100))
                defRPr = etree.SubElement(lpPr, qn("a:defRPr"), attrib=rpr_attribs)
                if legend_font_name:
                    etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": legend_font_name})
            except Exception:
                pass
    else:
        chart.has_legend = False
    try:
        if show_gridlines:
            chart.value_axis.has_major_gridlines = True
            if gridline_color:
                gridline_elem = chart.value_axis._element.find(qn("c:majorGridlines"))
                if gridline_elem is None:
                    gridline_elem = etree.SubElement(chart.value_axis._element, qn("c:majorGridlines"))
                spPr = gridline_elem.find(qn("c:spPr"))
                if spPr is None:
                    spPr = etree.SubElement(gridline_elem, qn("c:spPr"))
                ln = spPr.find(qn("a:ln"))
                if ln is None:
                    ln = etree.SubElement(spPr, qn("a:ln"))
                solidFill = etree.SubElement(ln, qn("a:solidFill"))
                etree.SubElement(solidFill, qn("a:srgbClr")).attrib["val"] = gridline_color
        else:
            chart.value_axis.has_major_gridlines = False
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass
    _style_axis_font(chart,
                     cat_font_name=cat_font_name,
                     cat_font_size=cat_font_size,
                     cat_bold=cat_bold,
                     cat_color=cat_color)
    return chart

MONTH_ORDER = ["January","February","March","April","May","June",
               "July","August","September","October","November","December"]

TA_MAP = {
    "diabetes":"DIA","dia":"DIA",
    "hac":"HAC",
    "onco":"ONCO","oncology":"ONCO",
    "vaccines":"VACC","vaccine":"VACC","vacc":"VACC",
}
TA_ORDER = ["DIA","HAC","ONCO","VACC"]

# ── Data loading ─────────────────────────────────────────────────────────────
def load_data(path):
    df = pd.read_excel(path, sheet_name=0)
    df.columns = [str(c).strip() for c in df.columns]
    print(f"   -> Raw columns: {list(df.columns)}")

    col_map, done = {}, set()
    for c in df.columns:
        lc = str(c).lower().replace(" ","").replace("_","").replace("/","")
        t = None
        if "month" in lc and "Month" not in done:                                        t="Month"
        elif ("landingpage" in lc or lc=="landingpage") and "Landing_Page" not in done:  t="Landing_Page"
        elif ("session" in lc
              and "manual"  not in lc
              and "source"  not in lc
              and "medium"  not in lc
              and "channel" not in lc
              and "Sessions" not in done):                                                 t="Sessions"
        elif "activeuser" in lc and "Active_Users" not in done:                          t="Active_Users"
        elif "newuser" in lc and "New_Users" not in done:                                t="New_Users"
        elif ("engagementtime" in lc or "engtime" in lc or "engagement" in lc) and "Eng_Time" not in done: t="Eng_Time"
        elif lc=="ta" and "TA" not in done:                                              t="TA"
        if t:
            col_map[c] = t
            done.add(t)

    df.rename(columns=col_map, inplace=True)
    df = df.loc[:, ~df.columns.duplicated()]
    print(f"   -> Mapped columns: {list(df.columns)}")

    if "TA" not in df.columns:
        df["TA"] = np.nan

    def map_ta(x):
        s = str(x).strip()
        if s in ("nan","","None"): return np.nan
        return TA_MAP.get(s.lower(), s).upper()

    df["TA"] = df["TA"].map(map_ta)

    for col in ["Sessions","Active_Users","New_Users","Eng_Time"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    df["Returning_Users"] = (df.get("Active_Users", 0) - df.get("New_Users", 0)).clip(lower=0)

    months_present = [m for m in MONTH_ORDER if m in df["Month"].values]
    latest = months_present[-1] if months_present else df["Month"].iloc[0]
    return df, latest

# ── PPT helpers ──────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=CARD, line=None):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_txt(slide, text, x, y, w, h, size=11, bold=False, color=NAVY,
            align=PP_ALIGN.LEFT, italic=False, font_name=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font_name if font_name else FONT
    return tb

def title_bar(slide, title):
    add_rect(slide, 0.1, 0.13, 0.07, 0.55, fill=TEAL)
    add_txt(slide, title, 0.28, 0.13, 12.8, 0.55, size=24, bold=True, color=NAVY)

def kpi_block(slide, label, value, x, y, w, icon_img_path=None,
              label_font="Roboto", value_font="Calibri"):
    add_txt(slide, label, x, y,      w, 0.28, size=10, color=GREY,
            align=PP_ALIGN.CENTER, font_name=label_font)
    add_txt(slide, value, x, y+0.27, w, 0.38, size=20, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, font_name=value_font)

# ── SLIDE 1: Title ───────────────────────────────────────────────────────────
def slide_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = TEAL

    circle = sl.shapes.add_shape(9, Inches(0.9), Inches(2.75), Inches(0.85), Inches(0.85))
    circle.fill.solid(); circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "2"
    r.font.size = Pt(22); r.font.bold = True
    r.font.color.rgb = TEAL; r.font.name = FONT

    add_txt(sl, "GCC Pulse",        2.0, 2.72, 9, 0.6, size=36, bold=True, color=WHITE)
    add_txt(sl, "Performance KPI's",2.0, 3.35, 9, 0.4, size=16, color=WHITE)

# ── SLIDE 2: Overview ────────────────────────────────────────────────────────

def slide_overview(prs, df):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, "GCC Pulse  Overview")

    # Year pill
    pill = sl.shapes.add_shape(1, Inches(5.5), Inches(0.18), Inches(1.3), Inches(0.38))
    pill.fill.solid(); pill.fill.fore_color.rgb = NAVY; pill.line.fill.background()
    tf = pill.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "2026"
    r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT

    # KPI card
    add_rect(sl, 0.2, 0.88, 12.9, 1.35, fill=CARD, line=RGBColor(0xD0,0xE8,0xE6))
    tot_users    = int(df["Active_Users"].sum())
    tot_sessions = int(df["Sessions"].sum())
    avg_time     = round(float(df["Eng_Time"].mean()), 2)
    kw = 12.9/3

    # KPI 1 - Active Users
    kx1 = 0.2
    kpi_block(sl, "Active Users",               f"{tot_users:,}",    kx1,       0.93, kw)

    # KPI 2 - Total Sessions
    kx2 = 0.2 + kw
    kpi_block(sl, "Total Sessions",             f"{tot_sessions:,}", kx2,       0.93, kw)

    # KPI 3 - Avg Time
    kx3 = 0.2 + kw*2
    kpi_block(sl, "Avg Time Spent on Page (s)", f"{avg_time}s",      kx3,       0.93, kw)

    # Bar chart by TA
    add_rect(sl, 0.2, 2.4, 12.9, 4.8, fill=CARD, line=RGBColor(0xD0,0xE8,0xE6))
    add_txt(sl, "Therapeutic Area Active Users Breakdown by Year",
            0.2, 2.48, 12.9, 0.3, size=11, color=GREY, align=PP_ALIGN.CENTER)

    # Native bar chart — TA Active Users breakdown (editable in PPT)
    ta_grp = df[df["TA"].notna()].groupby("TA")["Active_Users"].sum().reindex(TA_ORDER, fill_value=0)
    _add_bar_chart(sl,
        list(ta_grp.index),
        [("Active Users", [int(v) for v in ta_grp.values])],
        0.3, 2.75, 12.7, 4.3,
        colors=[RGBColor(0x00,0x96,0x88)],
        show_legend=False, data_labels=True,
        label_fill_rgb=None,
        label_font_color="FFFFFF",       # white numbers inside bar
        label_font_name="Calibri",
        cat_font_name="Calibri",
        cat_font_size=13,                # bigger category axis labels
        cat_bold=True,
        cat_color=RGBColor(0x1B, 0x2A, 0x4A),
        gridline_color="D8E8E7",         # faded gridlines
        show_gridlines=True,
        label_positions=["inEnd"],       # inside end of bar so white is visible on teal
        label_font_sizes=[1000],         # bigger font so numbers are readable
        low_bar_patch_series=[0])        # solid fill pill for short bars

# ── SLIDE 3: Overview TAs ────────────────────────────────────────────────────

def slide_ta_overview(prs, df):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, "GCC Pulse  Overview - TAs")

    all_months = [m for m in MONTH_ORDER if m in df["Month"].values]
    positions  = [(0.15, 0.85), (6.75, 0.85), (0.15, 4.1), (6.75, 4.1)]
    cw, ch     = 6.3, 3.1

    for i, ta in enumerate(TA_ORDER):
        x, y   = positions[i]
        ta_df  = df[df["TA"] == ta]
        active  = int(ta_df["Active_Users"].sum())
        sessions= int(ta_df["Sessions"].sum())
        avg_t   = round(float(ta_df["Eng_Time"].mean()), 2) if len(ta_df) > 0 else 0

        add_rect(sl, x, y, cw, ch, fill=CARD, line=RGBColor(0xD0,0xE8,0xE6))
        add_txt(sl, ta, x+0.1, y+0.06, cw-0.2, 0.3, size=13, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)

        kw3 = (cw-0.2)/3
        for ki,(lbl,val) in enumerate([
            ("Total Active Users", f"{active:,}"),
            ("Total Sessions",     f"{sessions:,}"),
            ("Avg Time on Page (s)", f"{avg_t}s"),
        ]):
            kx = x+0.1+ki*kw3
            add_txt(sl, lbl, kx, y+0.38, kw3, 0.2,  size=7.5, color=GREY,
                    align=PP_ALIGN.CENTER, font_name="Roboto")
            add_txt(sl, val, kx, y+0.57, kw3, 0.34, size=13, bold=True, color=NAVY,
                    align=PP_ALIGN.CENTER, font_name="Roboto Condensed")

        if len(ta_df) > 0:
            by_m = ta_df.groupby("Month", as_index=False).agg(
                sessions=("Sessions","sum"),
                active_sum=("Active_Users","sum"),
                new_sum=("New_Users","sum"))
            by_m["returning"] = (by_m["active_sum"] - by_m["new_sum"]).clip(lower=0).astype(int)
            by_m["Month"] = pd.Categorical(by_m["Month"], categories=MONTH_ORDER, ordered=True)
            by_m = by_m.sort_values("Month")
            m_labels = [m[:3] for m in by_m["Month"]]
            s_vals   = [int(v) for v in by_m["sessions"]]
            r_vals   = [int(v) for v in by_m["returning"]]
            _add_bar_chart(sl, m_labels,
                [("Returning Users", r_vals), ("Sessions", s_vals)],
                x+0.05, y+0.98, cw-0.1, ch-1.05,
                colors=[RGBColor(0x15,0x25,0x48), RGBColor(0x15,0x85,0x7C)],
                show_legend=True, data_labels=True,
                label_font_color="FFFFFF",
                label_font_name="Calibri",
                show_gridlines=False,
                label_positions=["outEnd", "inEnd"],
                label_font_sizes=[750, 750],
                cat_font_name="Calibri",
                cat_font_size=9,
                cat_bold=True,
                legend_font_name="Calibri",
                legend_font_size=9,
                always_fill_series={0: "152548"})  # Returning Users: always show label as navy pill above bar

# ── SLIDE 4: Most Engaging Pages ─────────────────────────────────────────────
def slide_engaging(prs, df, latest_month):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = WHITE
    title_bar(sl, f"GCC Pulse  Most Engaging Pages")

    # Exclude ANY landing page containing 'reference' or 'therapeutic'
    # anywhere — with or without a leading slash, then sort by Sessions
    # (the mapped column, which already excludes manual/source/medium).
    _lp = df["Landing_Page"].astype(str)
    latest = df[
        (df["Month"] == latest_month) &
        (df["TA"].notna()) &
        (~_lp.str.contains("reference",   case=False, na=False)) &
        (~_lp.str.contains("therapeutic", case=False, na=False))
    ].copy()

    positions = [(0.15, 0.85), (6.75, 0.85), (0.15, 4.1), (6.75, 4.1)]
    cw, ch = 6.3, 3.1

    for i, ta in enumerate(TA_ORDER):
        x, y  = positions[i]
        ta_df = latest[latest["TA"] == ta]

        # Aggregate and rank strictly by Sessions (not manual source/medium)
        grp = ta_df.groupby("Landing_Page", as_index=False).agg(
            active=("Active_Users","sum"),
            sessions=("Sessions","sum"),
            avg_eng=("Eng_Time","mean"),
        ).sort_values("sessions", ascending=False).head(3)

        add_rect(sl, x, y, cw, ch, fill=CARD, line=RGBColor(0xD0,0xE8,0xE6))

        # TA badge
        badge = sl.shapes.add_shape(1, Inches(x+0.1), Inches(y+0.08),
                                    Inches(0.7), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        tf = badge.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ta
        r.font.size = Pt(9); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT

        # Table headers
        hdrs = ["Page","Active Users","Sessions","Avg Eng (s)","Source"]
        cws  = [2.5, 0.85, 0.8, 0.82, 0.8]
        cx   = x+0.05
        for hi,(h,hw) in enumerate(zip(hdrs,cws)):
            add_txt(sl, h, cx, y+0.44, hw, 0.22, size=9, bold=True,
                    color=GREY, align=PP_ALIGN.LEFT if hi==0 else PP_ALIGN.CENTER)
            cx += hw

        # Divider line
        line = sl.shapes.add_shape(1, Inches(x+0.05), Inches(y+0.68),
                                   Inches(cw-0.1), Inches(0.01))
        line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCB,0xD5,0xE1)
        line.line.fill.background()

        # Rows
        for ri, row in enumerate(grp.itertuples()):
            ry = y + 0.72 + ri*0.62
            bg = WHITE if ri % 2 == 0 else TEAL_L
            add_rect(sl, x+0.05, ry-0.03, cw-0.1, 0.58, fill=bg)
            page_name = str(row.Landing_Page).lstrip("/")[:30]
            vals = [page_name, str(int(row.active)), str(int(row.sessions)),
                    f"{row.avg_eng:.2f}", "Organic"]
            cx = x+0.05
            for vi,(v,hw) in enumerate(zip(vals,cws)):
                add_txt(sl, v, cx, ry, hw, 0.52, size=8.5 if vi==0 else 9, color=NAVY,
                        align=PP_ALIGN.LEFT if vi==0 else PP_ALIGN.CENTER)
                cx += hw

        # Insight
        if len(grp) > 0:
            top = grp.iloc[0]
            insight = (f"{ta} led by {str(top['Landing_Page']).lstrip('/')[:25]} "
                       f"with {int(top['active'])} active users & {int(top['sessions'])} sessions. "
                       f"All traffic organic.")
            add_rect(sl, x+0.05, y+2.58, cw-0.1, 0.45, fill=RGBColor(0xB8,0xD4,0xE8),
                     line=RGBColor(0x7A,0xAF,0xD4))
            add_txt(sl, insight, x+0.12, y+2.61, cw-0.24, 0.4,
                    size=6.5, color=GREY, italic=True)

# ── EXCEL REPORT ─────────────────────────────────────────────────────────────
def build_excel(df, latest_month, out_path):
    wb = openpyxl.Workbook()
    H_FILL = PatternFill("solid", fgColor="008080")
    S_FILL = PatternFill("solid", fgColor="006666")
    ALT    = PatternFill("solid", fgColor="E0F5F5")
    WH     = PatternFill("solid", fgColor="FFFFFF")
    H_FONT = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    N_FONT = Font(name="Calibri", size=10, color="1B2A4A")
    CTR    = Alignment(horizontal="center", vertical="center", wrap_text=True)
    LFT    = Alignment(horizontal="left",   vertical="center", wrap_text=True)
    thin   = Side(style="thin", color="C0D8D8")
    BDR    = Border(left=thin, right=thin, top=thin, bottom=thin)

    def write_section(ws, row, title, headers, rows, widths=None):
        ws.merge_cells(start_row=row, start_column=1,
                       end_row=row, end_column=len(headers))
        c = ws.cell(row, 1, title)
        c.font = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
        c.fill = S_FILL; c.alignment = LFT
        for ci in range(2, len(headers)+1):
            ws.cell(row, ci).fill = S_FILL
        ws.row_dimensions[row].height = 24
        hr = row+1
        for ci,h in enumerate(headers,1):
            cell = ws.cell(hr,ci,h)
            cell.font=H_FONT; cell.fill=H_FILL
            cell.alignment=CTR; cell.border=BDR
        ws.row_dimensions[hr].height = 22
        for ri,r in enumerate(rows):
            dr = hr+1+ri
            for ci,val in enumerate(r,1):
                cell = ws.cell(dr,ci,val)
                cell.font=N_FONT
                cell.fill = ALT if ri%2==1 else WH
                cell.alignment = LFT if ci==1 else CTR
                cell.border=BDR
            ws.row_dimensions[dr].height = 20
        if widths:
            for ci,w in enumerate(widths,1):
                ws.column_dimensions[get_column_letter(ci)].width = w
        return hr+1+len(rows)

    # Sheet 1: Overview
    ws1 = wb.active; ws1.title = "Overview"
    ws1.sheet_view.showGridLines = False
    ws1.merge_cells("A1:F2")
    c = ws1["A1"]
    c.value = "GCC Pulse - Overview KPIs"
    c.font = Font(name="Calibri", bold=True, color="FFFFFF", size=16)
    c.fill = H_FILL; c.alignment = CTR
    ws1.row_dimensions[1].height = 36; ws1.row_dimensions[2].height = 8

    tot_u = int(df["Active_Users"].sum())
    tot_s = int(df["Sessions"].sum())
    avg_t = round(float(df["Eng_Time"].mean()),2)
    for ri,(lbl,val) in enumerate([
        ("Active Users", f"{tot_u:,}"),
        ("Total Sessions", f"{tot_s:,}"),
        ("Avg Time on Page (s)", f"{avg_t}s"),
    ],3):
        ws1.row_dimensions[ri].height = 30
        lc = ws1.cell(ri,1,lbl); vc = ws1.cell(ri,2,val)
        lc.font=Font(name="Calibri",bold=True,color="1B2A4A",size=12)
        vc.font=Font(name="Calibri",bold=True,color="008080",size=14)
        lc.alignment=LFT; vc.alignment=CTR
        fill = ALT if ri%2==0 else WH
        lc.fill=PatternFill("solid",fgColor="E0F5F5" if ri%2==0 else "FFFFFF")
        vc.fill=PatternFill("solid",fgColor="E0F5F5" if ri%2==0 else "FFFFFF")
        lc.border=BDR; vc.border=BDR

    # TA breakdown
    ta_grp = df[df["TA"].notna()].groupby("TA",as_index=False).agg(
        Active_Users=("Active_Users","sum"),
        Sessions=("Sessions","sum"),
        Avg_Time=("Eng_Time","mean"),
    )
    ta_rows = [[r.TA, f"{int(r.Active_Users):,}", f"{int(r.Sessions):,}", f"{r.Avg_Time:.2f}s"]
               for r in ta_grp.itertuples()]
    write_section(ws1, 7, "Active Users by Therapeutic Area",
                  ["TA","Active Users","Total Sessions","Avg Time (s)"],
                  ta_rows, [15,20,20,18])

    # Sheet 2: Monthly
    ws2 = wb.create_sheet("Monthly TA Breakdown")
    ws2.sheet_view.showGridLines = False
    mta = df[df["TA"].notna()].groupby(["Month","TA"],as_index=False).agg(
        Sessions=("Sessions","sum"),
        Returning=("Returning_Users","sum"),
    )
    m_rows = [[r.Month, r.TA, f"{int(r.Sessions):,}", f"{int(r.Returning):,}"]
              for r in mta.itertuples()]
    write_section(ws2, 1, "Monthly TA Breakdown",
                  ["Month","TA","Sessions","Returning Users"],
                  m_rows, [16,12,18,20])

    # Sheet 3: Most Engaging Pages
    ws3 = wb.create_sheet("Most Engaging Pages")
    ws3.sheet_view.showGridLines = False
    _lp_ex = df["Landing_Page"].astype(str)
    latest = df[
        (df["Month"]==latest_month) & df["TA"].notna() &
        (~_lp_ex.str.contains("reference",   case=False, na=False)) &
        (~_lp_ex.str.contains("therapeutic", case=False, na=False))
    ]
    row_cursor = 1
    for ta in TA_ORDER:
        sub = latest[latest["TA"]==ta].groupby("Landing_Page",as_index=False).agg(
            active=("Active_Users","sum"),
            sessions=("Sessions","sum"),
            avg_eng=("Eng_Time","mean"),
        ).sort_values("sessions",ascending=False).head(3)
        ep_rows = [[r.Landing_Page, int(r.active), int(r.sessions), round(r.avg_eng,2), "Organic"]
                   for r in sub.itertuples()]
        row_cursor = write_section(ws3, row_cursor, f"{ta} - Most Engaging Pages ({latest_month})",
            ["Landing Page","Active Users","Sessions","Avg Eng Time (s)","Source"],
            ep_rows, [45,16,14,20,14]) + 1

    wb.save(out_path)
    print(f"   Excel saved: {out_path}")

# ── MAIN ─────────────────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print("Usage: py -3.10 gcc_pulse_automation.py <input.xlsx>")
        sys.exit(1)

    inp  = sys.argv[1]
    base = os.path.splitext(os.path.basename(inp))[0]
    out_dir  = os.path.dirname(os.path.abspath(inp))
    pptx_out = os.path.join(out_dir, f"{base}_GCC_Pulse_Report.pptx")
    xlsx_out = os.path.join(out_dir, f"{base}_GCC_Pulse_Report.xlsx")

    print(f"Loading: {inp}")
    df, latest = load_data(inp)
    print(f"   -> {len(df)} rows | Months: {[m for m in MONTH_ORDER if m in df['Month'].values]}")
    print(f"   -> TAs found: {df['TA'].dropna().unique().tolist()}")
    print(f"   -> Latest month: {latest}")

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_title(prs);              print("  Slide 1: Title")
    slide_overview(prs, df);       print("  Slide 2: Overview")
    slide_ta_overview(prs, df);    print("  Slide 3: Overview TAs")
    slide_engaging(prs, df, latest); print("  Slide 4: Most Engaging Pages")

    prs.save(pptx_out)
    print(f"  PowerPoint saved: {pptx_out}")

    print("Building Excel...")
    build_excel(df, latest, xlsx_out)

    print(f"\nDone! Files in: {out_dir}")

if __name__ == "__main__":
    main()
