"""
REE Report Generator - Generates PPT and Excel from raw data Excel file.
Usage: python generate_ree_report.py <input_excel_file>
"""

import sys
import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
import warnings
warnings.filterwarnings('ignore')
import base64 as _base64
import tempfile as _tempfile

# ── Embedded icon PNGs (base64) — no external assets folder needed ─────────
_ICON_B64 = {
    "delivered": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAC8UlEQVR4nO2cwY6bQBQE21E+Ofco9/wzOUQjsaywAWOYeVV1WWkPrNdd048B48c0TREuP+5+AXIvCgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoABwFgKMAcBQAjgLAUQA4CgBHAeAoAJzNAjz+/vbTowOxNa9NArSDKcEY7Mnr59GDT7/+PPa/NPkkRxbo49VzAa768Xm2WJ+OAMOvwbMc3QXAWRXA1V+LtTxXBfAkrxZreT4dAUpQg8MngVIfBYCjAHB2Xwl8B88ptnPVLuzSBnBruY0r36dLGyDxXsIz7lggpzTAkTBtg68ceT/OWESnjQAlOM5d4Scnj4D2ovb8Q+SRcGfwjY+cBNoGr+kh/OSDuwAlWKeX8JMP7wIcCV/pKfjGJdcBbIM+w08uvBBElqDX8JObLgVTRkLPwTduuRlEaIMRwk9uvBtYWYJRwk9uuBcwp9pIGCn4RhefB6jQBiOGn3QiQDK2BKOGn9w8ApaMNhJGDr7RTQPMGaENKoSfdCpA0rcEVcJPOhsBS3obCZWCb3TbAHN6aIOK4SeDCJDcK0HV8JPOR8CSq0dC5eAbwzTAnCvCJISfDCpA8lkJKOEng42AJWePBFLwjWEbYM4ZbUAMPykiQPKeBNTwk8FHwJJ3RsLev1GFMg0w51MhVQs/KSpAcn5YFcNPio2AJUdGwtoxqlK2AeYcDbF6+AlEgGQ9zL2/r0bpEbBky0igBN/ANMCcFvLaTxJIAZLvYRPDT8ACyH8UAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBYCjAHAUAI4CwFEAOAoARwHgKAAcBYCjAHBOeTawl69tl/3YAHAUAI4CwHkpAPWp2Sq8ys8GgLNJAFtgTLbktrkBlGAstub1mKb9W3j3/f2yd6EeEkDq4EkgHAWAowBwFACOAsBRADgKAEcB4CgAHAWAowBwFACOAsBRADgKAEcB4CgAHAWAowBwFACOAsBRADj/AEG3dJq0AI3SAAAAAElFTkSuQmCC",
    "opens": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAADhklEQVR4nO2dzZKiMBRGYWoe2b3l3ndmFlZXOTaRBJLcn++cpdIY7ne8kGDrum3bArr8sR4A2IIA4iCAOAggDgKI89d6ALWsz3u46cp2e6zWYzgiRAeIGP6yxBi3ewEiFPEb3sfvWgDvxavF83G4FcBz0c7g9XjcCgBzQABxEEAcBBAnzELQHp4WWrxe5B1BBxAHAcRBAHEQQBwEECf0LOAdi6twT7OQs9ABxEEAcRBAHAQQBwHEQQBxEECcNOsAGebkFtABxJkiQNR75dbMqNtwAX4OAgnamFW3oQJ8Dh4J6phZt2EClAaNBN+ZXbchs4Cjwa7P+9b7qj3D3UCLunXvALVB0An+x6puXQVoHRwSvLCsWzcBzgyKxZsXZ+rQS4IuAhD+dawkuCwA4ffDQoJLAhB+f2ZLcFoAwh/HTAlOrQN4DD+bXNvtsZ6ZHbTWobkDeAw/KzM6QZMAhD+f0RKstd8WPjv8jItEs+tR83pVHYB3vj2jOsGhAITvhxESfBWA8P3RW4KiAITvl54S7ArgIfxsMnmoz16uuwK07nxUWFkk8FKfve2Lp4DanbPC9x0v9Slt92sdoKX9zwwn4rqA1/q8j+v0ZwJnvzOjd4LRnLl3sCwn7wYShk/O5MK/homDAOIggDgIIA4CiIMA4iCAOAggDgKIgwDiIIA4CCAOAoiDAOIggDiyAvAVdi/SfFdwLd+C/nlO6QMvMh1gfd630rt+73GVjpC+A1wJUqEjpO4ANeH32iYqKTvAiMCydoNUAsx4p2YTIYUAFi06iwihBfBwbo4uQkgBPAT/SVQRQgnQO/j3sHrtO5oIIQQYGfznY2oiuF8HmBF+y/OteDxdveO2A8wOfm9bhW7gTgDL4Et/m1kENwJ4Cr60r4wimAvgOfjSvjOJYCZApOBLr5VBhOkCRA6+9NqRRZg6DewZ/nZ7rB7OocvSfywzp47Tfjy6d/i99tWT3hLMEGHoKSBTu68l2mnB5KdjW4kQ/CdRROgqAMH/xrsIXQQg+GO8inBJAIJvx5sIl74ruBcKwZewrrf57WDl8JfF/vjNloKtD9wTvU8LLUwXgODLWIgQ/l4AXMP8GgBsQQBxEEAcBBDn8MejuWiLz7eZFx1AHAQQBwHEObwGgNzQAcRBAHEQQBwEEAcBxEEAcRBAHAQQBwHEQQBxEEAcBBAHAcRBAHEQQBwEEOcfLN80d7bPrM4AAAAASUVORK5CYII=",
    "pct": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAD/UlEQVR4nO2dy3KbQBQFj1L+5Oxd2fuflUXKjrCFmMd939N7w0A3I4wQ3O73O0hffnkPgPjCAJrDAJrDAJrDAJrDAJrDAJrzprnw28f75UWG++8/N80xkNfcpC8EjUg/gzHYIxbAjvjvMAQ7RM4BJOVrLI+cszUDWIjibKDL8gxgdZRyNtBlKQBrKYxAj+kAvGQwAh14ISg42uFPnQSuDObVSZz08qrxuH+0tns4gFlZMwPWXHZWnu0Tje1W+QiYHWgHoTOcHRAaHwdDAcyseFWm5oyRiattk9520Rlg90juPhOMypWMQCwAKXldI5iVKhXBZQBRp9uo41phdVsk9oHIDCB91HaaBXYl7v49LwQ5InEE7x4sDMCJCPIBBuBCFPkAAzAnknyAAZgSTT7AAMyIKB/gPYEmRJUPDAQQ9X/yqOP6TmT5gOBHgNRRW+nojy4fED4H8L6qFYkM8oHBACy+qrX4ytmKLPIBpf8CvL7ZikAm+QDvCRQlm3xg4ZdBnkcr5cvDC0ECZJUPLATgNdCoR39m+cDiDGA9YMrXY/kjwGrg3jvojAryAaEHRGicGEbYOWdUkQ/wnsBpKskH+IygKarJBxQCOCy80FPCKsoHlAOoQlX5AC8EXVJZPsAAXlJdPsAATukgH2AAT+kiH2AAP+gkH2AAB7rJBxjAFx3lAwwAQF/5AANoLR9oHkB3+UDjACj/Hy0DoPz/tAuA8o+0CoDyf9ImAMp/TosAKP+c8gFQ/mtKB0D515QNgPLHKBkA5Y9TLgDKn6NUAJQ/T5kAKH+NEgFQ/jrpA6D8PVIHQPn7pA2A8mVIGQDly5EuAMqXJVUAlC9PmgAoX4cUAVC+HuEDoHxdQgdA+fq8eQ/gjE7yPZ+lFPIZQR3kR3maWrgAqsuXfKhmuTeHUr798sLMAJXlW7xjYXXbQ8wAlO+3HvcAKN93fa4BUL7/et0CqCw/Ey4ngdXlZ3q7mnkAlH9E86WcI8s2/QioLn+W2W0J+/r4ETrIt3j9rfSMYRJAB/kz7G6L5L5QD4Dyj0hti9RyVAPoJD/qC7CvxqUWQCf5o0R8u1qI18c/o5r8qIgHQPm5EA2A8vMhFgDl50QkAMrPy3YAlJ+b7QAiXdWKTsR7Al3fHl5JftRtuRqX2ElghG+2MiA1C0gtR/TfwFGpXeV/sitP8qNE/ELQldzK8jVv7lj5O7cbQs5WXFn+CrMRaHzhpHpL2OOAO8nnPYGPK/h4v3eS/4nn18Mz+1v9hpCO8jPh/sugqniFP7teBqCIdQQr62MAylhFsLqeMD8P74DGieFuYJwBDIl4TyBnACf4jCDyBZ8SRtzgOUBzGEBzGEBzGEBzGEBzGEBz/gKFAOnHjQLKDAAAAABJRU5ErkJggg==",
    "clicks": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAADlklEQVR4nO2dzVakMBBGC4+P7Ho8rsd3ZhYelGlpSCAkX/Ldu1OhW6puKgk/YZrnObaYPt83/zC/fUybO+zsk7IvtOFl65d7idz7Gwnuj18CHLXiZZuU7c58NtTllwDz28eU2pJJaP9Mz8YAEXkJXkvzuB9dgy67AnxvlCjCkujp830m6X2QJEDE+WoA2iQL8L0DU72h2JwG7rGXYJLfH9kV4L+dV9WA5PdJdgVYsySd5PfLtwBn5/Qkv2+meZ4l5+2ccq7DpgARbYJ8pgohwzWm+PtHYlp35bQyEpznUICFO4Nc4poCEpwjeRZw14WfUp/LhalzJFeANaVa2x1JoxLkUfRMYA7qFcWFl15G+0qfPxIvEek3gVBex+N1/cP6ev7jhi2Sv/WdtO6ybI4BFKZ8z/4Hblcry9NB4LpbqN36j76Prqgch7MAgj02ly4HQ/8ggDkIYA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOYggDkIYA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOa8Hm8yHqxA9oONAKkPi7qtfmrRBZx9UtjhCePhBbiaxNElGFoAViA7ZmgBSjKqBDaDwBI8vgpnhNkEApwg9dV6Efoi0AUksiyZk9sVqHcdQ1eA3FXPjra7Mp1UrQSXXhlz6gsvrhJ25TNSyvKzZI26oHV1ASKuSVBCoFxKlnE1CaS7APX+cwSkBVBhbyyRs78iTbqAiHtb990Bb9EN3UUzASLGeF/A1smhHhK/0FSACL9+Xk2O5gJE+EmwoCBDcwFck7+mpQhNTwWT/C9axkH6RJAjtatB9QpA8vepHR+uBppTVQBafxo141RNAJKfR614yV4LUJgjl0axEVSZBeQc+IiJf0QpHlKDQIfkR2gdp4wASkGpgcrx3i6AYr/XE3fHT6ICqLSG2igct4QA0A4EMAcBzEEAcxDAHAQwBwHMQQBzEMAcBDAHAcxBAHMQwBwEMAcBzEEAcxDAHAQwBwHMQQBzJARwvXNY4bhvF0DhzteesXkySKE11ETleGUEiNAJyt0oHWe1JWJyD3rErkMxBrKPhyu1kpGp1gWM2KLvpFa8qo4BkCCNmnGSGgRCfaoLQBXYp/pi16OtFt4rrRpG0y6AavBFyzg0Xyx6wbEaKDQAGQHWjCyDQtLXSAoA9WAaaA4CmIMA5iCAOQhgDgKYgwDmIIA5CGAOApiDAOYggDn/AJsh0uzbsQ3aAAAAAElFTkSuQmCC",
    "bounced": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAAEOUlEQVR4nO2cy3HbMBRFQU/aUTZsIw2kD+092ruPNJA2tIkKYlbUyBRJAe+HR917ZrywLRIQ7xG+FIdpmgrB5aN3BUhfKAA4FACcH70rgMhtHKsHXqfrdfCsy8BBYAwtoW/hIQMFcMQi9C2sZKAADngGv0QrAgeBxkSGb1EeBTAkOnyLctkFGNAr+DVauwS2AEoyhV9Ke30ogIJs4c+01ItdgBBt+DVNdUQZXAkMpLV/fny9V2vDLkCAJAztfF1yfE09KUAj0k+i9hPsVS4FaCBriJrjKUAwrWF6zzQoQCWWQdSeK6JMCtCJV+FGrTFQgI5shRy5wEQBKqgNxGKq5jnFXDv3rgDD1yeXCSuZQ9BIoAlfus6wKcAcPiVox2vRxqKcJasCLENHlkDaH3vfzCk9//L9PG0GIYe9xr8/f1++Zi8MjwGdtrzH47+1AAxfxt51s24JrM93F4Dh64iQwKNb+SiF4VvhKYHXmOKjlFKm88V1wILCq+soDdFzQHnvAiiBjprr12s3cI9vg0BKIMMzfKvjt+BScABW4XlIQAGcsQ7N+nxuN4UeoTuxmv3cxnFaG6h5NduW5bm1ANmnlrX1+/n7l+j83lu6VkvUrl3A8PU5ZRMhok5RGzsWkoWMAbJI4FmPXlu6mnJLCd4M6jku0L6vmk0hCZFjh7WyRC2ANMheLYG0XG9htz7x3lvJj4i7gKNIcLTwa/9vhWoMoJHAWwRNGWvvSzobWKM2XEsJts6lHgRO58uQrTXQBN/7k699fStms4AsEng2+dpWoNdu4N7xptPA3hJE9PdSCXrdD/DqOPN1gF4SRA72JBL02A2skcZlIShagqwj/SW3cZzmH4/XSxAtBLVcOM9wIoLfK8NrcciC2i7DfSnYqzXI8Km3nBpa0jJeCNkLsJYgQ/gz2SRI+5xAKwkyhT9zul6HyOVbyzqE3hGklSBj+I/0kkBTbvhj4uYwWsOUhN9j9/F0vQ58WngF3uH03HqO6BKsyuj6oMjpfBk89gOy3I9o/aBHD6m6PynUWoIs4S+RyuDdknQXoBQ7CbKGv6T3jOGRNN8L0IZ3lPCzkUaAUvb34/f+zvDlpBJgZhno/PvW34mclAKU8jp0hm9DWgFK2W/2o+vyrqQWgPhDAcChAOBQAHAoADgUABwKAA4FAIcCgEMBwKEA4FAAcCgAOBQAHAoADgUAhwKAQwHAoQDgUABwKAA4FAAcCgCO6MuhWZ7/T/SwBQCHAoBDAcB5EoDfu2vj6NeLLQA4qwIc3eoo3uE6bbYA7/DmPHmX6/P0tPDVF3Hef+ddgp+pEoC8LxwEgkMBwKEA4FAAcCgAOBQAHAoADgUAhwKAQwHAoQDgUABwKAA4FAAcCgAOBQCHAoBDAcChAOBQAHAoADgUABwKAA4FAIcCgEMBwKEA4FAAcCgAOBQAHAoADgUA5z/bzWp+06I4wgAAAABJRU5ErkJggg==",
    "dropped": "iVBORw0KGgoAAAANSUhEUgAAAIAAAACACAYAAADDPmHLAAADlklEQVR4nO2d227bMBBER0U/Oc8N8tz8s/oQCJVtSZYokjvLmQMEBZo6JjiHy4vccJrnGUaXX9ENMLFYAHEsgDgWQJzTAkzfn14tVoSlP08JsDSWpdHZYerP31dfsDR6/via6jdnbBgCf2Z6dw7A2OjRiBxMh1OAw+9DZD97FyDOrgAe/X2J6u9dAbzI60tUfx9OAZagD7SLQDM+FkAcCyDO5ZPAO4y0phhll9RVgBGOkUcJfiFkCsjaiVnbfUQVAUpG9PT9OWfp0NK2Zqh01SrA/PE1lYpQqw0tKA0+Q/hAgylgJAlGHfVrmiwCl0640oFMC0SF4BeaLgIzVgOl8IEOu4BMEqiFD3Q6B2CfEhSDX+h6DsBYDZTDBwIOgljODEbe218h5CQw+sxg9L39FUKfBkZI4FH/SNeHQVv0WiA6+G1oPg/Qsho4/H1oBADaSODwjwmfAp6pNSU4+HNQVYA1d6qBwz8PrQBAv12CavgA4RTwTMmUcPVnK0NdAdbUDsvh/0BfAdbUqAYO/pE0FWBNaYgO/5WUAgD7YV79e3VSTQHPnJkSHPwxaSvAmiXkvT/NPkMIALyG7fDPMYwApgwLII4FEMcCiGMBxLEA4lgAcSyAOBZAHAsgjgUQxwKIYwHEsQDiWABxLIA4FkCcKp8JjP7NXqYcVwBxLIA4FkCctwL407Vtie5fVwBxTgkQbemoMPTr6QrA0NiRYOnPt7eHb77I+/5iWIJfKBLAjIMXgeJYAHEsgDgWQJzUvyHkym6EbfXNgiuAOBZAHAsgjgUQxwKIYwHEsQDiWABxaAXo/chZ9RE3pQB3rn7J8H5MUAmwdZ1r77uDW1xTywyNAGydztaeVlAI0OLuP+b3ZSJUgCvltsXt4Wf/3cgihAlQ0qmRt4ePKkGYAKXP5yNuDwfG/TxBuuvjgfIQHf4rFItAEwfF/wtgnl9HHv0ASQVg7WTWdtWEQgCAr7PZ2tMKGgEAnk5naUcPqAQA4js/+v17QyeA6QulAHdGYdRrs0KxDdyj5/ZQMXyAtAIs9ApFNXyAXACgfTjK4QMJBDBtSSFAq1GqPvqBJAIA9cNy+D+kEQCoF5rD/08qAYD74Tn8R9IJAJSH6PBfSSnAEQ75GtQnge9YnxRuBf/u+wbAPM+pv/D3z3zn++pfqSuAuc9wawBzDQsgjgUQxwKIYwHEsQDiWABxLIA4FkAcCyCOBRDHAohjAcT5B7HnPUxNpXevAAAAAElFTkSuQmCC",
}

def _icon_path(key):
    """Decode a base64 icon to a temp PNG file and return its path."""
    data = _base64.b64decode(_ICON_B64[key])
    tmp = _tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.write(data)
    tmp.close()
    return tmp.name



# ── Brand colors ──────────────────────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x00, 0x40, 0x3A)   # #00403A  bar / headers
TEAL        = RGBColor(0x00, 0x8B, 0x7E)   # #008B7E  secondary / line
LIGHT_GREEN = RGBColor(0x7B, 0xC6, 0x7E)   # #7BC67E  click-rate line
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)   # near-black
LIGHT_BG    = RGBColor(0xF5, 0xF5, 0xF5)
ICON_COLOR  = RGBColor(0x00, 0x8B, 0x7E)

# ── KPI Icons — resolved from embedded base64 at runtime ──────────────────
KPI_ICONS = {
    "Total Delivered": "delivered",
    "Total Opens":     "opens",
    "OR":              "pct",
    "Total Clicks":    "clicks",
    "CTR":             "pct",
    "Bounced":         "bounced",
    "Dropped":         "dropped",
}

# Chart bar colors (updated)
BAR_DELIVERED = RGBColor(0x15, 0x25, 0x48)  # #152548
BAR_OPENS     = RGBColor(0x15, 0x85, 0x7C)  # #15857C
BAR_CLICKS    = RGBColor(0x92, 0xD0, 0x50)  # #92D050

# hex versions for matplotlib
HEX_DARK_GREEN  = "#00403A"
HEX_TEAL        = "#008B7E"
HEX_LIGHT_GREEN = "#7BC67E"
HEX_NAVY        = "#1A1A2E"
HEX_ORANGE      = "#E8A838"   # for bounced/dropped icons

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

NUM_FONT = "Roboto Condensed"   # number / KPI font

# ══════════════════════════════════════════════════════════════════════════════
# NATIVE CHART HELPERS  (no matplotlib – fully editable in PowerPoint)
# ══════════════════════════════════════════════════════════════════════════════

def _set_datalabels(ser_elem, font_size=800, font_name="Calibri",
                    position="inEnd", font_color_hex="FFFFFF"):
    """Safe data-label XML. Forces font color explicitly so PowerPoint cannot override it.
    position: 'inEnd' for inside bar end (default), 'outEnd' for outside.
    font_color_hex: 6-char hex string for label text color (default white)."""
    dLbls = etree.SubElement(ser_elem, qn("c:dLbls"))
    etree.SubElement(dLbls, qn("c:numFmt"),
                     attrib={"formatCode": "#,##0", "sourceLinked": "0"})
    # txPr: sets the default run properties for ALL label text in this series.
    # We must set solidFill on defRPr AND add a solid-color spPr so PPT
    # cannot revert to its theme color.
    txPr = etree.SubElement(dLbls, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                               attrib={"sz": str(font_size), "b": "1",
                                       "dirty": "0"})
    # Explicit solid fill on the run — overrides any theme/accent color
    sf = etree.SubElement(defRPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": font_color_hex})
    etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": font_name})
    # Also add end-of-paragraph run so the color sticks
    endParaRPr = etree.SubElement(p, qn("a:endParaRPr"),
                                   attrib={"lang": "en-US", "dirty": "0"})
    sf2 = etree.SubElement(endParaRPr, qn("a:solidFill"))
    etree.SubElement(sf2, qn("a:srgbClr"), attrib={"val": font_color_hex})
    etree.SubElement(dLbls, qn("c:showLegendKey"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showVal"),         attrib={"val": "1"})
    etree.SubElement(dLbls, qn("c:showCatName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showSerName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showPercent"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showBubbleSize"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:dLblPos"),         attrib={"val": position})


def _set_datalabels_with_shortbar_fallback(ser_elem, values, bar_color_hex,
                                            font_size=800, threshold_pct=8,
                                            global_max=None):
    """
    Sets data labels white inside the bar (inEnd).
    For bars too short to show an inside label (value < threshold_pct% of
    the chart global max across all series), adds a per-point override:
    solid fill box in bar_color_hex so the white label is always readable.
    global_max: pass the max value across ALL series so threshold is relative
    to the chart scale, not just this series (critical for Clicks vs Delivered).
    """
    if not values:
        return
    series_max = max(float(v) for v in values) if values else 1
    series_max = series_max if series_max > 0 else 1
    # Use global_max for threshold so tiny-relative-to-chart bars get outEnd label
    max_val = float(global_max) if global_max and float(global_max) > 0 else series_max

    dLbls = etree.SubElement(ser_elem, qn("c:dLbls"))
    etree.SubElement(dLbls, qn("c:numFmt"),
                     attrib={"formatCode": "#,##0", "sourceLinked": "0"})
    # Series-level txPr — white font
    txPr = etree.SubElement(dLbls, qn("c:txPr"))
    etree.SubElement(txPr, qn("a:bodyPr"))
    etree.SubElement(txPr, qn("a:lstStyle"))
    p = etree.SubElement(txPr, qn("a:p"))
    pPr = etree.SubElement(p, qn("a:pPr"))
    defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                               attrib={"sz": str(font_size), "b": "1", "dirty": "0"})
    sf = etree.SubElement(defRPr, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
    etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": "Calibri"})
    endRPr = etree.SubElement(p, qn("a:endParaRPr"),
                               attrib={"lang": "en-US", "dirty": "0"})
    sf2 = etree.SubElement(endRPr, qn("a:solidFill"))
    etree.SubElement(sf2, qn("a:srgbClr"), attrib={"val": "FFFFFF"})

    etree.SubElement(dLbls, qn("c:showLegendKey"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showVal"),         attrib={"val": "1"})
    etree.SubElement(dLbls, qn("c:showCatName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showSerName"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showPercent"),     attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:showBubbleSize"),  attrib={"val": "0"})
    etree.SubElement(dLbls, qn("c:dLblPos"),         attrib={"val": "inEnd"})

    # Per-point overrides for short bars
    for idx, v in enumerate(values):
        pct = float(v) / max_val * 100 if max_val else 0
        if pct < threshold_pct:
            # Short bar: add solid-fill callout box in the bar's own color
            dLbl = etree.SubElement(dLbls, qn("c:dLbl"))
            etree.SubElement(dLbl, qn("c:idx"), attrib={"val": str(idx)})
            # spPr: solid fill with bar color so label box is visible
            spPr = etree.SubElement(dLbl, qn("c:spPr"))
            fillEl = etree.SubElement(spPr, qn("a:solidFill"))
            etree.SubElement(fillEl, qn("a:srgbClr"), attrib={"val": bar_color_hex})
            ln = etree.SubElement(spPr, qn("a:ln"))
            etree.SubElement(ln, qn("a:noFill"))
            # txPr override — white font
            dLbl_txPr = etree.SubElement(dLbl, qn("c:txPr"))
            etree.SubElement(dLbl_txPr, qn("a:bodyPr"))
            etree.SubElement(dLbl_txPr, qn("a:lstStyle"))
            dp = etree.SubElement(dLbl_txPr, qn("a:p"))
            dpPr = etree.SubElement(dp, qn("a:pPr"))
            ddefRPr = etree.SubElement(dpPr, qn("a:defRPr"),
                                        attrib={"sz": str(font_size), "b": "1", "dirty": "0"})
            dsf = etree.SubElement(ddefRPr, qn("a:solidFill"))
            etree.SubElement(dsf, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
            etree.SubElement(ddefRPr, qn("a:latin"), attrib={"typeface": "Calibri"})
            dendRPr = etree.SubElement(dp, qn("a:endParaRPr"),
                                        attrib={"lang": "en-US", "dirty": "0"})
            dsf2 = etree.SubElement(dendRPr, qn("a:solidFill"))
            etree.SubElement(dsf2, qn("a:srgbClr"), attrib={"val": "FFFFFF"})
            etree.SubElement(dLbl, qn("c:showLegendKey"),  attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showVal"),         attrib={"val": "1"})
            etree.SubElement(dLbl, qn("c:showCatName"),     attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showSerName"),     attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showPercent"),     attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:showBubbleSize"),  attrib={"val": "0"})
            etree.SubElement(dLbl, qn("c:dLblPos"),         attrib={"val": "outEnd"})


def _style_axis_font(chart, font_name="Calibri", font_size=9):
    try:
        for axis in [chart.category_axis, chart.value_axis]:
            tf = axis.tick_labels.font
            tf.name = font_name
            tf.size = Pt(font_size)
    except Exception:
        pass


def _set_gridline_color(chart, color_hex="D9D9D9", line_width_emu=9525):
    """Set value axis major gridlines to a lighter color via XML."""
    try:
        val_axis_elem = chart.value_axis._element
        majorGridlines = val_axis_elem.find(qn("c:majorGridlines"))
        if majorGridlines is None:
            majorGridlines = etree.SubElement(val_axis_elem, qn("c:majorGridlines"))
        spPr = majorGridlines.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(majorGridlines, qn("c:spPr"))
        # Remove existing ln if any
        ln_old = spPr.find(qn("a:ln"))
        if ln_old is not None:
            spPr.remove(ln_old)
        ln = etree.SubElement(spPr, qn("a:ln"), attrib={"w": str(line_width_emu)})
        sf = etree.SubElement(ln, qn("a:solidFill"))
        etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex})
    except Exception:
        pass


def _add_bar_chart(slide, labels, series_data, x, y, w, h,
                   colors=None, show_legend=True, data_labels=True,
                   gap_width=80):
    """Native clustered column chart — fully editable in PPT, no repair prompt."""
    if not labels or not series_data:
        return None
    cd = ChartData()
    cd.categories = labels
    for sname, svals in series_data:
        cd.add_series(sname, [float(v) for v in svals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd).chart
    default_colors = [BAR_DELIVERED, BAR_OPENS, BAR_CLICKS, DARK_TEXT]
    if colors is None:
        colors = default_colors

    # Build hex strings for each series color for the short-bar fallback
    def _rgb_to_hex(rgb): return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"

    # Global max across all series so short-bar threshold is relative to chart scale
    all_vals_flat = [float(v) for _, vlist in series_data for v in vlist]
    global_max = max(all_vals_flat) if all_vals_flat else 1

    for i, ser in enumerate(chart.series):
        c = colors[i % len(colors)]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = c
        if data_labels:
            vals = [v for _, vlist in [series_data[i]] for v in vlist]
            color_hex = _rgb_to_hex((c[0], c[1], c[2])) if hasattr(c, '__getitem__') else \
                        f"{c.red:02X}{c.green:02X}{c.blue:02X}"
            _set_datalabels_with_shortbar_fallback(
                ser._element, vals, color_hex, global_max=global_max)
    chart.has_title = False
    chart.plots[0].gap_width = gap_width
    if show_legend:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    try:
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.has_major_gridlines = False
    except Exception:
        pass
    _set_gridline_color(chart, color_hex="E0E0E0")
    _style_axis_font(chart)
    return chart


def _make_line_ser(idx, name, vals, color_hex, marker_color_hex=None):
    """Build a <c:ser> element for a lineChart. Used in combo chart.
    Follows sfmc_report pattern exactly to avoid repair prompt."""
    if marker_color_hex is None:
        marker_color_hex = color_hex
    ser = etree.Element(qn("c:ser"))
    etree.SubElement(ser, qn("c:idx"),   attrib={"val": str(idx)})
    etree.SubElement(ser, qn("c:order"), attrib={"val": str(idx)})
    # series name as literal string
    tx = etree.SubElement(ser, qn("c:tx"))
    strRef = etree.SubElement(tx, qn("c:strRef"))
    etree.SubElement(strRef, qn("c:f")).text = f'"{name}"'
    strCache = etree.SubElement(strRef, qn("c:strCache"))
    etree.SubElement(strCache, qn("c:ptCount"), attrib={"val": "1"})
    pt = etree.SubElement(strCache, qn("c:pt"), attrib={"idx": "0"})
    etree.SubElement(pt, qn("c:v")).text = name
    # spPr — line colour, no fill
    spPr = etree.SubElement(ser, qn("c:spPr"))
    etree.SubElement(spPr, qn("a:noFill"))
    ln = etree.SubElement(spPr, qn("a:ln"), attrib={"w": "25400"})  # 2pt line
    sf = etree.SubElement(ln, qn("a:solidFill"))
    etree.SubElement(sf, qn("a:srgbClr"), attrib={"val": color_hex})
    # marker — filled circle
    marker = etree.SubElement(ser, qn("c:marker"))
    etree.SubElement(marker, qn("c:symbol"), attrib={"val": "circle"})
    etree.SubElement(marker, qn("c:size"),   attrib={"val": "5"})
    mSpPr = etree.SubElement(marker, qn("c:spPr"))
    mFill = etree.SubElement(mSpPr, qn("a:solidFill"))
    etree.SubElement(mFill, qn("a:srgbClr"), attrib={"val": marker_color_hex})
    mLn = etree.SubElement(mSpPr, qn("a:ln"))
    mLnFill = etree.SubElement(mLn, qn("a:solidFill"))
    etree.SubElement(mLnFill, qn("a:srgbClr"), attrib={"val": marker_color_hex})
    # inline data values
    val_el = etree.SubElement(ser, qn("c:val"))
    numLit = etree.SubElement(val_el, qn("c:numLit"))
    etree.SubElement(numLit, qn("c:formatCode")).text = "0.00"
    etree.SubElement(numLit, qn("c:ptCount"), attrib={"val": str(len(vals))})
    for i, v in enumerate(vals):
        pt2 = etree.SubElement(numLit, qn("c:pt"), attrib={"idx": str(i)})
        etree.SubElement(pt2, qn("c:v")).text = str(round(float(v), 4))
    etree.SubElement(ser, qn("c:smooth"), attrib={"val": "0"})
    return ser


def _add_slide1_combo_chart(slide, labels, d_vals, or_vals, ct_vals):
    """
    Combo chart using full XML-replacement approach (same as SoMe report).
    - Delivered:    clustered columns, hex #15857C, labels inEnd (inside bar), white font
    - Open Rate %:  line, hex #152548, solid fill labels, white font, values as XX.XX%
    - Click Rate %: line, hex #92D050, solid fill labels, white font, values as XX.XX%
    OR/CTR values are already in percent (e.g. 25.6 → shown as 25.6%, not 25600%).
    Delivered on secondary (right, hidden) axis so lines scale independently.
    No repair prompt — chart XML and workbook replaced atomically.
    """
    import io
    import openpyxl as _openpyxl

    n = len(labels)

    COL_DEL = "15857C"
    COL_OR  = "152548"
    COL_CTR = "92D050"

    def _pts(vals):
        return " ".join(f'<c:pt idx="{i}"><c:v>{v}</c:v></c:pt>' for i, v in enumerate(vals))

    lbl_pts = _pts(labels)
    del_pts = _pts(d_vals)
    or_pts  = _pts(or_vals)
    ctr_pts = _pts(ct_vals)

    max_del = max(d_vals)   if d_vals   else 1
    all_pct_vals = (or_vals or [0]) + (ct_vals or [0])
    max_pct = max(all_pct_vals)
    min_pct = min(all_pct_vals)
    del_max = max_del * 1.18
    pct_pad = (max_pct - min_pct) * 0.35 if (max_pct - min_pct) > 0 else max_pct * 0.1
    pct_max = max_pct + pct_pad
    pct_min = max(0, min_pct - pct_pad)

    def _line_series(idx, order, name, col, pts, ref_col_letter):
        return f"""
      <c:ser>
        <c:idx val="{idx}"/><c:order val="{order}"/>
        <c:tx><c:strRef><c:f>Sheet1!${ref_col_letter}$1</c:f>
          <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>{name}</c:v></c:pt></c:strCache>
        </c:strRef></c:tx>
        <c:spPr>
          <a:ln w="22860"><a:solidFill><a:srgbClr val="{col}"/></a:solidFill></a:ln>
        </c:spPr>
        <c:marker><c:symbol val="circle"/><c:size val="5"/>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{col}"/></a:solidFill>
            <a:ln><a:solidFill><a:srgbClr val="{col}"/></a:solidFill></a:ln>
          </c:spPr>
        </c:marker>
        <c:dLbls>
          <c:numFmt formatCode='0.00"%"' sourceLinked="0"/>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{col}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
            <a:defRPr b="1" sz="900"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr>
          </a:pPr></a:p></c:txPr>
          <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
          <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
          <c:dLblPos val="t"/>
        </c:dLbls>
        <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
          <c:strCache><c:ptCount val="{n}"/>{lbl_pts}</c:strCache>
        </c:strRef></c:cat>
        <c:val><c:numRef><c:f>Sheet1!${ref_col_letter}$2:${ref_col_letter}${n+1}</c:f>
          <c:numCache><c:formatCode>0.00</c:formatCode><c:ptCount val="{n}"/>{pts}</c:numCache>
        </c:numRef></c:val>
        <c:smooth val="0"/>
      </c:ser>"""

    chart_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:date1904 val="0"/>
  <c:lang val="en-US"/>
  <c:chart>
    <c:autoTitleDeleted val="1"/>
    <c:plotArea>
      <c:layout/>

      <c:barChart>
        <c:barDir val="col"/>
        <c:grouping val="clustered"/>
        <c:varyColors val="0"/>
        <c:ser>
          <c:idx val="0"/><c:order val="0"/>
          <c:tx><c:strRef><c:f>Sheet1!$B$1</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>Delivered</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          <c:spPr>
            <a:solidFill><a:srgbClr val="{COL_DEL}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </c:spPr>
          <c:dLbls>
            <c:numFmt formatCode="#,##0" sourceLinked="0"/>
            <c:spPr>
              <a:solidFill><a:srgbClr val="{COL_DEL}"/></a:solidFill>
              <a:ln><a:noFill/></a:ln>
            </c:spPr>
            <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
              <a:defRPr b="1" sz="900"><a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill></a:defRPr>
            </a:pPr></a:p></c:txPr>
            <c:showLegendKey val="0"/><c:showVal val="1"/><c:showCatName val="0"/>
            <c:showSerName val="0"/><c:showPercent val="0"/><c:showBubbleSize val="0"/>
            <c:dLblPos val="ctr"/>
          </c:dLbls>
          <c:cat><c:strRef><c:f>Sheet1!$A$2:$A${n+1}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{lbl_pts}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>Sheet1!$B$2:$B${n+1}</c:f>
            <c:numCache><c:formatCode>#,##0</c:formatCode><c:ptCount val="{n}"/>{del_pts}</c:numCache>
          </c:numRef></c:val>
        </c:ser>
        <c:gapWidth val="60"/>
        <c:overlap val="0"/>
        <c:axId val="10"/>
        <c:axId val="20"/>
      </c:barChart>

      <c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>
        {_line_series(1, 1, "Open Rate %",  COL_OR,  or_pts,  "C")}
        {_line_series(2, 2, "Click Rate %", COL_CTR, ctr_pts, "D")}
        <c:marker val="1"/>
        <c:axId val="11"/>
        <c:axId val="30"/>
      </c:lineChart>

      <c:catAx>
        <c:axId val="10"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:solidFill><a:srgbClr val="DDDDDD"/></a:solidFill></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="900" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="20"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

      <c:valAx>
        <c:axId val="20"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="0"/><c:max val="{del_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="l"/>
        <c:majorGridlines>
          <c:spPr><a:ln w="6350"><a:solidFill><a:srgbClr val="E0E0E0"/></a:solidFill></a:ln></c:spPr>
        </c:majorGridlines>
        <c:numFmt formatCode="#,##0" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="10"/>
      </c:valAx>

      <c:valAx>
        <c:axId val="30"/>
        <c:scaling><c:orientation val="minMax"/><c:min val="{pct_min}"/><c:max val="{pct_max}"/></c:scaling>
        <c:delete val="0"/>
        <c:axPos val="r"/>
        <c:numFmt formatCode='0.00"%"' sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/>
        <c:spPr><a:ln><a:noFill/></a:ln></c:spPr>
        <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
          <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="666666"/></a:solidFill></a:defRPr>
        </a:pPr></a:p></c:txPr>
        <c:crossAx val="11"/>
      </c:valAx>

      <c:catAx>
        <c:axId val="11"/>
        <c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="1"/>
        <c:axPos val="b"/>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="none"/>
        <c:crossAx val="30"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:noMultiLvlLbl val="1"/>
      </c:catAx>

    </c:plotArea>
    <c:legend>
      <c:legendPos val="b"/>
      <c:spPr><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>
      <c:txPr><a:bodyPr/><a:lstStyle/><a:p><a:pPr>
        <a:defRPr sz="800" b="0"><a:solidFill><a:srgbClr val="444444"/></a:solidFill></a:defRPr>
      </a:pPr></a:p></c:txPr>
      <c:overlay val="0"/>
    </c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </c:spPr>
  <c:externalData r:id="rId1"><c:autoUpdate val="0"/></c:externalData>
</c:chartSpace>""".encode("utf-8")

    wb2 = _openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "Sheet1"
    ws2.append(["Month", "Delivered", "Open Rate %", "Click Rate %"])
    for i in range(n):
        ws2.append([labels[i], d_vals[i], or_vals[i], ct_vals[i]])
    buf = io.BytesIO()
    wb2.save(buf)
    wb_bytes = buf.getvalue()

    cd = ChartData()
    cd.categories = labels
    cd.add_series("Open Rate %",  [float(v) for v in or_vals])
    cd.add_series("Click Rate %", [float(v) for v in ct_vals])
    cd.add_series("Delivered",    [float(v) for v in d_vals])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3), Inches(1.95), Inches(12.7), Inches(5.3), cd)

    chart_part = chart_frame.chart.part
    chart_part._element = etree.fromstring(chart_xml)

    xlsx_part = chart_part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/package")
    xlsx_part._blob = wb_bytes

    return chart_frame

def normalise_df(df: pd.DataFrame) -> pd.DataFrame:
    """Normalise columns, add 'Status' alias if needed."""
    df.columns = df.columns.str.strip()
    # Normalise month order
    month_order = ["January","February","March","April","May","June",
                   "July","August","September","October","November","December"]
    df["Month"] = pd.Categorical(df["Month"], categories=month_order, ordered=True)
    # Unified Status column (handle uppercase STATUS)
    if "STATUS" in df.columns and "Status" not in df.columns:
        df["Status"] = df["STATUS"]
    elif "Status" not in df.columns:
        raise ValueError("No 'Status' or 'STATUS' column found in input file.")
    # Fill blank Campaign / TA
    df["Campaign"] = df["Campaign"].fillna("").astype(str).str.strip()
    df["TA"]       = df["TA"].fillna("").astype(str).str.strip()
    # Numeric cols
    for c in ["Opens", "Clicks", "TOTAL OPENS", "TOTAL CLICKS"]:
        if c in df.columns:
            df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)
    return df


def load_and_calc(path: str) -> pd.DataFrame:
    df = pd.read_excel(path)
    return normalise_df(df)


def kpis(df: pd.DataFrame) -> dict:
    """Top-level KPIs across entire dataset."""
    total_delivered = (df["Status"].str.lower() == "delivered").sum()
    total_bounced   = (df["Status"].str.lower() == "bounced").sum()
    total_dropped   = (df["Status"].str.lower() == "dropped").sum()
    # Opens and Clicks only from rows where Status == "delivered"
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), None)
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), None)
    delivered_rows = df[df["Status"].str.lower() == "delivered"]
    total_opens  = int(delivered_rows[opens_col].sum())  if opens_col  else 0
    total_clicks = int(delivered_rows[clicks_col].sum()) if clicks_col else 0
    or_pct  = total_opens  / total_delivered * 100 if total_delivered else 0
    ctr_pct = total_clicks / total_delivered * 100 if total_delivered else 0
    return dict(
        delivered=int(total_delivered),
        bounced=int(total_bounced),
        dropped=int(total_dropped),
        opens=int(total_opens),
        clicks=int(total_clicks),
        or_pct=or_pct,
        ctr_pct=ctr_pct,
    )


def monthly_kpis(df: pd.DataFrame) -> pd.DataFrame:
    # Only delivered rows contribute to Opens and Clicks
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")
    tmp = df.copy()
    tmp["_delivered"] = (tmp["Status"].str.lower() == "delivered").astype(int)
    # Zero out Opens/Clicks on non-delivered rows before aggregating
    tmp["_opens"]  = tmp[opens_col].where(tmp["Status"].str.lower() == "delivered", 0)
    tmp["_clicks"] = tmp[clicks_col].where(tmp["Status"].str.lower() == "delivered", 0)
    grp = tmp.groupby("Month", observed=True, as_index=False).agg(
        Delivered=("_delivered", "sum"),
        Opens=("_opens", "sum"),
        Clicks=("_clicks", "sum"),
    )
    grp = grp[grp["Delivered"] > 0].copy()
    grp["OR"]  = grp["Opens"]  / grp["Delivered"] * 100
    grp["CTR"] = grp["Clicks"] / grp["Delivered"] * 100
    return grp


def ta_kpis(df: pd.DataFrame) -> pd.DataFrame:
    """Per-TA summary."""
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")
    ta_list = df[df["TA"].str.len() > 0]["TA"].unique()
    rows = []
    for ta in ta_list:
        sub = df[df["TA"] == ta]
        sub_del = sub[sub["Status"].str.lower() == "delivered"]  # delivered rows only
        campaigns = sub[sub["Campaign"].str.len() > 0]["Campaign"].nunique()
        delivered = len(sub_del)
        bounced   = (sub["Status"].str.lower() == "bounced").sum()
        dropped   = (sub["Status"].str.lower() == "dropped").sum()
        opens     = sub_del[opens_col].sum()
        clicks    = sub_del[clicks_col].sum()
        or_pct    = opens  / delivered * 100 if delivered else 0
        ctr_pct   = clicks / delivered * 100 if delivered else 0
        rows.append(dict(TA=ta, Campaigns=int(campaigns), Delivered=int(delivered),
                         Bounced=int(bounced), Dropped=int(dropped),
                         Opens=int(opens), Clicks=int(clicks),
                         OR=or_pct, CTR=ctr_pct))
    return pd.DataFrame(rows)


def ta_campaign_kpis(df: pd.DataFrame, ta: str) -> pd.DataFrame:
    """Per-campaign breakdown for one TA. Normalises campaign names to avoid duplicates."""
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")
    sub = df[(df["TA"] == ta) & (df["Campaign"].str.len() > 0)].copy()
    # Normalise: strip whitespace and unify case for grouping, but keep original label
    sub["_camp_norm"] = sub["Campaign"].str.strip().str.upper()
    sub["_delivered"] = (sub["Status"].str.lower() == "delivered").astype(int)
    # Zero out Opens/Clicks on non-delivered rows
    sub["_opens"]  = sub[opens_col].where(sub["Status"].str.lower() == "delivered", 0)
    sub["_clicks"] = sub[clicks_col].where(sub["Status"].str.lower() == "delivered", 0)
    # Use the most common original spelling as the display label
    label_map = sub.groupby("_camp_norm")["Campaign"].agg(
        lambda x: x.str.strip().mode()[0] if len(x.str.strip().mode()) else x.iloc[0]
    )
    grp = sub.groupby("_camp_norm", as_index=False).agg(
        Delivered=("_delivered", "sum"),
        Opens=("_opens", "sum"),
        Clicks=("_clicks", "sum"),
    )
    grp["Campaign"] = grp["_camp_norm"].map(label_map)
    grp = grp.drop(columns=["_camp_norm"])
    grp = grp[grp["Delivered"] > 0]
    return grp


# Canonical TA config: real data name → display label → color
TA_CONFIG = [
    ("Diabetes",  "DIA",  RGBColor(0x15, 0x25, 0x48)),  # #152548
    ("HAC",       "HAC",  RGBColor(0x15, 0x85, 0x7C)),  # #15857C
    ("ONCO",      "ONCO", RGBColor(0x92, 0xD0, 0x50)),  # #92D050
    ("Vaccines",  "VACC", RGBColor(0x11, 0x9E, 0xDD)),  # #119EDD
]

def monthly_ta_kpis(df: pd.DataFrame) -> pd.DataFrame:
    """Delivered by Month × TA (using real TA names from data)."""
    rows = []
    for month in df["Month"].cat.categories:
        sub_m = df[df["Month"] == month]
        if (sub_m["Status"].str.lower() == "delivered").sum() == 0:
            continue
        row = {"Month": month}
        for ta_real, ta_label, _ in TA_CONFIG:
            sub_t = sub_m[sub_m["TA"] == ta_real]
            row[ta_label] = int((sub_t["Status"].str.lower() == "delivered").sum())
        rows.append(row)
    return pd.DataFrame(rows)


def best_campaigns(df: pd.DataFrame, min_delivered: int = 50) -> pd.DataFrame:
    """Top-3 by OR in latest month. Skips rows with blank Campaign or blank TA."""
    latest = df["Month"].max()
    sub = df[df["Month"] == latest].copy()
    # Only keep rows where both Campaign and TA are non-blank
    sub = sub[
        (sub["Campaign"].str.strip() != "") &
        (sub["TA"].str.strip() != "")
    ]
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")
    sub["_delivered"] = (sub["Status"].str.lower() == "delivered").astype(int)
    # Zero out Opens/Clicks on non-delivered rows
    sub["_opens"]  = sub[opens_col].where(sub["Status"].str.lower() == "delivered", 0)
    sub["_clicks"] = sub[clicks_col].where(sub["Status"].str.lower() == "delivered", 0)
    grp = sub.groupby("Campaign", as_index=False).agg(
        Delivered=("_delivered", "sum"),
        Opens=("_opens", "sum"),
        Clicks=("_clicks", "sum"),
    )
    # Add TA column: most common TA per campaign
    ta_map = sub.groupby("Campaign")["TA"].agg(lambda x: x.mode()[0] if len(x.mode()) else "")
    grp["TA"] = grp["Campaign"].map(ta_map)
    # Drop any that still ended up with blank TA after aggregation
    grp = grp[grp["TA"].str.strip() != ""]
    # Adaptive min_delivered
    filtered = grp[grp["Delivered"] >= min_delivered].copy()
    if len(filtered) < 1:
        filtered = grp[grp["Delivered"] >= 2].copy()
    if len(filtered) < 1:
        filtered = grp[grp["Delivered"] >= 1].copy()
    filtered["OR"]  = filtered["Opens"]  / filtered["Delivered"].replace(0, 1) * 100
    filtered["CTR"] = filtered["Clicks"] / filtered["Delivered"].replace(0, 1) * 100
    # Sort by OR descending — exact values, no rounding
    filtered = filtered.sort_values("OR", ascending=False).head(7).reset_index(drop=True)
    return filtered, latest


def best_campaigns_overall(df: pd.DataFrame, min_delivered: int = 50) -> pd.DataFrame:
    """Top-3 by OR across ALL months combined. Skips rows with blank Campaign or blank TA.
    Only keeps campaigns where total Delivered > min_delivered."""
    sub = df.copy()
    # Only keep rows where both Campaign and TA are non-blank
    sub = sub[
        (sub["Campaign"].str.strip() != "") &
        (sub["TA"].str.strip() != "")
    ]
    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")
    sub["_delivered"] = (sub["Status"].str.lower() == "delivered").astype(int)
    # Zero out Opens/Clicks on non-delivered rows
    sub["_opens"]  = sub[opens_col].where(sub["Status"].str.lower() == "delivered", 0)
    sub["_clicks"] = sub[clicks_col].where(sub["Status"].str.lower() == "delivered", 0)
    # Group by Campaign across ALL months
    grp = sub.groupby("Campaign", as_index=False).agg(
        Delivered=("_delivered", "sum"),
        Opens=("_opens", "sum"),
        Clicks=("_clicks", "sum"),
    )
    # Add TA column: most common TA per campaign
    ta_map = sub.groupby("Campaign")["TA"].agg(lambda x: x.mode()[0] if len(x.mode()) else "")
    grp["TA"] = grp["Campaign"].map(ta_map)
    # Drop any that still ended up with blank TA after aggregation
    grp = grp[grp["TA"].str.strip() != ""]
    # Total Delivered (across all months) must be > min_delivered
    filtered = grp[grp["Delivered"] > min_delivered].copy()
    filtered["OR"]  = filtered["Opens"]  / filtered["Delivered"].replace(0, 1) * 100
    filtered["CTR"] = filtered["Clicks"] / filtered["Delivered"].replace(0, 1) * 100
    # Sort by OR descending — exact values, no rounding
    filtered = filtered.sort_values("OR", ascending=False).head(3).reset_index(drop=True)
    return filtered


# ══════════════════════════════════════════════════════════════════════════════
# PPT HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    from pptx.util import Pt as Pt2
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = Pt2(line_width_pt)
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, font_size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, wrap=True, italic=False, num_font=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    # Use Invention for number/KPI values when requested
    run.font.name = NUM_FONT if num_font else "Calibri"
    if color:
        run.font.color.rgb = color
    return txb


def add_image(slide, img_path, l, t, w, h):
    return slide.shapes.add_picture(img_path, l, t, w, h)


def kpi_block(slide, label, value, icon_char, l, t, w=Inches(1.7), h=Inches(0.85),
              val_color=None, val_size=22):
    """A single KPI stat block: value + label (icon omitted when icon_char is empty)."""
    if val_color is None:
        val_color = DARK_TEXT

    _icon_key = KPI_ICONS.get(label)
    icon_path = _icon_path(_icon_key) if _icon_key else None
    has_img_icon = bool(icon_path)

    if has_img_icon:
        icon_sz = Inches(0.46)
        slide.shapes.add_picture(icon_path, l, t + Inches(0.02), icon_sz, icon_sz)
        text_l = l + icon_sz + Inches(0.08)
        text_w = w - icon_sz - Inches(0.08)
    elif icon_char:
        icon_sz = Inches(0.38)
        add_rect(slide, l, t + Inches(0.08), icon_sz, icon_sz,
                 fill_rgb=RGBColor(0xE8, 0xF5, 0xF4), line_rgb=None)
        add_text(slide, icon_char, l + Inches(0.04), t + Inches(0.08),
                 icon_sz, icon_sz, font_size=14, color=ICON_COLOR, align=PP_ALIGN.CENTER)
        text_l = l + icon_sz + Inches(0.06)
        text_w = w - icon_sz - Inches(0.06)
    else:
        text_l = l
        text_w = w

    add_text(slide, label, text_l, t, text_w, Inches(0.28),
             font_size=10, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.LEFT)
    add_text(slide, str(value), text_l, t + Inches(0.26), text_w, Inches(0.45),
             font_size=val_size, bold=True, color=val_color, align=PP_ALIGN.LEFT,
             num_font=True)


def add_slide_title(slide, title_text, subtitle=None):
    """Standard slide title in dark green, top-left."""
    add_text(slide, title_text, Inches(0.4), Inches(0.15), Inches(6), Inches(0.55),
             font_size=24, bold=True, color=DARK_GREEN)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.65), Inches(7), Inches(0.3),
                 font_size=10, color=RGBColor(0x55, 0x55, 0x55), italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – REE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════

def slide1_overview(prs, df, tmp_dir):
    k = kpis(df)
    monthly = monthly_kpis(df)

    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "REE Overview")

    # ── Top KPI row ──
    KPI_TOP = Inches(0.95)
    KPI_H   = Inches(0.85)

    _bd_icon_path = _icon_path("bounced")
    _dp_icon_path = _icon_path("dropped")
    _bd_off = Inches(0.42)

    slide.shapes.add_picture(_bd_icon_path, Inches(9.8), Inches(0.18), Inches(0.34), Inches(0.34))
    add_text(slide, "Bounced", Inches(9.8)+_bd_off, Inches(0.15), Inches(1.5), Inches(0.3),
             font_size=9, color=RGBColor(0x55,0x55,0x55))
    add_text(slide, f"{k['bounced']:,}", Inches(9.8)+_bd_off, Inches(0.38), Inches(1.5), Inches(0.45),
             font_size=22, bold=True, color=DARK_TEXT, num_font=True)

    slide.shapes.add_picture(_dp_icon_path, Inches(11.3), Inches(0.18), Inches(0.34), Inches(0.34))
    add_text(slide, "Dropped", Inches(11.3)+_bd_off, Inches(0.15), Inches(1.5), Inches(0.3),
             font_size=9, color=RGBColor(0x55,0x55,0x55))
    add_text(slide, f"{k['dropped']:,}", Inches(11.3)+_bd_off, Inches(0.38), Inches(1.5), Inches(0.45),
             font_size=22, bold=True, color=DARK_TEXT, num_font=True)

    kpi_data = [
        ("Total Delivered", f"{k['delivered']:,}", ""),
        ("Total Opens",     f"{k['opens']:,}",     ""),
        ("OR",              f"{k['or_pct']:.2f}%", ""),
        ("Total Clicks",    f"{k['clicks']:,}",    ""),
        ("CTR",             f"{k['ctr_pct']:.2f}%",""),
    ]
    kpi_w = Inches(2.1)
    for i, (lbl, val, icon) in enumerate(kpi_data):
        kpi_block(slide, lbl, val, icon,
                  l=Inches(0.35 + i * 2.55), t=KPI_TOP, w=kpi_w, h=KPI_H)

    # ── Native combo chart: bars=Delivered, lines=OR & CTR ──
    months    = [str(m)[:3] for m in monthly["Month"].tolist()]
    delivered = [int(v) for v in monthly["Delivered"].tolist()]   # integers, no .0
    or_vals   = monthly["OR"].tolist()
    ctr_vals  = monthly["CTR"].tolist()

    _add_slide1_combo_chart(slide, months, delivered, or_vals, ctr_vals)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – REE OVERVIEW – TAs
# ══════════════════════════════════════════════════════════════════════════════

def _set_legend_font(chart, font_size_pt=7, font_name="Calibri"):
    """Set legend entry font size via XML."""
    try:
        legend_elem = chart.legend._element
        txPr = legend_elem.find(qn("c:txPr"))
        if txPr is None:
            txPr = etree.SubElement(legend_elem, qn("c:txPr"))
        # Clear existing txPr children and rebuild
        for child in list(txPr):
            txPr.remove(child)
        etree.SubElement(txPr, qn("a:bodyPr"))
        etree.SubElement(txPr, qn("a:lstStyle"))
        p = etree.SubElement(txPr, qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        defRPr = etree.SubElement(pPr, qn("a:defRPr"),
                                   attrib={"sz": str(int(font_size_pt * 100)), "b": "0"})
        etree.SubElement(defRPr, qn("a:latin"), attrib={"typeface": font_name})
    except Exception:
        pass


def slide2_ta_overview(prs, df, tmp_dir):
    slide  = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)

    # Title: Invention bold
    title_txb = add_text(slide, "REE Overview – TAs",
                         Inches(0.4), Inches(0.15), Inches(8), Inches(0.55),
                         font_size=24, bold=True, color=DARK_GREEN)
    title_txb.text_frame.paragraphs[0].runs[0].font.name = "Invention"

    # Real TA names from data, display order: DIA TL, HAC TR, ONCO BL, VACC BR
    col_w = Inches(6.45)
    row_h = Inches(3.2)
    pad   = Inches(0.08)
    cols  = 2

    opens_col  = next((c for c in ["Opens", "TOTAL OPENS"] if c in df.columns), "Opens")
    clicks_col = next((c for c in ["Clicks", "TOTAL CLICKS"] if c in df.columns), "Clicks")

    for pos, (ta_real, ta_label, _) in enumerate(TA_CONFIG):
        col   = pos % cols
        r     = pos // cols
        box_l = Inches(0.15) + col * col_w
        box_t = Inches(0.85) + r * row_h

        # Build TA stats using real name
        sub = df[df["TA"] == ta_real]
        sub_del   = sub[sub["Status"].str.lower() == "delivered"]
        delivered = len(sub_del)
        bounced   = int((sub["Status"].str.lower() == "bounced").sum())
        dropped   = int((sub["Status"].str.lower() == "dropped").sum())
        opens     = int(sub_del[opens_col].sum())
        clicks    = int(sub_del[clicks_col].sum())
        or_pct    = opens  / delivered * 100 if delivered else 0
        ctr_pct   = clicks / delivered * 100 if delivered else 0
        campaigns = int(sub[sub["Campaign"].str.len() > 0]["Campaign"]
                        .str.strip().str.casefold().nunique())

        # Panel background
        add_rect(slide, box_l, box_t, col_w - pad, row_h - pad,
                 fill_rgb=RGBColor(0xF8, 0xFD, 0xFB),
                 line_rgb=RGBColor(0xCC, 0xE8, 0xE2), line_width_pt=0.75)

        # TA name — Trebuchet MS 13pt #0D6E6E, no solid fill badge
        add_text(slide, ta_label, box_l + Inches(0.12), box_t + Inches(0.1),
                 Inches(1.5), Inches(0.28), font_size=13, bold=True,
                 color=RGBColor(0x0D, 0x6E, 0x6E), align=PP_ALIGN.LEFT)

        add_text(slide, f"{campaigns} Campaigns",
                 box_l + Inches(0.9), box_t + Inches(0.12),
                 Inches(2.5), Inches(0.28), font_size=9, bold=True, color=DARK_TEXT)

        # Bounced: icon + text
        _ic_sz = Inches(0.20)
        slide.shapes.add_picture(_icon_path("bounced"),
                 box_l + col_w - Inches(2.5), box_t + Inches(0.08), _ic_sz, _ic_sz)
        add_text(slide, f"Bounced: {bounced}",
                 box_l + col_w - Inches(2.5) + _ic_sz + Inches(0.04), box_t + Inches(0.1),
                 Inches(1.1), Inches(0.28), font_size=8, color=RGBColor(0x55,0x55,0x55))
        # Dropped: icon + text
        slide.shapes.add_picture(_icon_path("dropped"),
                 box_l + col_w - Inches(1.25), box_t + Inches(0.08), _ic_sz, _ic_sz)
        add_text(slide, f"Dropped: {dropped}",
                 box_l + col_w - Inches(1.25) + _ic_sz + Inches(0.04), box_t + Inches(0.1),
                 Inches(1.1), Inches(0.28), font_size=8, color=RGBColor(0x55,0x55,0x55))

        # Mini KPI row — icon LEFT of label, value below
        mini_kpis = [
            ("Total Delivered", f"{delivered:,}", "delivered"),
            ("Total Opens",     f"{opens:,}",     "opens"),
            ("OR",              f"{or_pct:.2f}%", "pct"),
            ("Total Clicks",    f"{clicks:,}",    "clicks"),
            ("CTR",             f"{ctr_pct:.2f}%","pct"),
        ]
        mk_w    = (col_w - Inches(0.3)) / len(mini_kpis)
        icon_sz = Inches(0.20)   # icon size in slide units
        icon_gap = Inches(0.04)  # gap between icon and label text
        for mi, (ml, mv, icon_key) in enumerate(mini_kpis):
            ml_l = box_l + Inches(0.12) + mi * mk_w
            row_t = box_t + Inches(0.42)
            # Icon — vertically centered on label row, at left of cell
            try:
                slide.shapes.add_picture(
                    _icon_path(icon_key),
                    ml_l,
                    row_t + (Inches(0.22) - icon_sz) / 2,
                    icon_sz, icon_sz)
            except Exception:
                pass
            # Label to the right of icon
            lbl_txb = add_text(slide, ml,
                     ml_l + icon_sz + icon_gap, row_t,
                     mk_w - icon_sz - icon_gap, Inches(0.22),
                     font_size=7, color=RGBColor(0x55,0x55,0x55))
            lbl_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto"
            # Value below — aligned to start under the label (same offset as label, past the icon)
            val_txb = add_text(slide, mv,
                     ml_l + icon_sz + icon_gap, row_t + Inches(0.24),
                     mk_w - icon_sz - icon_gap, Inches(0.32),
                     font_size=10, bold=True, color=DARK_GREEN, num_font=False)
            val_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto Condensed"

        # Campaign chart: Delivered + Opens + Clicks
        # Filter: Delivered >= 50; if >5 campaigns, take top 5 by OR%
        camp_df = ta_campaign_kpis(df, ta_real)
        if len(camp_df) > 0:
            camp_df = camp_df.copy()
            camp_df["OR"] = camp_df.apply(
                lambda r: r["Opens"] / r["Delivered"] * 100 if r["Delivered"] > 0 else 0, axis=1)
            # Apply delivered >= 50 filter only if enough data remains
            filtered = camp_df[camp_df["Delivered"] >= 50].copy()
            if len(filtered) == 0:
                filtered = camp_df.copy()   # fallback: no filter if all below 50
            filtered = filtered.sort_values("OR", ascending=False)
            # Cap at top 5
            filtered = filtered.head(5)
            # Truncate labels to 18 chars for readability
            labels = [str(c)[:18] + ("…" if len(str(c)) > 18 else "")
                      for c in filtered["Campaign"].tolist()]
            chart = _add_bar_chart(
                slide, labels,
                [("Delivered", [int(v) for v in filtered["Delivered"]]),
                 ("Opens",     [int(v) for v in filtered["Opens"]]),
                 ("Clicks",    [int(v) for v in filtered["Clicks"]])],
                x=(box_l + Inches(0.1)) / 914400,
                y=(box_t + Inches(1.05)) / 914400,
                w=(col_w - Inches(0.3)) / 914400,
                h=(row_h - Inches(1.15)) / 914400,
                colors=[BAR_DELIVERED, BAR_OPENS, BAR_CLICKS],
                show_legend=True, data_labels=True, gap_width=55,
            )
            # Make legend font smaller (7pt)
            if chart and chart.has_legend:
                _set_legend_font(chart, font_size_pt=7)
            # Make category axis (campaign names) font Aptos 8pt
            if chart:
                try:
                    chart.category_axis.tick_labels.font.size = Pt(8)
                    chart.category_axis.tick_labels.font.name = "Aptos"
                except Exception:
                    pass

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – MONTHLY DELIVERY BY TA
# ══════════════════════════════════════════════════════════════════════════════

def slide3_monthly_ta(prs, df, tmp_dir):
    mta = monthly_ta_kpis(df)
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "REE Delivery Date Overview – TAs")

    # Native grouped bar: one series per TA using TA_CONFIG display labels
    months = [str(m)[:3] for m in mta["Month"].tolist()]
    ta_colors = [c for _, _, c in TA_CONFIG]

    series_data = []
    for _, ta_label, _ in TA_CONFIG:
        vals = [int(mta[ta_label].iloc[i]) if ta_label in mta.columns else 0
                for i in range(len(mta))]
        series_data.append((ta_label, vals))

    _add_bar_chart(
        slide, months, series_data,
        x=0.4, y=1.0, w=12.5, h=6.2,
        colors=ta_colors,
        show_legend=True, data_labels=True, gap_width=15,
    )
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – BEST PERFORMING CAMPAIGNS – ASSETS (top-3 cards)
# ══════════════════════════════════════════════════════════════════════════════

def slide4_best_assets(prs, df, tmp_dir):
    best, latest = best_campaigns(df)
    best = best.head(3)   # assets slide shows top 3
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "Best Performing Campaigns – REE")
    add_text(slide, f"Top 3 by Open Rate  |  Assets View ({latest})", Inches(0.4), Inches(0.62),
             Inches(8), Inches(0.28), font_size=9, italic=True,
             color=RGBColor(0x77, 0x77, 0x77))

    card_w  = Inches(3.9)
    card_h  = Inches(5.9)
    card_t  = Inches(1.05)
    gap     = Inches(0.3)
    total_w = 3 * card_w + 2 * gap
    start_l = (SLIDE_W - total_w) / 2

    for i, (_, row) in enumerate(best.iterrows()):
        c_l = start_l + i * (card_w + gap)

        # Card background
        add_rect(slide, c_l, card_t, card_w, card_h,
                 fill_rgb=RGBColor(0xF4, 0xFA, 0xF8),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.75)

        # Email asset placeholder (large, at top of card)
        add_rect(slide, c_l + Inches(0.2), card_t + Inches(0.12),
                 card_w - Inches(0.4), Inches(3.5),
                 fill_rgb=RGBColor(0xE0, 0xF0, 0xEC),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.5)
        add_text(slide, "📧\n(Email Preview)", c_l + Inches(0.2), card_t + Inches(1.4),
                 card_w - Inches(0.4), Inches(0.8),
                 font_size=11, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

        # Campaign name — below the asset box, prominent, full name visible
        name = str(row["Campaign"])
        add_text(slide, name, c_l + Inches(0.1), card_t + Inches(3.66),
                 card_w - Inches(0.2), Inches(0.7),
                 font_size=9, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER, wrap=True)

        # TA label
        ta_real = str(row["TA"]) if row["TA"] else ""
        ta_disp = next((lbl for rn, lbl, _ in TA_CONFIG if rn == ta_real), ta_real.upper())
        add_text(slide, ta_disp, c_l + Inches(0.1), card_t + Inches(4.42),
                 card_w - Inches(0.2), Inches(0.25),
                 font_size=9, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

        # OR badge — at the very bottom
        add_rect(slide, c_l + Inches(0.3), card_t + Inches(4.72),
                 card_w - Inches(0.6), Inches(0.45),
                 fill_rgb=DARK_GREEN)
        add_text(slide, f"OR  {row['OR']:.2f}%",
                 c_l + Inches(0.3), card_t + Inches(4.72),
                 card_w - Inches(0.6), Inches(0.45),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 num_font=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – BEST PERFORMING CAMPAIGNS (bars + insights)
# ══════════════════════════════════════════════════════════════════════════════

def generate_insights(row: pd.Series) -> list:
    insights = []
    or_val  = row["OR"]
    ctr_val = row["CTR"]
    if or_val >= 50:
        insights.append(f"Very strong {or_val:.1f}% open rate — subject line and sender name are highly effective.")
    elif or_val >= 30:
        insights.append(f"Solid {or_val:.1f}% open rate — message relevance resonates with recipients.")
    else:
        insights.append(f"Open rate at {or_val:.1f}% — consider A/B testing subject lines to improve reach.")

    if ctr_val == 0:
        insights.append("Opens are not converting to clicks — improving CTA placement can drive action.")
    elif ctr_val < 2:
        insights.append(f"CTR of {ctr_val:.2f}% — optimising subject line + preview text can push performance higher.")
    else:
        insights.append(f"Good CTR of {ctr_val:.2f}% — strong audience match and compelling content flow.")
    return insights


def slide5_best_bars(prs, df, tmp_dir):
    best, latest = best_campaigns(df)
    best = best.head(3)   # bars slide shows top 3
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "Best Performing Campaigns – REE")
    add_text(slide, f"Top 3 by Open Rate  |  Performance View ({latest})", Inches(0.4), Inches(0.62),
             Inches(8), Inches(0.28), font_size=9, italic=True,
             color=RGBColor(0x77, 0x77, 0x77))

    n = len(best)
    panel_w = (SLIDE_W - Inches(0.6)) / max(n, 1) - Inches(0.2)
    panel_h = Inches(5.9)
    panel_t = Inches(1.0)
    gap     = Inches(0.25)

    for i, (_, row) in enumerate(best.iterrows()):
        p_l = Inches(0.3) + i * (panel_w + gap)

        # Panel
        add_rect(slide, p_l, panel_t, panel_w, panel_h,
                 fill_rgb=RGBColor(0xF8, 0xFD, 0xFB),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.75)

        # TA label — Trebuchet MS, grey color
        ta_label = str(row["TA"]).upper()
        ta_txb = add_text(slide, ta_label, p_l, panel_t + Inches(0.08),
                 panel_w, Inches(0.22), font_size=9, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
        ta_txb.text_frame.paragraphs[0].runs[0].font.name = "Trebuchet MS"

        # Campaign name — Trebuchet MS, dark navy, ALL CAPS
        camp_name = str(row["Campaign"]).upper()
        camp_txb = add_text(slide, camp_name,
                 p_l + Inches(0.05), panel_t + Inches(0.28),
                 panel_w - Inches(0.1), Inches(0.42),
                 font_size=9, bold=True, color=RGBColor(0x15, 0x25, 0x48),
                 align=PP_ALIGN.CENTER, wrap=True)
        camp_txb.text_frame.paragraphs[0].runs[0].font.name = "Trebuchet MS"

        # Native bar chart — Delivered + Opens + Clicks, same brand colors
        month_label = str(latest)[:3]
        chart = _add_bar_chart(
            slide,
            [month_label],
            [("Delivered", [int(row["Delivered"])]),
             ("Opens",     [int(row["Opens"])]),
             ("Clicks",    [int(row["Clicks"])])],
            x=(p_l + Inches(0.1)) / 914400,
            y=(panel_t + Inches(0.75)) / 914400,
            w=(panel_w - Inches(0.2)) / 914400,
            h=2.4,
            colors=[BAR_DELIVERED, BAR_OPENS, BAR_CLICKS],
            show_legend=True, data_labels=True, gap_width=100,
        )
        if chart and chart.has_legend:
            _set_legend_font(chart, font_size_pt=7)

        # KPI mini row — Total Delivered + OR: Roboto title, Roboto Condensed number
        kpi_y = panel_t + Inches(3.28)
        kpi_data_b = [
            ("Total Delivered", f"{int(row['Delivered']):,}"),
            ("Open Rate",       f"{row['OR']:.2f}%"),
        ]
        mk_w = panel_w / 2
        for mi, (ml, mv) in enumerate(kpi_data_b):
            mk_l = p_l + mi * mk_w
            lbl_txb = add_text(slide, ml, mk_l, kpi_y, mk_w, Inches(0.25),
                     font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.CENTER)
            lbl_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto"
            val_txb = add_text(slide, mv, mk_l, kpi_y + Inches(0.23), mk_w, Inches(0.38),
                     font_size=13, bold=True, color=DARK_GREEN,
                     align=PP_ALIGN.CENTER, num_font=False)
            val_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto Condensed"

        # Insights
        insights_y = panel_t + Inches(3.82)
        add_text(slide, "INSIGHTS", p_l + Inches(0.1), insights_y,
                 panel_w - Inches(0.2), Inches(0.22),
                 font_size=7.5, bold=True, color=DARK_TEXT)
        for j, ins in enumerate(generate_insights(row)):
            add_text(slide, f"• {ins}", p_l + Inches(0.1),
                     insights_y + Inches(0.25 + j * 0.55),
                     panel_w - Inches(0.2), Inches(0.52),
                     font_size=7, color=DARK_TEXT, wrap=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – BEST PERFORMING CAMPAIGNS (Assets View, OVERALL — all months)
# ══════════════════════════════════════════════════════════════════════════════

def slide6_best_assets_overall(prs, df, tmp_dir):
    best = best_campaigns_overall(df)
    best = best.head(3)   # assets slide shows top 3
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "Best Performing Campaigns – REE")
    add_text(slide, "Top 3 by Open Rate  |  Assets View (Overall)", Inches(0.4), Inches(0.62),
             Inches(8), Inches(0.28), font_size=9, italic=True,
             color=RGBColor(0x77, 0x77, 0x77))

    card_w  = Inches(3.9)
    card_h  = Inches(5.9)
    card_t  = Inches(1.05)
    gap     = Inches(0.3)
    total_w = 3 * card_w + 2 * gap
    start_l = (SLIDE_W - total_w) / 2

    for i, (_, row) in enumerate(best.iterrows()):
        c_l = start_l + i * (card_w + gap)

        # Card background
        add_rect(slide, c_l, card_t, card_w, card_h,
                 fill_rgb=RGBColor(0xF4, 0xFA, 0xF8),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.75)

        # Email asset placeholder (large, at top of card)
        add_rect(slide, c_l + Inches(0.2), card_t + Inches(0.12),
                 card_w - Inches(0.4), Inches(3.5),
                 fill_rgb=RGBColor(0xE0, 0xF0, 0xEC),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.5)
        add_text(slide, "📧\n(Email Preview)", c_l + Inches(0.2), card_t + Inches(1.4),
                 card_w - Inches(0.4), Inches(0.8),
                 font_size=11, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

        # Campaign name — below the asset box, prominent, full name visible
        name = str(row["Campaign"])
        add_text(slide, name, c_l + Inches(0.1), card_t + Inches(3.66),
                 card_w - Inches(0.2), Inches(0.7),
                 font_size=9, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER, wrap=True)

        # TA label
        ta_real = str(row["TA"]) if row["TA"] else ""
        ta_disp = next((lbl for rn, lbl, _ in TA_CONFIG if rn == ta_real), ta_real.upper())
        add_text(slide, ta_disp, c_l + Inches(0.1), card_t + Inches(4.42),
                 card_w - Inches(0.2), Inches(0.25),
                 font_size=9, bold=True, color=DARK_GREEN, align=PP_ALIGN.CENTER)

        # OR badge — at the very bottom
        add_rect(slide, c_l + Inches(0.3), card_t + Inches(4.72),
                 card_w - Inches(0.6), Inches(0.45),
                 fill_rgb=DARK_GREEN)
        add_text(slide, f"OR  {row['OR']:.2f}%",
                 c_l + Inches(0.3), card_t + Inches(4.72),
                 card_w - Inches(0.6), Inches(0.45),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 num_font=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – BEST PERFORMING CAMPAIGNS (bars + insights, OVERALL — all months)
# ══════════════════════════════════════════════════════════════════════════════

def slide7_best_bars_overall(prs, df, tmp_dir):
    best = best_campaigns_overall(df)
    best = best.head(3)   # bars slide shows top 3
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_rgb=WHITE)
    add_slide_title(slide, "Best Performing Campaigns – REE")
    add_text(slide, "Top 3 by Open Rate  |  Performance View (Overall)", Inches(0.4), Inches(0.62),
             Inches(8), Inches(0.28), font_size=9, italic=True,
             color=RGBColor(0x77, 0x77, 0x77))

    n = len(best)
    panel_w = (SLIDE_W - Inches(0.6)) / max(n, 1) - Inches(0.2)
    panel_h = Inches(5.9)
    panel_t = Inches(1.0)
    gap     = Inches(0.25)

    for i, (_, row) in enumerate(best.iterrows()):
        p_l = Inches(0.3) + i * (panel_w + gap)

        # Panel
        add_rect(slide, p_l, panel_t, panel_w, panel_h,
                 fill_rgb=RGBColor(0xF8, 0xFD, 0xFB),
                 line_rgb=RGBColor(0xCC, 0xE4, 0xDE), line_width_pt=0.75)

        # TA label — Trebuchet MS, grey color
        ta_label = str(row["TA"]).upper()
        ta_txb = add_text(slide, ta_label, p_l, panel_t + Inches(0.08),
                 panel_w, Inches(0.22), font_size=9, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
        ta_txb.text_frame.paragraphs[0].runs[0].font.name = "Trebuchet MS"

        # Campaign name — Trebuchet MS, dark navy, ALL CAPS
        camp_name = str(row["Campaign"]).upper()
        camp_txb = add_text(slide, camp_name,
                 p_l + Inches(0.05), panel_t + Inches(0.28),
                 panel_w - Inches(0.1), Inches(0.42),
                 font_size=9, bold=True, color=RGBColor(0x15, 0x25, 0x48),
                 align=PP_ALIGN.CENTER, wrap=True)
        camp_txb.text_frame.paragraphs[0].runs[0].font.name = "Trebuchet MS"

        # Native bar chart — Delivered + Opens + Clicks, summed across all months
        chart = _add_bar_chart(
            slide,
            ["Overall"],
            [("Delivered", [int(row["Delivered"])]),
             ("Opens",     [int(row["Opens"])]),
             ("Clicks",    [int(row["Clicks"])])],
            x=(p_l + Inches(0.1)) / 914400,
            y=(panel_t + Inches(0.75)) / 914400,
            w=(panel_w - Inches(0.2)) / 914400,
            h=2.4,
            colors=[BAR_DELIVERED, BAR_OPENS, BAR_CLICKS],
            show_legend=True, data_labels=True, gap_width=100,
        )
        if chart and chart.has_legend:
            _set_legend_font(chart, font_size_pt=7)

        # KPI mini row — Total Delivered + OR: Roboto title, Roboto Condensed number
        kpi_y = panel_t + Inches(3.28)
        kpi_data_b = [
            ("Total Delivered", f"{int(row['Delivered']):,}"),
            ("Open Rate",       f"{row['OR']:.2f}%"),
        ]
        mk_w = panel_w / 2
        for mi, (ml, mv) in enumerate(kpi_data_b):
            mk_l = p_l + mi * mk_w
            lbl_txb = add_text(slide, ml, mk_l, kpi_y, mk_w, Inches(0.25),
                     font_size=7, color=RGBColor(0x55,0x55,0x55), align=PP_ALIGN.CENTER)
            lbl_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto"
            val_txb = add_text(slide, mv, mk_l, kpi_y + Inches(0.23), mk_w, Inches(0.38),
                     font_size=13, bold=True, color=DARK_GREEN,
                     align=PP_ALIGN.CENTER, num_font=False)
            val_txb.text_frame.paragraphs[0].runs[0].font.name = "Roboto Condensed"

        # Insights
        insights_y = panel_t + Inches(3.82)
        add_text(slide, "INSIGHTS", p_l + Inches(0.1), insights_y,
                 panel_w - Inches(0.2), Inches(0.22),
                 font_size=7.5, bold=True, color=DARK_TEXT)
        for j, ins in enumerate(generate_insights(row)):
            add_text(slide, f"• {ins}", p_l + Inches(0.1),
                     insights_y + Inches(0.25 + j * 0.55),
                     panel_w - Inches(0.2), Inches(0.52),
                     font_size=7, color=DARK_TEXT, wrap=True)

    return slide



# ══════════════════════════════════════════════════════════════════════════════

TEAL_FILL  = PatternFill("solid", fgColor="006B63")
DGREEN_FILL= PatternFill("solid", fgColor="00403A")
LTEAL_FILL = PatternFill("solid", fgColor="E8F5F4")
WHITE_FILL = PatternFill("solid", fgColor="FFFFFF")

def hdr_font(sz=11):  return Font(name="Calibri", bold=True, color="FFFFFF", size=sz)
def subhdr_font():    return Font(name="Calibri", bold=True, color="00403A", size=10)
def body_font(sz=10): return Font(name="Calibri", size=sz, color="1A1A2E")
def pct_fmt():        return "0.00%"

thin = Side(style="thin", color="CCE8E2")
border = Border(left=thin, right=thin, top=thin, bottom=thin)

def set_hdr(ws, row, col, val):
    c = ws.cell(row=row, column=col, value=val)
    c.font = hdr_font()
    c.fill = DGREEN_FILL
    c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    c.border = border
    return c

def set_body(ws, row, col, val, fmt=None, bold=False):
    c = ws.cell(row=row, column=col, value=val)
    c.font = Font(name="Calibri", size=10, color="1A1A2E", bold=bold)
    c.alignment = Alignment(horizontal="center", vertical="center")
    c.border = border
    if fmt:
        c.number_format = fmt
    return c


def build_excel(df, out_path):
    wb = Workbook()

    # ── Sheet 1: Overview ──────────────────────────────────────────────────
    ws1 = wb.active
    ws1.title = "Overview"
    ws1.freeze_panes = "A3"

    hdrs = ["Month", "Total Delivered", "Total Opens", "OR (%)", "Total Clicks", "CTR (%)", "Bounced", "Dropped"]
    for ci, h in enumerate(hdrs, 1):
        set_hdr(ws1, 1, ci, h)

    monthly = monthly_kpis(df)
    monthly_bounced = df.groupby("Month", observed=True).apply(
        lambda g: (g["Status"].str.lower() == "bounced").sum()).reset_index(name="Bounced")
    monthly_dropped = df.groupby("Month", observed=True).apply(
        lambda g: (g["Status"].str.lower() == "dropped").sum()).reset_index(name="Dropped")
    monthly = monthly.merge(monthly_bounced, on="Month").merge(monthly_dropped, on="Month")

    for ri, row in enumerate(monthly.itertuples(), 2):
        ws1.cell(ri, 1, str(row.Month)).font = body_font()
        ws1.cell(ri, 1).alignment = Alignment(horizontal="center")
        ws1.cell(ri, 1).border = border
        set_body(ws1, ri, 2, int(row.Delivered))
        set_body(ws1, ri, 3, int(row.Opens))
        set_body(ws1, ri, 4, row.OR / 100, fmt=pct_fmt())
        set_body(ws1, ri, 5, int(row.Clicks))
        set_body(ws1, ri, 6, row.CTR / 100, fmt=pct_fmt())
        set_body(ws1, ri, 7, int(row.Bounced))
        set_body(ws1, ri, 8, int(row.Dropped))

    # Totals row
    last = len(monthly) + 2
    set_body(ws1, last, 1, "TOTAL", bold=True)
    for ci in range(2, 9):
        col_letter = get_column_letter(ci)
        if ci in [4, 6]:
            set_body(ws1, last, ci,
                     f"=AVERAGE({col_letter}2:{col_letter}{last-1})", fmt=pct_fmt(), bold=True)
        else:
            set_body(ws1, last, ci,
                     f"=SUM({col_letter}2:{col_letter}{last-1})", bold=True)

    for ci in range(1, 9):
        ws1.column_dimensions[get_column_letter(ci)].width = 16
    ws1.row_dimensions[1].height = 30

    # ── Sheet 2: TA Overview ───────────────────────────────────────────────
    ws2 = wb.create_sheet("TA Overview")
    ws2.freeze_panes = "A2"

    hdrs2 = ["TA", "Campaigns", "Total Delivered", "Total Opens", "OR (%)",
             "Total Clicks", "CTR (%)", "Bounced", "Dropped"]
    for ci, h in enumerate(hdrs2, 1):
        set_hdr(ws2, 1, ci, h)

    ta_df = ta_kpis(df)
    for ri, row in enumerate(ta_df.itertuples(), 2):
        ws2.cell(ri, 1, row.TA).font = Font(name="Calibri", bold=True, size=10, color="00403A")
        ws2.cell(ri, 1).alignment = Alignment(horizontal="center")
        ws2.cell(ri, 1).border = border
        set_body(ws2, ri, 2, int(row.Campaigns))
        set_body(ws2, ri, 3, int(row.Delivered))
        set_body(ws2, ri, 4, int(row.Opens))
        set_body(ws2, ri, 5, row.OR / 100, fmt=pct_fmt())
        set_body(ws2, ri, 6, int(row.Clicks))
        set_body(ws2, ri, 7, row.CTR / 100, fmt=pct_fmt())
        set_body(ws2, ri, 8, int(row.Bounced))
        set_body(ws2, ri, 9, int(row.Dropped))

    for ci in range(1, 10):
        ws2.column_dimensions[get_column_letter(ci)].width = 17
    ws2.row_dimensions[1].height = 30

    # ── Sheet 3: Monthly × TA ─────────────────────────────────────────────
    ws3 = wb.create_sheet("Monthly by TA")
    ws3.freeze_panes = "A2"
    tas = ["DIA","HAC","ONCO","VACC"]
    hdrs3 = ["Month"] + tas + ["Total"]
    for ci, h in enumerate(hdrs3, 1):
        set_hdr(ws3, 1, ci, h)

    mta = monthly_ta_kpis(df)
    for ri, row in enumerate(mta.itertuples(), 2):
        ws3.cell(ri, 1, str(row.Month)).font = body_font()
        ws3.cell(ri, 1).alignment = Alignment(horizontal="center")
        ws3.cell(ri, 1).border = border
        row_sum = 0
        for ci, ta in enumerate(tas, 2):
            val = int(getattr(row, ta, 0))
            row_sum += val
            set_body(ws3, ri, ci, val)
        set_body(ws3, ri, len(tas) + 2, row_sum, bold=True)

    for ci in range(1, len(hdrs3) + 1):
        ws3.column_dimensions[get_column_letter(ci)].width = 15
    ws3.row_dimensions[1].height = 30

    # ── Sheet 4: Best Campaigns ───────────────────────────────────────────
    ws4 = wb.create_sheet("Best Campaigns")
    ws4.freeze_panes = "A2"
    hdrs4 = ["Month", "TA", "Campaign", "Total Delivered", "Total Opens",
             "OR (%)", "Total Clicks", "CTR (%)"]
    for ci, h in enumerate(hdrs4, 1):
        set_hdr(ws4, 1, ci, h)

    best, latest = best_campaigns(df)
    for ri, row in enumerate(best.itertuples(), 2):
        ws4.cell(ri, 1, str(latest)).font = body_font()
        ws4.cell(ri, 1).alignment = Alignment(horizontal="center")
        ws4.cell(ri, 1).border = border
        ws4.cell(ri, 2, row.TA).font = Font(name="Calibri", bold=True, size=10, color="00403A")
        ws4.cell(ri, 2).alignment = Alignment(horizontal="center")
        ws4.cell(ri, 2).border = border
        ws4.cell(ri, 3, row.Campaign).font = body_font()
        ws4.cell(ri, 3).border = border
        set_body(ws4, ri, 4, int(row.Delivered))
        set_body(ws4, ri, 5, int(row.Opens))
        set_body(ws4, ri, 6, row.OR / 100, fmt=pct_fmt())
        set_body(ws4, ri, 7, int(row.Clicks))
        set_body(ws4, ri, 8, row.CTR / 100, fmt=pct_fmt())

    ws4.column_dimensions["A"].width = 12
    ws4.column_dimensions["B"].width = 10
    ws4.column_dimensions["C"].width = 30
    for ci in range(4, 9):
        ws4.column_dimensions[get_column_letter(ci)].width = 16
    ws4.row_dimensions[1].height = 30

    wb.save(out_path)
    print(f"  ✓ Excel saved: {out_path}")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_ree_report.py <input_excel_or_csv>")
        sys.exit(1)

    input_file = sys.argv[1]
    base = os.path.splitext(os.path.basename(input_file))[0]
    out_dir  = os.path.dirname(os.path.abspath(input_file))
    tmp_dir  = out_dir   # no longer used for images, kept for signature compat

    ppt_out  = os.path.join(out_dir, f"{base}_REE_Report.pptx")
    xlsx_out = os.path.join(out_dir, f"{base}_REE_Report.xlsx")

    print(f"Loading data: {input_file}")
    if input_file.endswith(".csv"):
        df = pd.read_csv(input_file)
        df = normalise_df(df)
    else:
        df = load_and_calc(input_file)

    print(f"  Rows loaded: {len(df)}")

    # ── Build PPT ──
    print("Building PowerPoint…")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide1_overview(prs, df, tmp_dir)
    print("  Slide 1 ✓ REE Overview")
    slide2_ta_overview(prs, df, tmp_dir)
    print("  Slide 2 ✓ TA Overview")
    slide3_monthly_ta(prs, df, tmp_dir)
    print("  Slide 3 ✓ Monthly by TA")
    slide4_best_assets(prs, df, tmp_dir)
    print("  Slide 4 ✓ Best Campaigns – Assets")
    slide5_best_bars(prs, df, tmp_dir)
    print("  Slide 5 ✓ Best Campaigns – Bars + Insights")
    slide6_best_assets_overall(prs, df, tmp_dir)
    print("  Slide 6 ✓ Best Campaigns – Assets (Overall)")
    slide7_best_bars_overall(prs, df, tmp_dir)
    print("  Slide 7 ✓ Best Campaigns – Bars + Insights (Overall)")

    prs.save(ppt_out)
    print(f"  ✓ PPT saved: {ppt_out}")

    # ── Build Excel ──
    print("Building Excel report…")
    build_excel(df, xlsx_out)

    print("\nDone! Files saved:")
    print(f"  PPT:   {ppt_out}")
    print(f"  Excel: {xlsx_out}")


if __name__ == "__main__":
    main()
