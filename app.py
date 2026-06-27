import streamlit as st
import subprocess
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import tempfile
import shutil
import time

# ── Page config ──────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Performance Deck Generator",
    page_icon="📊",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ── Custom CSS ────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Space+Grotesk:wght@400;500;600;700&display=swap');

/* Base */
html, body, [class*="css"] {
    font-family: 'Inter', sans-serif;
}

.stApp {
    background: linear-gradient(135deg, #020b18 0%, #041425 50%, #051e30 100%);
    min-height: 100vh;
}

/* Hide default streamlit elements */
#MainMenu, footer, header { visibility: hidden; }
.block-container { padding-top: 2rem; padding-bottom: 3rem; max-width: 780px; }

/* Hero header */
.hero {
    text-align: center;
    padding: 2.5rem 1rem 1.5rem;
    margin-bottom: 0.5rem;
}
.hero-badge {
    display: inline-block;
    background: linear-gradient(90deg, #0d4f3c, #0a7c5e);
    color: #2fffc2;
    font-size: 0.7rem;
    font-weight: 700;
    letter-spacing: 0.2em;
    text-transform: uppercase;
    padding: 0.3rem 1rem;
    border-radius: 20px;
    margin-bottom: 1.2rem;
    border: 1px solid #0a7c5e55;
}
.hero-title {
    font-family: 'Space Grotesk', sans-serif;
    font-size: 2.6rem;
    font-weight: 700;
    color: #ffffff;
    line-height: 1.15;
    margin: 0 0 0.6rem;
    letter-spacing: -0.02em;
}
.hero-title span {
    background: linear-gradient(90deg, #2fffc2 0%, #00b4d8 60%, #0077b6 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}
.hero-sub {
    color: #7a9bb5;
    font-size: 0.95rem;
    font-weight: 400;
    line-height: 1.6;
}

/* Divider */
.grad-divider {
    height: 1px;
    background: linear-gradient(90deg, transparent, #0a7c5e55, #00b4d855, transparent);
    margin: 1.8rem 0;
}

/* Platform cards */
.platform-grid {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 0.9rem;
    margin: 1.2rem 0;
}
.platform-card {
    background: linear-gradient(135deg, #071c2e 0%, #0a2540 100%);
    border: 1px solid #0e3a5a;
    border-radius: 12px;
    padding: 1rem 1.2rem;
    display: flex;
    align-items: center;
    gap: 0.75rem;
    transition: border-color 0.2s;
}
.platform-card:hover { border-color: #0a7c5e; }
.platform-icon {
    font-size: 1.4rem;
    width: 2.2rem;
    text-align: center;
    flex-shrink: 0;
}
.platform-name {
    font-weight: 600;
    color: #cde7f5;
    font-size: 0.88rem;
    letter-spacing: 0.01em;
}
.platform-file {
    color: #4a7c99;
    font-size: 0.72rem;
    margin-top: 0.15rem;
    font-family: 'Courier New', monospace;
}

/* Order badge */
.order-badge {
    display: flex;
    gap: 0.5rem;
    align-items: center;
    justify-content: center;
    margin: 1rem 0 1.5rem;
    flex-wrap: wrap;
}
.order-pill {
    background: #071c2e;
    border: 1px solid #0e3a5a;
    border-radius: 20px;
    padding: 0.3rem 0.85rem;
    color: #7abfdf;
    font-size: 0.75rem;
    font-weight: 500;
    display: flex;
    align-items: center;
    gap: 0.4rem;
}
.order-pill .num {
    background: linear-gradient(135deg, #0a7c5e, #00b4d8);
    color: white;
    border-radius: 50%;
    width: 1.2rem;
    height: 1.2rem;
    display: inline-flex;
    align-items: center;
    justify-content: center;
    font-size: 0.65rem;
    font-weight: 700;
    flex-shrink: 0;
}
.arrow { color: #1e4d6a; font-size: 0.8rem; }

/* Generate button */
.stButton > button {
    width: 100%;
    background: linear-gradient(135deg, #0a7c5e 0%, #00897b 50%, #0077b6 100%);
    color: white;
    border: none;
    border-radius: 10px;
    padding: 0.85rem 2rem;
    font-size: 1rem;
    font-weight: 700;
    font-family: 'Space Grotesk', sans-serif;
    letter-spacing: 0.03em;
    cursor: pointer;
    transition: all 0.25s ease;
    box-shadow: 0 4px 24px #0a7c5e33;
    text-transform: uppercase;
}
.stButton > button:hover {
    transform: translateY(-2px);
    box-shadow: 0 8px 32px #0a7c5e55;
    background: linear-gradient(135deg, #0d9e78 0%, #00a693 50%, #0096d6 100%);
}
.stButton > button:active { transform: translateY(0); }

/* Download button */
.stDownloadButton > button {
    width: 100%;
    background: linear-gradient(135deg, #041e30 0%, #071c2e 100%);
    color: #2fffc2;
    border: 1.5px solid #0a7c5e;
    border-radius: 10px;
    padding: 0.85rem 2rem;
    font-size: 1rem;
    font-weight: 700;
    font-family: 'Space Grotesk', sans-serif;
    letter-spacing: 0.03em;
    cursor: pointer;
    transition: all 0.25s ease;
    box-shadow: 0 4px 16px #0a7c5e22;
    text-transform: uppercase;
}
.stDownloadButton > button:hover {
    background: linear-gradient(135deg, #071c2e 0%, #0a2540 100%);
    border-color: #2fffc2;
    box-shadow: 0 6px 24px #0a7c5e44;
    transform: translateY(-1px);
}

/* Log box */
.log-box {
    background: #020d16;
    border: 1px solid #0e3a5a;
    border-radius: 10px;
    padding: 1rem 1.2rem;
    font-family: 'Courier New', monospace;
    font-size: 0.78rem;
    color: #4a9e7a;
    max-height: 240px;
    overflow-y: auto;
    line-height: 1.7;
    margin: 1rem 0;
}
.log-line-ok  { color: #2fffc2; }
.log-line-err { color: #ff6b6b; }
.log-line-inf { color: #00b4d8; }
.log-line-dim { color: #2a5a7a; }

/* Status cards */
.status-row {
    display: grid;
    grid-template-columns: repeat(4, 1fr);
    gap: 0.6rem;
    margin: 1rem 0;
}
.status-card {
    background: #071c2e;
    border: 1px solid #0e3a5a;
    border-radius: 8px;
    padding: 0.6rem 0.5rem;
    text-align: center;
}
.status-card .s-name { color: #4a7c99; font-size: 0.68rem; font-weight: 600; text-transform: uppercase; letter-spacing: 0.08em; }
.status-card .s-icon { font-size: 1.2rem; margin: 0.2rem 0; }
.status-card.done  { border-color: #0a7c5e; background: #041e18; }
.status-card.done .s-name { color: #2fffc2; }
.status-card.error { border-color: #7c1a0a; background: #1e0707; }
.status-card.error .s-name { color: #ff6b6b; }
.status-card.active { border-color: #00b4d8; background: #041a28; }
.status-card.active .s-name { color: #00e5ff; }

/* Success box */
.success-box {
    background: linear-gradient(135deg, #041e18 0%, #051a28 100%);
    border: 1px solid #0a7c5e;
    border-radius: 12px;
    padding: 1.5rem;
    text-align: center;
    margin: 1.2rem 0;
}
.success-box .s-emoji { font-size: 2.5rem; margin-bottom: 0.5rem; }
.success-box .s-title { color: #2fffc2; font-size: 1.1rem; font-weight: 700; font-family: 'Space Grotesk', sans-serif; margin-bottom: 0.3rem; }
.success-box .s-sub { color: #4a9e7a; font-size: 0.82rem; }

/* Error box */
.error-box {
    background: #160404;
    border: 1px solid #7c1a0a;
    border-radius: 10px;
    padding: 1rem 1.2rem;
    color: #ff8a8a;
    font-size: 0.83rem;
    margin: 0.8rem 0;
}
</style>
""", unsafe_allow_html=True)

# ── Constants ─────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).parent

PLATFORMS = [
    {
        "name": "SoMe",
        "script": "generate_some_report(o).py",
        "output_pptx": "SoMe_Report.pptx",
        "excel": "SoMe(Cons4).xlsx",
        "icon": "📱",
        "order": 1,
    },
    {
        "name": "GCC Pulse",
        "script": "gcc_pulse_automation.py",
        "output_pptx": "GCC-pulse data_GCC_Pulse_Report.pptx",
        "excel": "GCC-pulse data.xlsx",
        "icon": "📡",
        "order": 2,
    },
    {
        "name": "SFMC",
        "script": "sfmc_reportwt(1).py",
        "output_pptx": "SFMC_Report.pptx",
        "excel": "SFMC-data.xlsx",
        "icon": "✉️",
        "order": 3,
    },
    {
        "name": "REE",
        "script": "generate_ree_reportnw.py",
        "output_pptx": "REE-data_REE_Report.pptx",
        "excel": "REE-data.xlsx",
        "icon": "⚡",
        "order": 4,
    },
]

MERGED_OUTPUT = "Performance_Deck.pptx"

# ── Helper: merge PPTs ────────────────────────────────────────────────────────
def merge_presentations(pptx_paths: list[Path], output_path: Path):
    """Merge multiple .pptx files into one, preserving all slides."""
    merged = Presentation()
    # Match slide dimensions from the first deck
    first = Presentation(str(pptx_paths[0]))
    merged.slide_width  = first.slide_width
    merged.slide_height = first.slide_height

    from pptx.util import Emu
    from lxml import etree
    import copy

    for pptx_path in pptx_paths:
        src = Presentation(str(pptx_path))
        for slide in src.slides:
            # Add blank slide then replace XML
            blank_layout = merged.slide_layouts[6]  # blank
            new_slide = merged.slides.add_slide(blank_layout)

            # Copy slide XML
            sp_tree = slide.shapes._spTree
            new_sp_tree = new_slide.shapes._spTree
            for elem in list(new_sp_tree):
                new_sp_tree.remove(elem)
            for elem in sp_tree:
                new_sp_tree.append(copy.deepcopy(elem))

            # Copy slide-level relationships (images, etc.)
            for rel in slide.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        img_part = rel.target_part
                        new_slide.part.relate_to(img_part, rel.reltype)
                    except Exception:
                        pass

    merged.save(str(output_path))

# ── Helper: run script ────────────────────────────────────────────────────────
def run_script(script_name: str) -> tuple[bool, str]:
    script_path = BASE_DIR / script_name
    if not script_path.exists():
        return False, f"Script not found: {script_name}"
    try:
        result = subprocess.run(
            [sys.executable, str(script_path)],
            cwd=str(BASE_DIR),
            capture_output=True,
            text=True,
            timeout=300,
        )
        if result.returncode != 0:
            err = result.stderr.strip() or result.stdout.strip()
            return False, err[-800:] if len(err) > 800 else err
        return True, result.stdout.strip()
    except subprocess.TimeoutExpired:
        return False, "Script timed out after 5 minutes."
    except Exception as e:
        return False, str(e)

# ── UI ────────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="hero">
    <div class="hero-badge">📊 Automated Reporting Suite</div>
    <div class="hero-title">Performance Deck<br><span>Generator</span></div>
    <div class="hero-sub">Runs all four platform scripts, merges the decks,<br>and delivers one polished PowerPoint — in one click.</div>
</div>
""", unsafe_allow_html=True)

st.markdown('<div class="grad-divider"></div>', unsafe_allow_html=True)

# Platform overview
st.markdown("""
<div class="platform-grid">
    <div class="platform-card">
        <div class="platform-icon">📱</div>
        <div>
            <div class="platform-name">SoMe</div>
            <div class="platform-file">generate_some_report(o).py</div>
        </div>
    </div>
    <div class="platform-card">
        <div class="platform-icon">📡</div>
        <div>
            <div class="platform-name">GCC Pulse</div>
            <div class="platform-file">gcc_pulse_automation.py</div>
        </div>
    </div>
    <div class="platform-card">
        <div class="platform-icon">✉️</div>
        <div>
            <div class="platform-name">SFMC</div>
            <div class="platform-file">sfmc_reportwt(1).py</div>
        </div>
    </div>
    <div class="platform-card">
        <div class="platform-icon">⚡</div>
        <div>
            <div class="platform-name">REE</div>
            <div class="platform-file">generate_ree_reportnw.py</div>
        </div>
    </div>
</div>
""", unsafe_allow_html=True)

# Merge order
st.markdown("""
<div style="text-align:center; color:#4a7c99; font-size:0.73rem; font-weight:600; letter-spacing:0.12em; text-transform:uppercase; margin: 0.5rem 0 0.4rem;">Merge Order</div>
<div class="order-badge">
    <div class="order-pill"><span class="num">1</span>SoMe</div>
    <span class="arrow">›</span>
    <div class="order-pill"><span class="num">2</span>GCC Pulse</div>
    <span class="arrow">›</span>
    <div class="order-pill"><span class="num">3</span>SFMC</div>
    <span class="arrow">›</span>
    <div class="order-pill"><span class="num">4</span>REE</div>
</div>
""", unsafe_allow_html=True)

st.markdown('<div class="grad-divider"></div>', unsafe_allow_html=True)

# ── Session state ─────────────────────────────────────────────────────────────
if "generated" not in st.session_state:
    st.session_state.generated = False
if "merged_bytes" not in st.session_state:
    st.session_state.merged_bytes = None
if "log_lines" not in st.session_state:
    st.session_state.log_lines = []

# ── Generate button ───────────────────────────────────────────────────────────
if st.button("🚀  Generate Performance Deck", key="gen_btn"):
    st.session_state.generated = False
    st.session_state.merged_bytes = None
    st.session_state.log_lines = []

    log_placeholder   = st.empty()
    status_placeholder = st.empty()

    def render_status(statuses):
        cards = ""
        for p in PLATFORMS:
            s = statuses.get(p["name"], "pending")
            cls  = {"done": "done", "error": "error", "running": "active"}.get(s, "")
            icon = {"done": "✅", "error": "❌", "running": "⏳", "pending": "⬜"}.get(s, "⬜")
            cards += f'<div class="status-card {cls}"><div class="s-name">{p["name"]}</div><div class="s-icon">{icon}</div></div>'
        status_placeholder.markdown(f'<div class="status-row">{cards}</div>', unsafe_allow_html=True)

    def render_log(lines):
        html_lines = ""
        for line in lines[-30:]:
            if line.startswith("✅"):
                html_lines += f'<div class="log-line-ok">{line}</div>'
            elif line.startswith("❌"):
                html_lines += f'<div class="log-line-err">{line}</div>'
            elif line.startswith("ℹ"):
                html_lines += f'<div class="log-line-inf">{line}</div>'
            else:
                html_lines += f'<div class="log-line-dim">{line}</div>'
        log_placeholder.markdown(f'<div class="log-box">{html_lines}</div>', unsafe_allow_html=True)

    logs    = []
    statuses = {p["name"]: "pending" for p in PLATFORMS}
    errors  = []
    pptx_paths = []

    logs.append("ℹ  Starting Performance Deck generation…")
    render_log(logs)
    render_status(statuses)

    for platform in PLATFORMS:
        name   = platform["name"]
        script = platform["script"]
        out_pptx = BASE_DIR / platform["output_pptx"]

        statuses[name] = "running"
        logs.append(f"ℹ  [{name}] Running {script} …")
        render_log(logs)
        render_status(statuses)

        ok, msg = run_script(script)

        if ok:
            if out_pptx.exists():
                pptx_paths.append(out_pptx)
                statuses[name] = "done"
                logs.append(f"✅ [{name}] Done → {platform['output_pptx']}")
            else:
                statuses[name] = "error"
                err = f"Script ran but output not found: {platform['output_pptx']}"
                logs.append(f"❌ [{name}] {err}")
                errors.append((name, err))
        else:
            statuses[name] = "error"
            logs.append(f"❌ [{name}] FAILED")
            for l in msg.splitlines()[-5:]:
                logs.append(f"   {l}")
            errors.append((name, msg))

        render_log(logs)
        render_status(statuses)

    # Merge
    if pptx_paths:
        logs.append(f"ℹ  Merging {len(pptx_paths)} deck(s)…")
        render_log(logs)
        try:
            output_path = BASE_DIR / MERGED_OUTPUT
            merge_presentations(pptx_paths, output_path)
            with open(output_path, "rb") as f:
                st.session_state.merged_bytes = f.read()
            logs.append(f"✅ Merged → {MERGED_OUTPUT}  ({len(st.session_state.merged_bytes)//1024} KB)")
            st.session_state.generated = True
        except Exception as e:
            logs.append(f"❌ Merge failed: {e}")
            errors.append(("Merge", str(e)))
    else:
        logs.append("❌ No decks to merge — all scripts failed.")

    render_log(logs)
    render_status(statuses)
    st.session_state.log_lines = logs

    if errors:
        err_html = "".join(f"<b>{n}:</b> {m[:200]}<br>" for n, m in errors)
        st.markdown(f'<div class="error-box">⚠️ Some steps had errors:<br><br>{err_html}</div>', unsafe_allow_html=True)

# ── Show result ───────────────────────────────────────────────────────────────
if st.session_state.generated and st.session_state.merged_bytes:
    st.markdown("""
    <div class="success-box">
        <div class="s-emoji">🎉</div>
        <div class="s-title">Performance Deck Ready</div>
        <div class="s-sub">All slides merged in order: SoMe → GCC Pulse → SFMC → REE</div>
    </div>
    """, unsafe_allow_html=True)

    st.download_button(
        label="⬇️  Download Performance_Deck.pptx",
        data=st.session_state.merged_bytes,
        file_name="Performance_Deck.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="dl_btn",
    )
